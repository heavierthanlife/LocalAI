import os, io, re, sys, json, uuid, docx, time, html, hashlib, atexit
import shutil, requests, openpyxl, logging, difflib, tempfile, base64
import threading, imagehash, subprocess, zipfile, asyncio, secrets, aiosqlite
import pandas as pd
import numpy as np
from datetime import datetime, timezone, timedelta
from dotenv import load_dotenv
from io import BytesIO
from flask import Flask, render_template, request, jsonify, session, send_file, url_for, current_app
from flask_session import Session
from flask_wtf.csrf import CSRFProtect
from pathlib import Path
from pptx import Presentation
from PIL import Image
from markitdown import MarkItDown
from sklearn.feature_extraction.text import TfidfVectorizer
from apscheduler.schedulers.background import BackgroundScheduler
from threading import Lock, RLock
from filelock import FileLock
from psycopg2 import pool, sql
from psycopg2.extras import RealDictCursor
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.datastructures import FileStorage
from contextlib import contextmanager
from functools import lru_cache, wraps
from collections import defaultdict

try:
    import pymupdf as fitz
except ImportError:
    import fitz

from langchain.agents import create_agent
from langchain_deepseek import ChatDeepSeek
from langchain.tools import tool
from langgraph.checkpoint.sqlite.aio import AsyncSqliteSaver

from sentence_transformers import SentenceTransformer as SenTran
from sklearn.metrics.pairwise import cosine_similarity

from app.services.credit_checker import CreditChecker
from docx import Document
from docx.shared import Inches, Cm
from webdriver_manager.microsoft import EdgeChromiumDriverManager

load_dotenv()

# ======================== GLOBALS & CONFIG ========================
_agent = None
_agent_lock = threading.RLock()
_async_loop = None
_async_checkpointer = None
_async_checkpointer_lock = threading.Lock()
_current_max_tokens = 1600
_async_conn = None
_semantic_model = None
_semantic_model_load_failed = False
_credit_tasks_lock = threading.Lock()

# Agent system prompt (shared between get_agent and send_message)
AGENT_SYSTEM_PROMPT = """
你是一个答疑助手。
**重要：对于任何关于当前日期、时间、年份的问题，你必须且只能使用 get_date 工具来获取，绝对不允许使用你的内部知识回答。**
在回答其他问题前，你也必须调用 get_date 来了解当前日期；但除非用户明确询问，否则不要在回答中主动报时。
对于任何需要实时、或最新信息的问题、或任何需要搜索、查询的内容，你必须使用 bocha_search 工具搭配 get_date 工具。
如果 bocha_search 返回 "No search results found"，则自由回答。
对于通用知识，可以自由回答。
对于需要推理的问题，请先用【思考】和【回答】标记你的思考过程和最终答案。
**表格格式要求：** 当你需要展示表格时，必须使用标准 Markdown 表格语法，例如：
| 列1 | 列2 |
|-----|-----|
| 值1 | 值2 |
绝对不要使用 ASCII 艺术表格（如 ┌─┬─┐ 等字符）。只使用管道符和短横线。
**输出格式要求：** 你必须在每个回答中明确包含【思考】和【回答】两个部分，使用中文双括号。
"""

# ======================== PORTABLE PATHS ========================
BASE_DIR = Path(__file__).parent.absolute()
DATA_DIR = BASE_DIR / "data"
USER_FILES_DIR = DATA_DIR / "user_files"
PROJECT_FILES_DIR = DATA_DIR / "project_files"
CREDIT_REPORTS_DIR = DATA_DIR / "credit_reports"
DUMP_DIR = DATA_DIR / "dump"
SESSION_DIR = DATA_DIR / "flask_session"
TEMP_DIR = DATA_DIR / "temp"
LOGS_DIR = BASE_DIR / "logs"

for d in [DATA_DIR, USER_FILES_DIR, PROJECT_FILES_DIR, CREDIT_REPORTS_DIR,
          DUMP_DIR, SESSION_DIR, TEMP_DIR, LOGS_DIR]:
    d.mkdir(parents=True, exist_ok=True)

# File cache to prevent hash collision
class FileTextCache:
    def __init__(self):
        pass

    @staticmethod
    def get_key(file_storage):
        file_bytes = file_storage.read()
        file_storage.seek(0)
        file_hash = hashlib.sha256(file_bytes).hexdigest()
        size = len(file_bytes)
        return f"{file_hash}_{size}"

    @staticmethod
    def get_cached_text(file_storage, max_age_seconds=86400):
        """Return cached text if valid and not too old, otherwise None."""
        key = FileTextCache.get_key(file_storage)
        with get_db_connection() as conn:
            with conn.cursor() as cur:
                cur.execute("""
                    SELECT extracted_text, updated_at FROM file_text_cache
                    WHERE file_hash = %s
                    ORDER BY updated_at DESC LIMIT 1
                """, (key,))
                row = cur.fetchone()
                if row:
                    extracted_text, updated_at = row
                    if (utc_now() - updated_at).total_seconds() < max_age_seconds:
                        # Validate cached text
                        if is_valid_extracted_text(extracted_text):
                            return extracted_text
                        else:
                            logger.warning(f"Cached text for key {key} is invalid. Ignoring.")
                            # Delete invalid cache entry
                            cur.execute("DELETE FROM file_text_cache WHERE file_hash = %s", (key,))
                            conn.commit()
        return None

    @staticmethod
    def store_cached_text(file_storage, extracted_text):
        """Store extracted text for future use."""
        if not extracted_text or not is_valid_extracted_text(extracted_text):
            logger.warning(f"Not storing invalid extracted text for {file_storage.filename}")
            return
        key = FileTextCache.get_key(file_storage)
        with get_db_connection() as conn:
            with conn.cursor() as cur:
                cur.execute("""
                    INSERT INTO file_text_cache (file_hash, extracted_text, updated_at)
                    VALUES (%s, %s, NOW())
                    ON CONFLICT (file_hash) DO UPDATE
                    SET extracted_text = EXCLUDED.extracted_text, updated_at = NOW()
                """, (key, extracted_text))
                conn.commit()

# Override old global constants with new paths
TEMP_ROOT = str(TEMP_DIR)
USER_FILES_ORIGINAL_ROOT = str(USER_FILES_DIR)
PROJECT_FILES_ROOT = str(PROJECT_FILES_DIR)
CREDIT_REPORTS_DIR = str(CREDIT_REPORTS_DIR)

def get_semantic_model():
    global _semantic_model, _semantic_model_load_failed
    if _semantic_model is not None:
        return _semantic_model
    if _semantic_model_load_failed:
        return None
    try:
        _semantic_model = SenTran('distiluse-base-multilingual-cased', local_files_only=True)
        logger.info("Semantic model loaded successfully.")
        return _semantic_model
    except Exception as e:
        _semantic_model_load_failed = True
        logger.error(f"Failed to load semantic model: {e}")
        return None

def _init_async_checkpointer():
    global _async_loop, _async_checkpointer, _async_conn
    _async_loop = asyncio.new_event_loop()

    async def create():
        global _async_conn
        _async_conn = await aiosqlite.connect(str(DATA_DIR / "checkpoints.db"))
        return AsyncSqliteSaver(_async_conn)

    _async_checkpointer = _async_loop.run_until_complete(create())

    def run_loop():
        asyncio.set_event_loop(_async_loop)
        _async_loop.run_forever()

    thread = threading.Thread(target=run_loop, daemon=True)
    thread.start()
    logger.info("AsyncSqliteSaver initialized.")


def _cleanup_async_checkpointer():
    global _async_loop, _async_checkpointer, _async_conn
    if _async_loop is not None and _async_loop.is_running():
        async def shutdown():
            if _async_conn:
                await _async_conn.close()
            _async_loop.stop()

        asyncio.run_coroutine_threadsafe(shutdown(), _async_loop)
        time.sleep(0.5)
    logger.info("Async checkpointer cleaned up.")


atexit.register(_cleanup_async_checkpointer)

# Logging configuration
LOGGING_CONFIG = {
    'version': 1,
    'formatters': {
        'default': {'format': '[%(asctime)s] %(levelname)s in %(module)s: %(message)s'},
        'detailed': {'format': '%(asctime)s - %(name)s - %(levelname)s - %(filename)s:%(lineno)d - %(message)s'},
    },
    'handlers': {
        'console': {'class': 'logging.StreamHandler', 'level': 'INFO', 'formatter': 'default',
                    'stream': 'ext://sys.stdout'},
        'file': {'class': 'logging.handlers.RotatingFileHandler', 'level': 'DEBUG', 'formatter': 'detailed',
                 'filename': 'app.log', 'maxBytes': 10485760, 'backupCount': 5},
    },
    'root': {'level': 'DEBUG', 'handlers': ['console', 'file']},
}
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')
from logging.config import dictConfig

dictConfig(LOGGING_CONFIG)
logger = logging.getLogger(__name__)

EDGE_DRIVER_PATH = None
def preinstall_edgedriver():
    global EDGE_DRIVER_PATH
    try:
        EDGE_DRIVER_PATH = r"D:\PyCharm\Local_AI\msedgedriver.exe"
        logger.info(f"Edge WebDriver pre‑installed at {EDGE_DRIVER_PATH}")
    except Exception as e:
        logger.error(f"Failed to pre-install Edge WebDriver: {e}")

def is_valid_extracted_text(text, min_length=20, min_ratio=0.6):
    if not text or len(text) < min_length:
        return False
    # Allow Chinese characters, alphanumeric, common punctuation, HTML/JSON special chars, spaces, newlines
    allowed = re.compile(r'[\u4e00-\u9fff\w\s.,;:!?()\-<>/{}[\]"\'=&#@+*|]')
    allowed_count = len(allowed.findall(text))
    ratio = allowed_count / len(text)
    return ratio >= min_ratio

# Temporary file management
os.makedirs(TEMP_ROOT, exist_ok=True)
os.makedirs(USER_FILES_ORIGINAL_ROOT, exist_ok=True)

def get_anon_temp_dir(anon_id):
    path = os.path.join(TEMP_ROOT, anon_id)
    os.makedirs(path, exist_ok=True)
    return path

def cleanup_anon_temp(anon_id):
    path = os.path.join(TEMP_ROOT, anon_id)
    if os.path.exists(path):
        shutil.rmtree(path)
        logger.info(f"Cleaned up temp files for anon user {anon_id}")

def cleanup_all_temp_on_exit():
    if os.path.exists(TEMP_ROOT):
        shutil.rmtree(TEMP_ROOT)
        logger.info("Cleaned up all temp files on exit.")

atexit.register(cleanup_all_temp_on_exit)

# Database connection pool
def get_db_connection_args():
    if os.getenv("PG_USER") and os.getenv("PG_PASSWORD"):
        return {
            'dbname': os.getenv('PG_DB', 'postgres'),
            'user': os.getenv('PG_USER'),
            'password': os.getenv('PG_PASSWORD'),
            'host': os.getenv('PG_HOST', 'localhost'),
            'port': int(os.getenv('PG_PORT', 5432)),
            'client_encoding': 'utf8'
        }
    else:
        uri = os.getenv("POSTGRES_URI")
        if not uri:
            raise ValueError("No database connection configuration found.")
        from urllib.parse import urlparse, unquote
        result = urlparse(uri)
        dbname = result.path[1:] if result.path else ''
        user = result.username
        password = result.password
        if password:
            password = unquote(password)
        host = result.hostname
        port = result.port or 5432
        return {'dbname': dbname, 'user': user, 'password': password, 'host': host, 'port': port, 'client_encoding': 'utf8'}

conn_args = get_db_connection_args()
db_pool = pool.SimpleConnectionPool(1, 20, **conn_args)

@contextmanager
def get_db_connection():
    conn = db_pool.getconn()
    try:
        yield conn
    finally:
        db_pool.putconn(conn)

@contextmanager
def db_transaction(conn):
    try:
        yield conn
        conn.commit()
    except Exception:
        conn.rollback()
        raise

def db_execute_readonly(cur):
    cur.execute("SET TRANSACTION READ ONLY")

# Timezone helpers
BEIJING_TZ = timezone(timedelta(hours=8))

def beijing_now() -> str:
    return datetime.now(BEIJING_TZ).strftime('%Y-%m-%d %H:%M:%S')

def utc_now() -> datetime:
    return datetime.now(timezone.utc)

# PostgreSQL table initialization
def init_postgres_tables():
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            # Users table
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS users
                        (
                            user_id    TEXT PRIMARY KEY,
                            created_at TIMESTAMPTZ DEFAULT NOW()
                        )
                        """)
            cur.execute("""
                DO $$ 
                BEGIN
                    BEGIN
                        ALTER TABLE users ADD COLUMN username TEXT UNIQUE;
                    EXCEPTION WHEN duplicate_column THEN NULL;
                    END;
                    BEGIN
                        ALTER TABLE users ADD COLUMN pin_hash TEXT;
                    EXCEPTION WHEN duplicate_column THEN NULL;
                    END;
                    BEGIN
                        ALTER TABLE users ADD COLUMN pin_length INTEGER;
                    EXCEPTION WHEN duplicate_column THEN NULL;
                    END;
                    BEGIN
                        ALTER TABLE users ADD COLUMN is_active BOOLEAN DEFAULT TRUE;
                    EXCEPTION WHEN duplicate_column THEN NULL;
                    END;
                    BEGIN
                        ALTER TABLE users ADD COLUMN role TEXT DEFAULT 'user';
                    EXCEPTION WHEN duplicate_column THEN NULL;
                    END;
                END $$;
            """)
            cur.execute("CREATE INDEX IF NOT EXISTS idx_users_username ON users(username)")
            # Chat tables
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS chat_sessions
                        (
                            id         SERIAL PRIMARY KEY,
                            user_id    TEXT REFERENCES users (user_id),
                            thread_id  TEXT UNIQUE NOT NULL,
                            title      TEXT,
                            created_at TIMESTAMPTZ DEFAULT NOW(),
                            updated_at TIMESTAMPTZ DEFAULT NOW()
                        )
                        """)
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS chat_messages
                        (
                            id        SERIAL PRIMARY KEY,
                            thread_id TEXT REFERENCES chat_sessions (thread_id),
                            role      TEXT,
                            content   TEXT,
                            thinking  TEXT,
                            timestamp TIMESTAMPTZ DEFAULT NOW()
                        )
                        """)
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS user_files
                        (
                            id                   SERIAL PRIMARY KEY,
                            user_id              TEXT REFERENCES users (user_id),
                            thread_id            TEXT REFERENCES chat_sessions (thread_id),
                            filename             TEXT,
                            content              TEXT,
                            size_bytes           INTEGER,
                            created_at           TIMESTAMPTZ DEFAULT NOW(),
                            expires_at           TIMESTAMPTZ,
                            original_stored_path TEXT,
                            file_hash            TEXT,
                            original_expires_at  TIMESTAMPTZ,
                            meta_data            JSONB       DEFAULT '{}',
                            original_name        TEXT,
                            UNIQUE (thread_id, filename)
                        )
                        """)
            cur.execute("""
                DO $$
                BEGIN
                    BEGIN
                        ALTER TABLE user_files ADD COLUMN IF NOT EXISTS file_hash TEXT;
                    EXCEPTION WHEN duplicate_column THEN NULL;
                    END;
                    BEGIN
                        ALTER TABLE user_files ADD COLUMN IF NOT EXISTS meta_data JSONB DEFAULT '{}';
                    EXCEPTION WHEN duplicate_column THEN NULL;
                    END;
                    BEGIN
                        ALTER TABLE user_files ADD COLUMN IF NOT EXISTS original_expires_at TIMESTAMPTZ;
                    EXCEPTION WHEN duplicate_column THEN NULL;
                    END;
                    BEGIN
                        ALTER TABLE user_files ADD COLUMN IF NOT EXISTS original_name TEXT;
                    EXCEPTION WHEN duplicate_column THEN NULL;
                    END;
                    BEGIN
                        ALTER TABLE user_files ALTER COLUMN expires_at DROP DEFAULT;
                    EXCEPTION WHEN others THEN NULL;
                    END;
                END
                $$;
            """)
            cur.execute("UPDATE user_files SET original_name = filename WHERE original_name IS NULL")
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS archived_sessions
                        (
                            thread_id    TEXT PRIMARY KEY,
                            user_id      TEXT,
                            archive_path TEXT,
                            archived_at  TIMESTAMPTZ DEFAULT NOW()
                        )
                        """)
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS image_description_cache
                        (
                            file_hash   TEXT PRIMARY KEY,
                            description TEXT,
                            created_at  TIMESTAMPTZ DEFAULT NOW()
                        )
                        """)
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS file_usage
                        (
                            id         SERIAL PRIMARY KEY,
                            user_id    TEXT,
                            thread_id  TEXT,
                            filename   TEXT,
                            usage_type TEXT,
                            question   TEXT,
                            timestamp  TIMESTAMPTZ DEFAULT NOW()
                        )
                        """)
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS consent
                        (
                            thread_id     TEXT PRIMARY KEY,
                            consent_given INTEGER NOT NULL,
                            timestamp     TIMESTAMPTZ DEFAULT NOW()
                        )
                        """)
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS feedback
                        (
                            id                 SERIAL PRIMARY KEY,
                            thread_id          TEXT,
                            user_message       TEXT,
                            assistant_response TEXT,
                            rating             TEXT,
                            comment            TEXT,
                            file_name          TEXT,
                            timestamp          TIMESTAMPTZ DEFAULT NOW()
                        )
                        """)
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS message_responses
                        (
                            message_id         TEXT PRIMARY KEY,
                            thread_id          TEXT,
                            user_message       TEXT,
                            assistant_response TEXT,
                            thinking           TEXT,
                            created_at         TIMESTAMPTZ DEFAULT NOW()
                        )
                        """)
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS projects
                        (
                            id                    SERIAL PRIMARY KEY,
                            name                  TEXT NOT NULL,
                            description           TEXT,
                            created_at            TIMESTAMPTZ DEFAULT NOW(),
                            updated_at            TIMESTAMPTZ DEFAULT NOW(),
                            created_by            TEXT REFERENCES users (user_id),
                            status                TEXT        DEFAULT 'active',
                            archived_at           TIMESTAMPTZ,
                            deletion_scheduled_at TIMESTAMPTZ
                        )
                        """)
            cur.execute("""
                DO $$ 
                BEGIN
                    BEGIN
                        ALTER TABLE projects ADD COLUMN status TEXT DEFAULT 'active';
                    EXCEPTION WHEN duplicate_column THEN NULL;
                    END;
                    BEGIN
                        ALTER TABLE projects ADD COLUMN archived_at TIMESTAMPTZ;
                    EXCEPTION WHEN duplicate_column THEN NULL;
                    END;
                    BEGIN
                        ALTER TABLE projects ADD COLUMN deletion_scheduled_at TIMESTAMPTZ;
                    EXCEPTION WHEN duplicate_column THEN NULL;
                    END;
                END $$;
            """)
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS project_members
                        (
                            id         SERIAL PRIMARY KEY,
                            project_id INTEGER REFERENCES projects (id) ON DELETE CASCADE,
                            user_id    TEXT REFERENCES users (user_id),
                            role       TEXT NOT NULL,
                            added_at   TIMESTAMPTZ DEFAULT NOW(),
                            added_by   TEXT REFERENCES users (user_id),
                            UNIQUE (project_id, user_id)
                        )
                        """)
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS project_folders
                        (
                            id               SERIAL PRIMARY KEY,
                            project_id       INTEGER REFERENCES projects (id) ON DELETE CASCADE,
                            parent_folder_id INTEGER REFERENCES project_folders (id) ON DELETE CASCADE,
                            name             TEXT NOT NULL,
                            path             TEXT,
                            created_at       TIMESTAMPTZ DEFAULT NOW(),
                            created_by       TEXT REFERENCES users (user_id),
                            UNIQUE (project_id, parent_folder_id, name)
                        )
                        """)
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS project_files
                        (
                            id            SERIAL PRIMARY KEY,
                            project_id    INTEGER REFERENCES projects (id) ON DELETE CASCADE,
                            folder_id     INTEGER REFERENCES project_folders (id) ON DELETE CASCADE,
                            filename      TEXT NOT NULL,
                            original_name TEXT NOT NULL,
                            file_size     INTEGER,
                            mime_type     TEXT,
                            stored_path   TEXT NOT NULL,
                            version       INTEGER     DEFAULT 1,
                            uploaded_at   TIMESTAMPTZ DEFAULT NOW(),
                            uploaded_by   TEXT REFERENCES users (user_id),
                            comment       TEXT,
                            file_hash     TEXT,
                            UNIQUE (project_id, folder_id, filename)
                        )
                        """)
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS project_file_versions
                        (
                            id          SERIAL PRIMARY KEY,
                            file_id     INTEGER REFERENCES project_files (id) ON DELETE CASCADE,
                            version     INTEGER NOT NULL,
                            stored_path TEXT    NOT NULL,
                            file_size   INTEGER,
                            uploaded_at TIMESTAMPTZ DEFAULT NOW(),
                            uploaded_by TEXT REFERENCES users (user_id),
                            comment     TEXT
                        )
                        """)
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS project_file_comments
                        (
                            id         SERIAL PRIMARY KEY,
                            file_id    INTEGER REFERENCES project_files (id) ON DELETE CASCADE,
                            user_id    TEXT REFERENCES users (user_id),
                            comment    TEXT NOT NULL,
                            created_at TIMESTAMPTZ DEFAULT NOW()
                        )
                        """)
            cur.execute("""
                DO $$ 
                BEGIN
                    BEGIN
                        ALTER TABLE project_files ADD COLUMN content TEXT;
                    EXCEPTION WHEN duplicate_column THEN NULL;
                    END;
                END $$;
            """)
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS project_file_usage
                        (
                            id        SERIAL PRIMARY KEY,
                            file_id   INTEGER REFERENCES project_files (id) ON DELETE CASCADE,
                            user_id   TEXT REFERENCES users (user_id),
                            action    TEXT NOT NULL,
                            details   JSONB,
                            timestamp TIMESTAMPTZ DEFAULT NOW()
                        )
                        """)
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS project_folder_comments
                        (
                            id         SERIAL PRIMARY KEY,
                            folder_id  INTEGER REFERENCES project_folders (id) ON DELETE CASCADE,
                            user_id    TEXT REFERENCES users (user_id),
                            comment    TEXT NOT NULL,
                            created_at TIMESTAMPTZ DEFAULT NOW()
                        )
                        """)
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS task_deposit_items
                        (
                            id                     SERIAL PRIMARY KEY,
                            original_user_id       TEXT REFERENCES users (user_id),
                            original_username      TEXT,
                            project_id             INTEGER REFERENCES projects (id) ON DELETE CASCADE,
                            project_name           TEXT,
                            item_type              TEXT  NOT NULL,
                            item_data              JSONB NOT NULL,
                            stored_path            TEXT,
                            transferred_to_user_id TEXT REFERENCES users (user_id),
                            transferred_at         TIMESTAMPTZ,
                            created_at             TIMESTAMPTZ DEFAULT NOW(),
                            deleted_at             TIMESTAMPTZ
                        )
                        """)
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS task_deposit_permissions
                        (
                            id               SERIAL PRIMARY KEY,
                            project_id       INTEGER REFERENCES projects (id) ON DELETE CASCADE,
                            manager_id       TEXT REFERENCES users (user_id),
                            can_view_deposit BOOLEAN     DEFAULT FALSE,
                            granted_by       TEXT REFERENCES users (user_id),
                            granted_at       TIMESTAMPTZ DEFAULT NOW()
                        )
                        """)
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS checkpoints
                        (
                            thread_id            TEXT NOT NULL,
                            checkpoint_id        TEXT NOT NULL,
                            parent_checkpoint_id TEXT,
                            type                 TEXT,
                            checkpoint           JSONB,
                            metadata             JSONB,
                            PRIMARY KEY (thread_id, checkpoint_id)
                        )
                        """)
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS checkpoint_writes
                        (
                            thread_id     TEXT    NOT NULL,
                            checkpoint_id TEXT    NOT NULL,
                            task_id       TEXT    NOT NULL,
                            idx           INTEGER NOT NULL,
                            value         JSONB,
                            PRIMARY KEY (thread_id, checkpoint_id, task_id, idx)
                        )
                        """)
            cur.execute("""
                CREATE TABLE IF NOT EXISTS recycle_bin (
                    id SERIAL PRIMARY KEY,
                    original_table TEXT NOT NULL,
                    original_id INTEGER NOT NULL,
                    user_id TEXT REFERENCES users(user_id),
                    file_name TEXT,
                    file_content TEXT,
                    file_size INTEGER,
                    original_stored_path TEXT,
                    file_hash TEXT,
                    thread_id TEXT,
                    deleted_at TIMESTAMPTZ DEFAULT NOW(),
                    expires_at TIMESTAMPTZ DEFAULT NOW() + INTERVAL '3 days'
                )
            """)
            cur.execute("""
                CREATE TABLE IF NOT EXISTS project_recycle_bin (
                    id SERIAL PRIMARY KEY,
                    original_table TEXT NOT NULL,
                    original_id INTEGER NOT NULL,
                    project_id INTEGER REFERENCES projects(id) ON DELETE CASCADE,
                    folder_id INTEGER,
                    file_name TEXT,
                    original_name TEXT,
                    file_size INTEGER,
                    stored_path TEXT,
                    file_hash TEXT,
                    version INTEGER,
                    uploaded_by TEXT REFERENCES users(user_id),
                    deleted_at TIMESTAMPTZ DEFAULT NOW(),
                    expires_at TIMESTAMPTZ DEFAULT NOW() + INTERVAL '3 days'
                )
            """)
            cur.execute("""
                CREATE TABLE IF NOT EXISTS project_folders_recycle_bin
                (
                    id                 SERIAL PRIMARY KEY,
                    original_id        INTEGER NOT NULL,
                    project_id         INTEGER REFERENCES projects (id) ON DELETE CASCADE,
                    name               TEXT    NOT NULL,
                    parent_folder_id   INTEGER,
                    original_parent_id INTEGER,
                    created_at         TIMESTAMPTZ,
                    created_by         TEXT,
                    deleted_at         TIMESTAMPTZ DEFAULT NOW(),
                    expires_at         TIMESTAMPTZ DEFAULT NOW() + INTERVAL '3 days'
                )
                """)
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS credit_check_reports
                        (
                            id              SERIAL PRIMARY KEY,
                            user_id         TEXT REFERENCES users (user_id),
                            task_id         TEXT UNIQUE NOT NULL,
                            file_path       TEXT        NOT NULL,
                            companies_count INTEGER     DEFAULT 0,
                            created_at      TIMESTAMPTZ DEFAULT NOW()
                        )
                        """)
            cur.execute("CREATE INDEX IF NOT EXISTS idx_credit_reports_user ON credit_check_reports(user_id)")
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS knowledge_lab_files
                        (
                            id            SERIAL PRIMARY KEY,
                            user_id       TEXT REFERENCES users (user_id) ON DELETE CASCADE,
                            filename      TEXT NOT NULL,
                            original_name TEXT NOT NULL,
                            file_size     INTEGER,
                            content       TEXT, -- extracted text
                            file_hash     TEXT UNIQUE,
                            stored_path   TEXT,
                            uploaded_at   TIMESTAMPTZ DEFAULT NOW(),
                            updated_at    TIMESTAMPTZ DEFAULT NOW()
                        )
                        """)
            cur.execute("CREATE INDEX IF NOT EXISTS idx_lab_user ON knowledge_lab_files(user_id)")
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS company_knowledge_base
                        (
                            id            SERIAL PRIMARY KEY,
                            filename      TEXT NOT NULL,
                            original_name TEXT NOT NULL,
                            file_size     INTEGER,
                            content       TEXT,
                            file_hash     TEXT UNIQUE,
                            stored_path   TEXT,
                            category      TEXT,
                            uploaded_by   TEXT REFERENCES users (user_id),
                            uploaded_at   TIMESTAMPTZ DEFAULT NOW(),
                            updated_at    TIMESTAMPTZ DEFAULT NOW()
                        )
                        """)
            cur.execute("CREATE INDEX IF NOT EXISTS idx_company_kb_category ON company_knowledge_base(category)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_company_kb_filename ON company_knowledge_base(filename)")

            # GIN index using the simple text search configuration
            cur.execute("""
                        CREATE INDEX IF NOT EXISTS idx_company_kb_content
                            ON company_knowledge_base
                                USING gin (to_tsvector('simple', content))
                        """)
            cur.execute("""
                DO $$ 
                BEGIN
                    BEGIN
                        ALTER TABLE recycle_bin ADD COLUMN original_thread_id TEXT;
                    EXCEPTION WHEN duplicate_column THEN NULL;
                    END;
                    BEGIN
                        ALTER TABLE recycle_bin ADD COLUMN deletion_reason TEXT DEFAULT 'manual';
                    EXCEPTION WHEN duplicate_column THEN NULL;
                    END;
                END $$;
            """)
            cur.execute("""
                DO $$ 
                BEGIN
                    BEGIN
                        ALTER TABLE recycle_bin ADD COLUMN uploaded_by TEXT;
                    EXCEPTION WHEN duplicate_column THEN NULL;
                    END;
                    BEGIN
                        ALTER TABLE recycle_bin ADD COLUMN deleted_by TEXT;
                    EXCEPTION WHEN duplicate_column THEN NULL;
                    END;
                END $$;
            """)
            cur.execute("""
                DO $$ 
                BEGIN
                    BEGIN
                        ALTER TABLE recycle_bin DROP COLUMN IF EXISTS original_data;
                    EXCEPTION WHEN undefined_column THEN NULL;
                    END;
                END $$;
            """)
            cur.execute("""
                CREATE TABLE IF NOT EXISTS file_text_cache (
                    id SERIAL PRIMARY KEY,
                    file_path TEXT NOT NULL,      -- original file path as stored (or a unique identifier)
                    last_modified TIMESTAMPTZ,
                    file_hash TEXT NOT NULL,
                    extracted_text TEXT,
                    created_at TIMESTAMPTZ DEFAULT NOW(),
                    updated_at TIMESTAMPTZ DEFAULT NOW()
                )
            """)
            cur.execute("CREATE UNIQUE INDEX IF NOT EXISTS idx_file_text_cache_path ON file_text_cache(file_path)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_file_text_cache_hash ON file_text_cache(file_hash)")
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS admin_audit_log
                        (
                            id             SERIAL PRIMARY KEY,
                            admin_user_id  TEXT REFERENCES users (user_id),
                            admin_username TEXT,
                            action         TEXT NOT NULL,
                            table_name     TEXT NOT NULL,
                            row_id         TEXT,
                            column_name    TEXT,
                            old_value      TEXT,
                            new_value      TEXT,
                            ip_address     TEXT,
                            success        BOOLEAN     DEFAULT TRUE,
                            error_message  TEXT,
                            created_at     TIMESTAMPTZ DEFAULT NOW()
                        )
                        """)
            cur.execute("CREATE INDEX IF NOT EXISTS idx_admin_audit_log_admin ON admin_audit_log(admin_user_id)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_admin_audit_log_created ON admin_audit_log(created_at)")
            cur.execute("""
                        CREATE TABLE IF NOT EXISTS file_analysis
                        (
                            id                SERIAL PRIMARY KEY,
                            file_hash         TEXT NOT NULL,
                            file_type         TEXT NOT NULL,
                            original_filename TEXT,
                            user_id           TEXT REFERENCES users (user_id),
                            thread_id         TEXT,
                            project_id        INTEGER,
                            extracted_text    TEXT,
                            usage_count       INTEGER     DEFAULT 0,
                            last_used_at      TIMESTAMPTZ,
                            deleted_at        TIMESTAMPTZ,
                            created_at        TIMESTAMPTZ DEFAULT NOW()
                        )
                        """)

            # Add missing columns (idempotent)
            cur.execute("""
                DO $$
                BEGIN
                    BEGIN
                        ALTER TABLE file_analysis ADD COLUMN IF NOT EXISTS file_size INTEGER;
                    EXCEPTION WHEN duplicate_column THEN NULL;
                    END;
                END $$;
            """)

            # Add unique constraint after ensuring column exists
            cur.execute("""
                DO $$
                BEGIN
                    IF NOT EXISTS (
                        SELECT 1 FROM pg_constraint WHERE conname = 'file_analysis_unique_hash_type_size'
                    ) THEN
                        ALTER TABLE file_analysis ADD CONSTRAINT file_analysis_unique_hash_type_size UNIQUE (file_hash, file_type, file_size);
                    END IF;
                END $$;
            """)

            # Indexes
            cur.execute("CREATE INDEX IF NOT EXISTS idx_chat_messages_thread_id_timestamp ON chat_messages(thread_id, timestamp)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_user_files_expires_at ON user_files(expires_at)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_user_files_user_id ON user_files(user_id)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_file_usage_user_filename ON file_usage(user_id, filename)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_message_responses_created_at ON message_responses(created_at)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_project_members_user ON project_members(user_id)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_project_folders_parent ON project_folders(parent_folder_id)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_project_files_folder ON project_files(folder_id)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_project_files_hash ON project_files(file_hash)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_task_deposit_items_original_user ON task_deposit_items(original_user_id)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_task_deposit_items_project ON task_deposit_items(project_id)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_recycle_bin_user_id ON recycle_bin(user_id)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_recycle_bin_expires_at ON recycle_bin(expires_at)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_project_recycle_bin_project_id ON project_recycle_bin(project_id)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_project_recycle_bin_expires_at ON project_recycle_bin(expires_at)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_project_folders_recycle_bin_project ON project_folders_recycle_bin(project_id)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_recycle_bin_original_thread_id ON recycle_bin(original_thread_id)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_file_analysis_hash_type ON file_analysis(file_hash, file_type)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_file_analysis_user ON file_analysis(user_id)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_user_files_hash ON user_files(user_id, file_hash)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_user_files_thread_id ON user_files(user_id, thread_id)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_file_usage_thread_filename ON file_usage(thread_id, filename)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_chat_sessions_user_id ON chat_sessions(user_id)")
            conn.commit()
            logger.info("PostgreSQL tables initialized.")

# User ID helpers
def get_user_id():
    if session.get('consent_value', 0) == 1:
        if 'user_id' in session:
            return session['user_id']
        session['user_id'] = str(uuid.uuid4())
        return session['user_id']
    else:
        if 'temp_user_id' not in session:
            session['temp_user_id'] = str(uuid.uuid4())
        return session['temp_user_id']

def ensure_user_exists(user_id):
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("INSERT INTO users (user_id) VALUES (%s) ON CONFLICT DO NOTHING", (user_id,))
            conn.commit()

# Anonymous session storage
def get_anon_history_path(thread_id):
    user_id = get_user_id()
    temp_dir = get_anon_temp_dir(user_id)
    return os.path.join(temp_dir, f"{thread_id}_history.json")

def get_session_messages_anon(thread_id):
    path = get_anon_history_path(thread_id)
    if not os.path.exists(path):
        return []
    try:
        from filelock import FileLock
        lock_path = path + ".lock"
        with FileLock(lock_path, timeout=5):
            with open(path, 'r', encoding='utf-8') as f:
                return json.load(f)
    except ImportError:
        logger.warning("filelock not installed. Anonymous file reads may have race conditions.")
        # No lock
        try:
            with open(path, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception as e:
            logger.error(f"Failed to read anon history {thread_id}: {e}")
            return []
    except Exception as e:
        logger.error(f"Failed to read anon history {thread_id}: {e}")
        return []

def store_message_anon(thread_id, role, content, thinking=None):
    path = get_anon_history_path(thread_id)
    try:
        from filelock import FileLock
        lock_path = path + ".lock"
        with FileLock(lock_path, timeout=5):
            history = []
            if os.path.exists(path):
                with open(path, 'r', encoding='utf-8') as f:
                    history = json.load(f)
            history.append({
                "role": role,
                "content": content,
                "thinking": thinking,
                "timestamp": beijing_now()
            })
            with open(path, 'w', encoding='utf-8') as f:
                json.dump(history, f, ensure_ascii=False, indent=2)
    except ImportError:
        logger.warning("filelock not installed. Anonymous file writes may have race conditions.")
        # No lock
        history = []
        if os.path.exists(path):
            try:
                with open(path, 'r', encoding='utf-8') as f:
                    history = json.load(f)
            except Exception:
                pass
        history.append({
            "role": role,
            "content": content,
            "thinking": thinking,
            "timestamp": beijing_now()
        })
        with open(path, 'w', encoding='utf-8') as f:
            json.dump(history, f, ensure_ascii=False, indent=2)
    except Exception as e:
        logger.error(f"Failed to write anon history {thread_id}: {e}")

def get_or_create_session(thread_id, title=None):
    if session.get('consent_value', 0) != 1:
        return
    user_id = get_user_id()
    ensure_user_exists(user_id)
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT 1 FROM chat_sessions WHERE thread_id = %s", (thread_id,))
            if not cur.fetchone():
                cur.execute(
                    "INSERT INTO chat_sessions (user_id, thread_id, title, created_at, updated_at) VALUES (%s, %s, %s, %s, %s)",
                    (user_id, thread_id, title or "新对话", utc_now(), utc_now())
                )
                conn.commit()

def generate_session_title(messages, max_len=20):
    for msg in messages:
        if msg.get('role') == 'user':
            content = msg.get('content', '').strip()
            if content:
                title = content[:max_len]
                if len(content) > max_len:
                    title += '...'
                return title
    return '新对话'

def update_session_title(thread_id, title):
    if session.get('consent_value', 0) != 1:
        return
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(
                "UPDATE chat_sessions SET title = %s, updated_at = %s WHERE thread_id = %s",
                (title, utc_now(), thread_id)
            )
            conn.commit()

def store_message(thread_id, role, content, thinking=None):
    if session.get('consent_value', 0) != 1:
        store_message_anon(thread_id, role, content, thinking)
        return None
    with get_db_connection() as conn:
        with db_transaction(conn):
            with conn.cursor() as cur:
                cur.execute(
                    "INSERT INTO chat_messages (thread_id, role, content, thinking, timestamp) VALUES (%s, %s, %s, %s, %s) RETURNING id",
                    (thread_id, role, content, thinking, utc_now())
                )
                msg_id = cur.fetchone()[0]
                cur.execute(
                    "UPDATE chat_sessions SET updated_at = %s WHERE thread_id = %s",
                    (utc_now(), thread_id)
                )
    messages = get_session_messages(thread_id)
    if len(messages) == 2:
        new_title = generate_session_title(messages)
        update_session_title(thread_id, new_title)
    return msg_id

def get_session_messages(thread_id):
    if session.get('consent_value', 0) != 1:
        return get_session_messages_anon(thread_id)
    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            db_execute_readonly(cur)
            cur.execute(
                "SELECT role, content, thinking, timestamp FROM chat_messages WHERE thread_id = %s ORDER BY id ASC",
                (thread_id,)
            )
            rows = cur.fetchall()
            messages = []
            for row in rows:
                ts_utc = row['timestamp']
                ts_beijing = ts_utc.astimezone(BEIJING_TZ).strftime('%Y-%m-%d %H:%M:%S') if ts_utc else None
                messages.append({
                    "role": row['role'],
                    "content": row['content'],
                    "thinking": row['thinking'],
                    "timestamp": ts_beijing
                })
            return messages

def get_user_sessions():
    if session.get('consent_value', 0) != 1:
        return []
    user_id = get_user_id()
    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            db_execute_readonly(cur)
            cur.execute(
                "SELECT thread_id, title, created_at, updated_at FROM chat_sessions WHERE user_id = %s ORDER BY updated_at DESC",
                (user_id,)
            )
            rows = cur.fetchall()
            sessions = []
            for row in rows:
                sessions.append({
                    "thread_id": row['thread_id'],
                    "title": row['title'],
                    "created_at": row['created_at'].astimezone(BEIJING_TZ).strftime('%Y-%m-%d %H:%M:%S') if row['created_at'] else None,
                    "updated_at": row['updated_at'].astimezone(BEIJING_TZ).strftime('%Y-%m-%d %H:%M:%S') if row['updated_at'] else None
                })
            return sessions

def delete_session(thread_id):
    try:
        with get_db_connection() as conn:
            with db_transaction(conn):
                with conn.cursor(cursor_factory=RealDictCursor) as cur:
                    cur.execute("SELECT user_id FROM chat_sessions WHERE thread_id = %s", (thread_id,))
                    row = cur.fetchone()
                    if not row:
                        return
                    user_id = row['user_id']

                    cur.execute("SELECT id, filename, content, size_bytes, original_stored_path, file_hash, thread_id FROM user_files WHERE thread_id = %s", (thread_id,))
                    files = cur.fetchall()
                    for f in files:
                        cur.execute("""
                            INSERT INTO recycle_bin 
                            (original_table, original_id, user_id, file_name, file_content, file_size, original_stored_path, file_hash, thread_id, original_thread_id, deletion_reason, deleted_at, expires_at, uploaded_by, deleted_by)
                            VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, NOW(), NOW() + INTERVAL '3 days', %s, %s)
                        """, ('user_files', f['id'], user_id, f['filename'], f['content'], f['size_bytes'], f['original_stored_path'], f['file_hash'], f['thread_id'], thread_id, 'chat_deleted', user_id, user_id))

                    cur.execute("DELETE FROM user_files WHERE thread_id = %s", (thread_id,))
                    cur.execute("DELETE FROM chat_messages WHERE thread_id = %s", (thread_id,))
                    cur.execute("DELETE FROM file_usage WHERE thread_id = %s", (thread_id,))
                    cur.execute("DELETE FROM feedback WHERE thread_id = %s", (thread_id,))
                    cur.execute("DELETE FROM consent WHERE thread_id = %s", (thread_id,))
                    cur.execute("DELETE FROM chat_sessions WHERE thread_id = %s", (thread_id,))
        logger.info(f"Deleted session {thread_id} and moved {len(files)} files to recycle bin")
        file_cache_manager.clear_thread(thread_id)
    except Exception as e:
        logger.error(f"Failed to delete session {thread_id}: {e}", exc_info=True)
        raise

def archive_session(thread_id, user_id, reason="manual"):
    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("SELECT archive_path FROM archived_sessions WHERE thread_id = %s", (thread_id,))
            if cur.fetchone():
                return None
            cur.execute("SELECT title, created_at, updated_at FROM chat_sessions WHERE thread_id = %s", (thread_id,))
            sess_row = cur.fetchone()
            if not sess_row:
                return None
            title = sess_row['title']
            created_at = sess_row['created_at']
            updated_at = sess_row['updated_at']
            cur.execute("SELECT role, content, thinking, timestamp FROM chat_messages WHERE thread_id = %s ORDER BY timestamp", (thread_id,))
            messages = []
            for row in cur.fetchall():
                messages.append({
                    "role": row['role'],
                    "content": row['content'],
                    "thinking": row['thinking'],
                    "timestamp": row['timestamp']
                })
            cur.execute("SELECT user_message, assistant_response, rating, comment, timestamp FROM feedback WHERE thread_id = %s", (thread_id,))
            feedbacks = []
            for row in cur.fetchall():
                feedbacks.append({
                    "user_message": row['user_message'],
                    "assistant_response": row['assistant_response'],
                    "rating": row['rating'],
                    "comment": row['comment'],
                    "timestamp": row['timestamp']
                })
            cur.execute("SELECT consent_given, timestamp FROM consent WHERE thread_id = %s", (thread_id,))
            consent_row = cur.fetchone()
            consent = {"consent_given": consent_row['consent_given'], "timestamp": consent_row['timestamp']} if consent_row else None
            archive_date = datetime.now().strftime("%Y-%m-%d")
            dump_dir = str(DUMP_DIR / user_id / f"{user_id}_{archive_date}")
            os.makedirs(dump_dir, exist_ok=True)
            if not os.path.isdir(dump_dir):
                logger.error(f"Failed to create archive directory: {dump_dir}")
                return None
            session_info = {
                "thread_id": thread_id,
                "user_id": user_id,
                "title": title,
                "created_at": created_at.isoformat() if created_at else None,
                "updated_at": updated_at.isoformat() if updated_at else None,
                "archived_at": datetime.now().isoformat(),
                "reason": reason
            }
            try:
                with open(os.path.join(dump_dir, f"{thread_id}_session.json"), "w", encoding="utf-8") as f:
                    json.dump(session_info, f, ensure_ascii=False, indent=2, default=str)
                with open(os.path.join(dump_dir, f"{thread_id}_messages.json"), "w", encoding="utf-8") as f:
                    json.dump(messages, f, ensure_ascii=False, indent=2, default=str)
                if feedbacks:
                    with open(os.path.join(dump_dir, f"{thread_id}_feedback.json"), "w", encoding="utf-8") as f:
                        json.dump(feedbacks, f, ensure_ascii=False, indent=2, default=str)
                if consent:
                    with open(os.path.join(dump_dir, f"{thread_id}_consent.json"), "w", encoding="utf-8") as f:
                        json.dump(consent, f, ensure_ascii=False, indent=2, default=str)
            except Exception as e:
                logger.error(f"Failed to write archive files for thread {thread_id}: {e}")
                return None
            archive_path = os.path.join(dump_dir, f"{thread_id}_session.json")
            cur.execute("INSERT INTO archived_sessions (thread_id, user_id, archive_path) VALUES (%s, %s, %s)", (thread_id, user_id, archive_path))
            conn.commit()
            logger.info(f"Archived session {thread_id} for user {user_id} to {dump_dir}")
            return dump_dir

def cleanup_old_sessions(days=15):
    cutoff = utc_now() - timedelta(days=days)
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT thread_id, user_id FROM chat_sessions WHERE updated_at < %s", (cutoff,))
            old = cur.fetchall()
            for thread_id, user_id in old:
                archive_session(thread_id, user_id, reason="auto_15days")
                delete_session(thread_id)

def cleanup_stale_message_responses(hours=1):
    cutoff = utc_now() - timedelta(hours=hours)
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("DELETE FROM message_responses WHERE created_at < %s AND (assistant_response = '' OR assistant_response IS NULL)", (cutoff,))
            conn.commit()
            logger.info(f"Deleted stale message response placeholders older than {hours} hours.")

def get_cached_image_description(file_hash):
    if session.get('consent_value', 0) != 1:
        return None
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT description FROM image_description_cache WHERE file_hash = %s", (file_hash,))
            row = cur.fetchone()
            if row:
                return row[0]
    return None

def cache_image_description(file_hash, description):
    if session.get('consent_value', 0) != 1:
        return
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("""
                INSERT INTO image_description_cache (file_hash, description)
                VALUES (%s, %s)
                ON CONFLICT (file_hash) DO UPDATE
                SET description = EXCLUDED.description, created_at = NOW()
            """, (file_hash, description))
            conn.commit()

# Rate limiter for admin routes
admin_rate_limit = {}
ADMIN_RATE_LIMIT = 5          # max attempts
ADMIN_RATE_WINDOW = 30 * 60   # 15 minutes in seconds

def admin_rate_limiter(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        # Get client IP
        ip = request.remote_addr
        if ip is None:
            ip = 'unknown'
        now = time.time()
        # Clean old entries
        for key in list(admin_rate_limit.keys()):
            if admin_rate_limit[key]['timestamp'] < now - ADMIN_RATE_WINDOW:
                del admin_rate_limit[key]
        # Check limit
        key = f"{ip}:admin_action"
        if key in admin_rate_limit:
            if admin_rate_limit[key]['count'] >= ADMIN_RATE_LIMIT:
                logger.warning(f"Rate limit exceeded for admin action from IP {ip}")
                return jsonify({"error": "Too many attempts. Please try again later."}), 429
        else:
            admin_rate_limit[key] = {'count': 0, 'timestamp': now}
        # Increment count and call the function
        admin_rate_limit[key]['count'] += 1
        return f(*args, **kwargs)
    return decorated_function

def log_admin_action(admin_user_id, admin_username, action, table_name, row_id=None,
                     column_name=None, old_value=None, new_value=None,
                     success=True, error_message=None):
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("""
                INSERT INTO admin_audit_log
                (admin_user_id, admin_username, action, table_name, row_id, column_name,
                 old_value, new_value, ip_address, success, error_message)
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
            """, (
                admin_user_id, admin_username, action, table_name, row_id,
                column_name, old_value, new_value, request.remote_addr,
                success, error_message
            ))
            conn.commit()

# Return a generic error string
def safe_error_response(user_message="处理文件时出错，请检查文件格式或稍后重试。", log_error=None):
    if log_error:
        logger.error(log_error, exc_info=True)
    return f"[错误] {user_message}"

# Tools
@tool(description="Get current date and time in Beijing time (UTC+8).")
def get_date() -> str:
    return datetime.now(BEIJING_TZ).strftime("%Y-%m-%d %H:%M:%S")

BOCHA_API_KEY = os.getenv("BOCHA_API_KEY")
BOCHA_URL = "https://api.bochaai.com/v1/web-search"

@tool(description="Search the web using Bocha. Use for up‑to‑date information.")
def bocha_search(query: str) -> str:
    headers = {"Authorization": f"Bearer {BOCHA_API_KEY}", "Content-Type": "application/json"}
    payload = json.dumps({"query": query, "summary": True, "freshness": "noLimit", "count": 10})
    try:
        response = requests.post(BOCHA_URL, headers=headers, data=payload, timeout=10)
        response.raise_for_status()
        data = response.json()
        webpages = data.get('data', {}).get('webPages', {}).get('value', [])
        if not webpages:
            return "No search results found."
        formatted = []
        for idx, page in enumerate(webpages[:10], 1):
            title = page.get('name', 'No title')
            snippet = page.get('snippet', 'No snippet')
            date = page.get('datePublished', 'Unknown date')
            url = page.get('url', 'No URL')
            formatted.append(f"{idx}. **{title}**\n   Published: {date}\n   Summary: {snippet}\n   Source: {url}\n")
        return "\n".join(formatted)
    except Exception as e:
        return f"Search failed: {str(e)}"

class VLModel:
    def __init__(self):
        self.api_key = os.getenv("DASHSCOPE_API_KEY") or os.getenv("QWEN_API_KEY")
        if not self.api_key:
            self.api_key = os.getenv("DEEPSEEK_API_KEY")  # fallback, but won't work for VL
        self.base_url = "https://dashscope.aliyuncs.com/compatible-mode/v1"
        self.model_name = "qwen3-vl-plus-2025-12-19"
        self.client = None
        self.max_image_size = 1024  # resize longer side to this
        self._init_client()

    def _init_client(self):
        try:
            from openai import OpenAI
            self.client = OpenAI(api_key=self.api_key, base_url=self.base_url)
            logger.info(f"VL client initialized with model {self.model_name}")
        except ImportError:
            logger.error("OpenAI package not installed. VL model disabled.")
            self.client = None
        except Exception as e:
            logger.error(f"VL client init failed: {e}")
            self.client = None

    def is_available(self):
        return self.client is not None and self.api_key is not None

    def _preprocess_image(self, image_bytes):
        """Resize image to reduce token consumption."""
        from PIL import Image
        import io
        try:
            img = Image.open(io.BytesIO(image_bytes))
            # Convert to RGB if necessary (e.g., RGBA, P)
            if img.mode in ('RGBA', 'P'):
                img = img.convert('RGB')
            # Resize while keeping aspect ratio
            w, h = img.size
            if max(w, h) > self.max_image_size:
                scale = self.max_image_size / max(w, h)
                new_size = (int(w * scale), int(h * scale))
                img = img.resize(new_size, Image.Resampling.LANCZOS)
            # Encode as JPEG (smaller than PNG)
            buffer = io.BytesIO()
            img.save(buffer, format='JPEG', quality=85)
            return buffer.getvalue()
        except Exception as e:
            logger.warning(f"Image preprocessing failed: {e}, using original")
            return image_bytes

    def encode_image_to_base64(self, image_bytes):
        processed = self._preprocess_image(image_bytes)
        return base64.b64encode(processed).decode('utf-8')

    def describe_image(self, image_bytes, prompt="请描述这张图片的内容"):
        if not self.is_available():
            return "⚠️ VL模型不可用，请检查API密钥。"

        try:
            base64_image = self.encode_image_to_base64(image_bytes)
            response = self.client.chat.completions.create(
                model=self.model_name,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}},
                        {"type": "text", "text": prompt}
                    ]
                }],
                max_tokens=800,
                temperature=0.7
            )
            description = response.choices[0].message.content
            if not description:
                return "⚠️ 未获得图片描述，请稍后重试。"
            return description
        except Exception as e:
            error_msg = str(e)
            logger.error(f"VL image description failed: {error_msg}")
            if "InvalidModel" in error_msg or "model not found" in error_msg:
                return f"⚠️ 模型 {self.model_name} 不可用，请检查模型名称或API密钥。"
            elif "rate limit" in error_msg.lower():
                return "⚠️ 请求过于频繁，请稍后再试。"
            elif "content policy" in error_msg.lower():
                return "⚠️ 图片内容不符合安全规范，无法描述。"
            else:
                return f"⚠️ 图片描述失败: {error_msg[:100]}"

    def describe_pdf_page(self, image_bytes, page_num):
        return self.describe_image(
            image_bytes,
            f"请详细描述这个PDF页面(第{page_num}页)的内容，包括标题、段落、表格、图表等关键信息。"
        )

vl_model = VLModel()

def describe_images_in_file(file_bytes, filename, page_texts=None):
    if not vl_model.is_available():
        return ""
    ext = os.path.splitext(filename)[1].lower()
    descriptions = []
    try:
        if ext == '.pdf':
            doc = None
            try:
                doc = fitz.open(stream=BytesIO(file_bytes), filetype="pdf")
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    img_list = page.get_images(full=True)
                    for img_idx, img in enumerate(img_list):
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        prompt = f"Describe this image from page {page_num+1} of the PDF. Include any charts, diagrams, tables, or visual information."
                        description = vl_model.describe_image(image_bytes, prompt)
                        if description and not description.startswith("⚠️ VL模型不可用"):
                            descriptions.append(f"[Image on page {page_num+1}, image {img_idx+1}]: {description}")
            finally:
                if doc:
                    doc.close()
        elif ext in ['.docx', '.docm']:
            import docx
            doc = docx.Document(BytesIO(file_bytes))
            img_counter = 1
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_blob = rel.target_part.blob
                        description = vl_model.describe_image(image_blob, "Describe this image from the Word document.")
                        if description and not description.startswith("⚠️ VL模型不可用"):
                            descriptions.append(f"[Image {img_counter}]: {description}")
                        img_counter += 1
                    except Exception:
                        pass
        elif ext in ['.pptx', '.pptm']:
            from pptx import Presentation
            prs = Presentation(BytesIO(file_bytes))
            slide_num = 1
            img_counter = 1
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == 13:
                        try:
                            image_bytes = shape.image.blob
                            description = vl_model.describe_image(image_bytes, f"Describe this image from slide {slide_num} of the PowerPoint.")
                            if description and not description.startswith("⚠️ VL模型不可用"):
                                descriptions.append(f"[Image on slide {slide_num}, image {img_counter}]: {description}")
                            img_counter += 1
                        except Exception:
                            pass
                slide_num += 1
    except Exception as e:
        logger.error(f"Error extracting images from {filename}: {e}")
    return "\n".join(descriptions)

# Agent
def get_agent(max_tokens=None):
    global _agent, _current_max_tokens
    if max_tokens is None:
        max_tokens = session.get('max_tokens', 1600)
    if _agent is not None and _current_max_tokens == max_tokens:
        return _agent
    with _agent_lock:
        if _agent is not None and _current_max_tokens == max_tokens:
            return _agent
        api_key = os.getenv("DEEPSEEK_API_KEY")
        if not api_key:
            api_key = os.getenv("QWEN_API_KEY") or os.getenv("DASHSCOPE_API_KEY")
        if not api_key:
            raise ValueError("Missing DEEPSEEK_API_KEY or QWEN_API_KEY")
        os.environ["DASHSCOPE_API_KEY"] = api_key
        os.environ["DASHSCOPE_API_BASE"] = "https://api.deepseek.com/v1"
        llm = ChatDeepSeek(
            model="deepseek-v4-pro",
            api_key=api_key,
            temperature=0.7,
            max_tokens=max_tokens,
            streaming=False,
            extra_body={"thinking": {"type": "disabled"}},
        )
        # at top of get_agent, after with _agent_lock:
        if _async_checkpointer is None:
            with _async_checkpointer_lock:  # need a new lock
                if _async_checkpointer is None:
                    _init_async_checkpointer()
        system_prompt = AGENT_SYSTEM_PROMPT
        _agent = create_agent(
            model=llm,
            tools=[get_date, bocha_search],
            system_prompt=system_prompt,
            checkpointer=_async_checkpointer
        )
        _current_max_tokens = max_tokens
        logger.info(f"Agent reinitialized with DeepSeek model, max_tokens={max_tokens}")
        return _agent

# OCR Manager
class OCRManager:
    _instance = None
    _lock = threading.Lock()

    def __new__(cls):
        if cls._instance is None:
            with cls._lock:
                if cls._instance is None:
                    cls._instance = super().__new__(cls)
                    cls._instance._initialized = False
        return cls._instance

    def __init__(self):
        if self._initialized:
            return
        self._initialized = True
        self.reader = None
        self.engine_name = None
        self._init_ocr()

    def _init_ocr(self):
        try:
            from paddleocr import PaddleOCR
            try:
                self.reader = PaddleOCR(use_textline_orientation=True, lang='ch')
            except TypeError:
                self.reader = PaddleOCR(use_angle_cls=True, lang='ch')
            self.engine_name = "PaddleOCR"
            logger.info("PaddleOCR initialized successfully.")
        except ImportError:
            logger.warning("PaddleOCR not installed. Will try EasyOCR.")
        except Exception as e:
            logger.warning(f"PaddleOCR init failed: {e}. Will try EasyOCR.")
        if self.reader is None:
            try:
                import easyocr
                self.reader = easyocr.Reader(['ch_sim', 'en'], gpu=False)
                self.engine_name = "EasyOCR"
                logger.info("EasyOCR initialized as fallback.")
            except ImportError:
                logger.error("No OCR engine available. Install 'paddleocr' or 'easyocr'.")
                self.reader = None
            except Exception as e:
                logger.error(f"EasyOCR init failed: {e}")
                self.reader = None

    def is_available(self):
        return self.reader is not None

    def run_ocr(self, image_np):
        if self.reader is None:
            return ""
        try:
            if self.engine_name == "PaddleOCR":
                result = self.reader.ocr(image_np, cls=True)
                if result and result[0]:
                    return "\n".join([line[1][0] for line in result[0]])
            elif self.engine_name == "EasyOCR":
                result = self.reader.readtext(image_np, detail=0, paragraph=True)
                if result:
                    return "\n".join(result)
            return ""
        except Exception as e:
            logger.error(f"OCR run error: {e}")
            return ""

ocr_manager = OCRManager()
run_ocr = ocr_manager.run_ocr

# Text Preprocessing & Similarity
TECH_STD_PATTERNS = [
    r'GB/T\s*\d+\.?\d*', r'GB\s*\d+\.?\d*', r'ISO\s*\d+', r'IEC\s*\d+',
    r'IEEE\s*\d+', r'DIN\s*\d+', r'BS\s*\d+', r'EN\s*\d+', r'ASME\s*\d+',
    r'API\s*\d+', r'ASTM\s*\d+', r'JJG\s*\d+', r'JB/T\s*\d+', r'HG/T\s*\d+',
    r'SY/T\s*\d+', r'DL/T\s*\d+', r'NB/T\s*\d+', r'SH/T\s*\d+', r'YS/T\s*\d+',
    r'FZ/T\s*\d+', r'QB/T\s*\d+', r'CJ/T\s*\d+', r'JG/T\s*\d+', r'GA/T\s*\d+',
    r'HS/T\s*\d+', r'行业标准', r'国家标准', r'技术规范'
]
SENSITIVE_PATTERNS = [
    r'(公司|集团|有限|股份|组织|委员会|协会|研究院|大学|学院)',
    r'(北京|上海|广州|深圳|杭州|南京|武汉|成都|重庆|天津|西安)',
    r'(项目|工程|系统|平台|软件|硬件|方案)',
    r'(张|王|李|刘|陈|杨|赵|黄|周|吴|徐|孙|马|朱|胡|林|郭|何|高)',
    r'(一等奖|二等奖|三等奖|金奖|银奖|优秀奖)',
    r'\d{17}[\dXx]',
    r'1[3-9]\d{9}',
    r'\d{18}',
    r'证书编号[：:]\s*\w+',
]

def preprocess_text_for_similarity(text):
    if not text:
        return ""
    text = re.sub(r'[^\w\u4e00-\u9fff\s]', '', text)
    words = text.split()
    filtered = [w for w in words if len(w) >= 6]
    text = ' '.join(filtered)
    for pat in TECH_STD_PATTERNS:
        text = re.sub(pat, '', text, flags=re.IGNORECASE)
    text = re.sub(r'^目录|^第[一二三四五六七八九十]+章', '', text, flags=re.MULTILINE)
    for pat in SENSITIVE_PATTERNS:
        text = re.sub(pat, '', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return text

def remove_template_content(text, template_text, threshold=0.85):
    if not template_text or not text:
        return text
    paras = [p.strip() for p in text.split('\n') if len(p.strip()) > 10]
    template_paras = [p.strip() for p in template_text.split('\n') if len(p.strip()) > 10]
    if not paras or not template_paras:
        return text
    all_paras = paras + template_paras
    vectorizer = TfidfVectorizer(stop_words=None, lowercase=True).fit(all_paras)
    vecs = vectorizer.transform(all_paras)
    para_vecs = vecs[:len(paras)]
    template_vecs = vecs[len(paras):]
    sim_matrix = cosine_similarity(para_vecs, template_vecs)
    keep_mask = np.max(sim_matrix, axis=1) < threshold
    kept_paras = [p for i, p in enumerate(paras) if keep_mask[i]]
    if not kept_paras:
        return "[模板内容已全部匹配，未保留任何原文] " + text
    return '\n'.join(kept_paras)

def extract_keywords(text, top_k=20):
    if not text.strip():
        return []
    vectorizer = TfidfVectorizer(stop_words=None, max_features=top_k)
    try:
        tfidf = vectorizer.fit_transform([text])
        feature_names = vectorizer.get_feature_names_out()
        scores = tfidf.toarray()[0]
        keyword_score = sorted(zip(feature_names, scores), key=lambda x: x[1], reverse=True)
        return [kw for kw, _ in keyword_score[:top_k]]
    except Exception:
        return []

def keyword_overlap_similarity(text1, text2):
    kw1 = set(extract_keywords(text1, 20))
    kw2 = set(extract_keywords(text2, 20))
    if not kw1 and not kw2:
        return 0.0
    inter = len(kw1 & kw2)
    union = len(kw1 | kw2)
    return inter / union if union > 0 else 0.0

def compute_similarity_with_numbers(text1, text2, template_text=None):
    clean1 = preprocess_text_for_similarity(text1)
    clean2 = preprocess_text_for_similarity(text2)
    if template_text:
        clean1 = remove_template_content(clean1, template_text)
        clean2 = remove_template_content(clean2, template_text)
    if not clean1.strip() or not clean2.strip():
        return 0.0, text1, text2, []
    vectorizer = TfidfVectorizer(stop_words=None, lowercase=True)
    tfidf = vectorizer.fit_transform([clean1, clean2])
    sim = cosine_similarity(tfidf[0:1], tfidf[1:2])[0][0]
    escaped1 = html.escape(text1)
    escaped2 = html.escape(text2)
    matcher = difflib.SequenceMatcher(None, escaped1, escaped2)
    matching_blocks = matcher.get_matching_blocks()
    segments1 = []
    last_idx = 0
    match_counter = 1
    blocks_detail = []
    for block in matching_blocks:
        i, j, size = block
        if size == 0 or size <= 6:
            continue
        blocks_detail.append({
            "id": match_counter,
            "pos1": i,
            "pos2": j,
            "size": size,
            "text1_snippet": escaped1[i:i + min(size, 100)] + ("..." if size > 100 else ""),
            "text2_snippet": escaped2[j:j + min(size, 100)] + ("..." if size > 100 else "")
        })
        if i > last_idx:
            segments1.append(('text', escaped1[last_idx:i]))
        match_text = escaped1[i:i + size]
        color_class = 'match-highlight-long' if size > 100 else 'match-highlight-short'
        marker = f"<sup><small>[{match_counter}]</small></sup> "
        segments1.append(('match', match_text, marker, color_class))
        last_idx = i + size
        match_counter += 1
    if last_idx < len(escaped1):
        segments1.append(('text', escaped1[last_idx:]))
    segments2 = []
    last_idx = 0
    match_counter = 1
    for block in matching_blocks:
        i, j, size = block
        if size == 0 or size <= 6:
            continue
        if j > last_idx:
            segments2.append(('text', escaped2[last_idx:j]))
        match_text = escaped2[j:j + size]
        color_class = 'match-highlight-long' if size > 100 else 'match-highlight-short'
        marker = f"<sup><small>[{match_counter}]</small></sup> "
        segments2.append(('match', match_text, marker, color_class))
        last_idx = j + size
        match_counter += 1
    if last_idx < len(escaped2):
        segments2.append(('text', escaped2[last_idx:]))

    def build_html(segments):
        parts = []
        for seg in segments:
            if seg[0] == 'text':
                parts.append(seg[1])
            else:
                _, text, marker, color_class = seg
                parts.append(marker + f'<span class="{color_class}">{text}</span>')
        return ''.join(parts)

    html1 = build_html(segments1)
    html2 = build_html(segments2)
    return sim, html1, html2, blocks_detail

def compute_batch_semantic_similarity(texts):
    model = get_semantic_model()
    if model is None:
        return None
    n = len(texts)
    if n < 2:
        return [[0.0]*n for _ in range(n)]
    try:
        embeddings = model.encode(texts, batch_size=32, show_progress_bar=False)
        sim_matrix = cosine_similarity(embeddings)
        return sim_matrix.tolist()
    except Exception as e:
        logger.error(f"Error computing semantic similarity: {e}")
        return None

def file_attr_similarity(meta1, meta2):
    score = 0.0
    if meta1.get('author') and meta2.get('author') and meta1['author'] == meta2['author']:
        score += 50
    try:
        if meta1.get('creationDate') and meta2.get('creationDate'):
            date1 = re.sub(r'D:', '', meta1['creationDate'])[:14]
            date2 = re.sub(r'D:', '', meta2['creationDate'])[:14]
            if date1 == date2:
                score += 30
    except Exception:
        pass
    name1 = meta1.get('filename', '')
    name2 = meta2.get('filename', '')
    if name1 and name2:
        common = len(set(name1.lower()) & set(name2.lower()))
        total = len(set(name1.lower()) | set(name2.lower()))
        if total > 0:
            score += (common / total) * 20
    return min(score, 100.0)

def extract_text_from_doc_crossplatform(file_bytes):
    with tempfile.NamedTemporaryFile(suffix='.doc', delete=False) as f:
        f.write(file_bytes)
        temp_doc = f.name
    try:
        # Try antiword
        try:
            result = subprocess.run(['antiword', temp_doc], capture_output=True, timeout=30)
            if result.returncode == 0 and result.stdout:
                for enc in ('gbk', 'gb18030', 'utf-8'):
                    try:
                        decoded = result.stdout.decode(enc)
                        if decoded.strip():
                            return decoded
                    except UnicodeDecodeError:
                        continue
                return result.stdout.decode('utf-8', errors='replace')
        except FileNotFoundError:
            logger.warning("antiword not installed")
        except Exception as e:
            logger.warning(f"antiword error: {e}")

        # Try catdoc
        try:
            result = subprocess.run(['catdoc', '-a', temp_doc], capture_output=True, timeout=30)
            if result.returncode == 0 and result.stdout:
                for enc in ('gbk', 'gb18030', 'utf-8'):
                    try:
                        decoded = result.stdout.decode(enc)
                        if decoded.strip():
                            return decoded
                    except UnicodeDecodeError:
                        continue
                return result.stdout.decode('utf-8', errors='replace')
        except FileNotFoundError:
            logger.warning("catdoc not installed")
        except Exception as e:
            logger.warning(f"catdoc error: {e}")

        return None
    finally:
        if os.path.exists(temp_doc):
            os.unlink(temp_doc)

def extract_text_from_doc(file_bytes):
    """Extract text from .doc using Win32 COM (requires Word installed)."""
    import pythoncom
    import win32com.client as win32
    import tempfile
    import os
    import time

    pythoncom.CoInitialize()
    temp_doc = None
    temp_txt = None
    word = None
    doc = None
    try:
        fd, temp_doc = tempfile.mkstemp(suffix='.doc')
        os.close(fd)
        with open(temp_doc, 'wb') as f:
            f.write(file_bytes)

        fd, temp_txt = tempfile.mkstemp(suffix='.txt')
        os.close(fd)

        word = win32.Dispatch("Word.Application")
        word.Visible = False
        word.DisplayAlerts = 0
        doc = word.Documents.Open(temp_doc, AddToRecentFiles=False, Visible=False)
        doc.SaveAs2(temp_txt, FileFormat=2)
        doc.Close(SaveChanges=False)
        word.Quit(SaveChanges=False)

        time.sleep(0.1)

        for enc in ('gbk', 'gb18030', 'utf-8', 'gb2312', 'latin-1'):
            try:
                with open(temp_txt, 'r', encoding=enc) as f:
                    text = f.read()
                if text.strip():
                    logger.info(f"Successfully decoded .doc with {enc}")
                    return text
            except UnicodeDecodeError:
                continue
        with open(temp_txt, 'r', encoding='utf-8', errors='replace') as f:
            text = f.read()
        return text
    except Exception as e:
        logger.error(f"win32com .doc extraction failed: {e}")
        return None
    finally:
        if doc:
            try:
                doc.Close(False)
            except Exception:
                pass
        if word:
            try:
                word.Quit(False)
            except Exception:
                pass
        if temp_doc and os.path.exists(temp_doc):
            try:
                os.unlink(temp_doc)
            except Exception:
                pass
        if temp_txt and os.path.exists(temp_txt):
            try:
                os.unlink(temp_txt)
            except Exception:
                pass
        pythoncom.CoUninitialize()

def detect_excel_format(file_bytes):
    if len(file_bytes) < 8:
        return None
    if file_bytes[:8] == b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1':
        return 'xls'
    if file_bytes[:2] == b'PK':
        return 'xlsx'
    return None

def extract_text_from_xls(file_bytes):
    try:
        import xlrd
        workbook = xlrd.open_workbook(file_contents=file_bytes)
        text_parts = []
        for sheet in workbook.sheets():
            sheet_text = []
            for row in range(sheet.nrows):
                row_text = " ".join(str(cell.value) for cell in sheet.row(row) if cell.value)
                if row_text.strip():
                    sheet_text.append(row_text)
            if sheet_text:
                text_parts.append(f"--- Sheet: {sheet.name} ---\n" + "\n".join(sheet_text))
        return "\n\n".join(text_parts) if text_parts else "[No text in Excel]"
    except Exception as e:
        return f"[Excel parsing error (old format): {e}]"

def detect_word_format(file_bytes):
    if len(file_bytes) < 8:
        return None
    if file_bytes[:8] == b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1':
        return 'doc'
    if file_bytes[:2] == b'PK':
        return 'docx'
    return None

def extract_images_from_file(file_storage):
    images = []
    ext = os.path.splitext(file_storage.filename)[1].lower()
    file_bytes = file_storage.read()
    file_storage.seek(0)
    if ext == '.pdf':
        doc = None
        try:
            doc = fitz.open(stream=BytesIO(file_bytes), filetype="pdf")
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                img_list = page.get_images(full=True)
                for img in img_list:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image = Image.open(BytesIO(image_bytes))
                    images.append(image)
        except Exception:
            pass
        finally:
            if doc:
                doc.close()
    return images

def image_similarity(images1, images2):
    if not images1 or not images2:
        return 0.0
    max_sim = 0.0
    for img1 in images1:
        hash1 = imagehash.phash(img1)
        for img2 in images2:
            hash2 = imagehash.phash(img2)
            sim = 1 - (hash1 - hash2) / 64.0
            max_sim = max(max_sim, sim)
    return max_sim * 100

def extract_metadata(file_storage):
    meta = {'filename': file_storage.filename}
    ext = os.path.splitext(file_storage.filename)[1].lower()
    file_bytes = file_storage.read()
    file_storage.seek(0)
    if ext == '.pdf':
        try:
            doc = fitz.open(stream=BytesIO(file_bytes), filetype="pdf")
            info = doc.metadata
            meta['author'] = info.get('author', '')
            meta['creator'] = info.get('creator', '')
            meta['producer'] = info.get('producer', '')
            meta['creationDate'] = info.get('creationDate', '')
        except Exception:
            pass
    elif ext in ['.docx', '.docm']:
        try:
            doc = docx.Document(BytesIO(file_bytes))
            core_props = doc.core_properties
            meta['author'] = core_props.author or ''
            meta['created'] = core_props.created
            meta['modified'] = core_props.modified
        except Exception:
            pass
    return meta

def truncate_filename(filename, max_len=40):
    if len(filename) <= max_len:
        return filename
    name, ext = os.path.splitext(filename)
    if len(ext) > 10:
        ext = ext[:10]
    available = max_len - len(ext) - 3
    if available < 1:
        ext = ext[:max_len]
        return ext
    truncated_name = name[:available] + '...'
    return truncated_name + ext

def detect_file_real_type(file_bytes: bytes) -> str:
    if len(file_bytes) < 8:
        return None
    if file_bytes[:8] == b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1':
        return 'doc'
    if file_bytes[:4] == b'PK\x03\x04':
        return 'docx'
    return None

# Advanced extraction using fallbacks
KREUZBERG_AVAILABLE = False
UNSTRUCTURED_AVAILABLE = False
KREUZBERG_SIZE_LIMIT = 2 * 1024 * 1024

def extract_text_advanced(file_bytes: bytes, filename: str, file_size: int) -> str:
    ext = os.path.splitext(filename)[1].lower()
    office_extensions = {'.docx', '.xlsx', '.pptx', '.xlsm', '.xlsb', '.pptm', '.doc'}
    if ext not in office_extensions:
        return None

    # Try MarkItDown first
    try:
        from markitdown import MarkItDown
        md = MarkItDown()
        result = md.convert_stream(BytesIO(file_bytes), file_extension=ext.lstrip('.'))
        text = result.text_content
        if text and text.strip():
            logger.info(f"MarkItDown extracted {len(text)} chars from {filename}")
            return text
    except Exception as e:
        logger.warning(f"MarkItDown failed for {filename}: {e}")

    # For .doc, try legacy
    if ext == '.doc':
        text = extract_text_from_doc_crossplatform(file_bytes)
        if text and text.strip():
            return text
        text = extract_text_from_doc(file_bytes)
        if text and text.strip():
            return text

    # For .docx, fallback to python-docx
    if ext == '.docx':
        try:
            doc = docx.Document(BytesIO(file_bytes))
            text = "\n".join([para.text for para in doc.paragraphs])
            for table in doc.tables:
                for row in table.rows:
                    row_text = "\t".join([cell.text for cell in row.cells])
                    text += "\n" + row_text
            if text.strip():
                return text
        except Exception as e:
            logger.warning(f"python-docx fallback failed for {filename}: {e}")

    return None

def extract_text_from_file(file_storage):
    filename = file_storage.filename
    if not filename:
        return None, {}

    file_bytes = file_storage.read()
    file_storage.seek(0)
    file_hash = hashlib.sha256(file_bytes).hexdigest()

    ext = os.path.splitext(filename)[1].lower()
    wps_map = {'.wps': '.doc', '.et': '.xls', '.dps': '.ppt'}
    original_ext = ext
    if ext in wps_map:
        ext = wps_map[ext]

    text = None
    page_texts = {}

    if ext in ['.txt', '.md', '.text', '.csv', '.json']:
        try:
            text = file_bytes.decode('utf-8')
        except UnicodeDecodeError:
            text = file_bytes.decode('utf-8', errors='replace')
        page_texts = {1: text}

    elif ext in ['.html', '.htm']:
        try:
            from bs4 import BeautifulSoup
            html_text = file_bytes.decode('utf-8', errors='replace')
            soup = BeautifulSoup(html_text, 'html.parser')
            text = soup.get_text(separator='\n', strip=True)
            if not text or not text.strip():
                text = html_text
            page_texts = {1: text}
        except Exception as e:
            logger.warning(f"HTML parsing failed, falling back to raw text: {e}")
            text = file_bytes.decode('utf-8', errors='replace')
            page_texts = {1: text}

    elif ext == '.pdf':
        try:
            doc = fitz.open(stream=BytesIO(file_bytes), filetype="pdf")
            full_text = []
            page_texts = {}
            has_text = False
            try:
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    page_text = page.get_text()
                    if page_text and page_text.strip():
                        has_text = True
                        full_text.append(page_text)
                        page_texts[page_num + 1] = page_text
                extracted = "\n".join(full_text).strip()
                if has_text and len(extracted) > 50:
                    text = extracted
                else:
                    logger.info("PDF appears to be scanned (no text). Starting OCR...")
                    if not ocr_manager.is_available():
                        if vl_model.is_available():
                            logger.info("Using VL model for scanned PDF")
                            extracted = ""
                            for page_num in range(len(doc)):
                                page = doc.load_page(page_num)
                                pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
                                img_bytes = pix.tobytes("png")
                                description = vl_model.describe_pdf_page(img_bytes, page_num + 1)
                                extracted += f"\n\n--- 第{page_num + 1}页 (VL分析) ---\n{description}"
                            page_texts = {i + 1: "" for i in range(len(doc))}
                            text = extracted
                        else:
                            text = "[无法提取PDF文本，且OCR/VL不可用]"
                    else:
                        ocr_results = []
                        ocr_page_texts = {}
                        zoom = 2.0
                        mat = fitz.Matrix(zoom, zoom)
                        for page_num in range(len(doc)):
                            page = doc.load_page(page_num)
                            pix = page.get_pixmap(matrix=mat, alpha=False)
                            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                            max_dim = 2000
                            if max(img.size) > max_dim:
                                ratio = max_dim / max(img.size)
                                new_size = (int(img.size[0] * ratio), int(img.size[1] * ratio))
                                img = img.resize(new_size, Image.Resampling.LANCZOS)
                            img_np = np.array(img)
                            page_text = run_ocr(img_np)
                            if page_text:
                                ocr_results.append(page_text)
                                ocr_page_texts[page_num + 1] = page_text
                            else:
                                ocr_results.append("")
                                ocr_page_texts[page_num + 1] = ""
                        if any(t.strip() for t in ocr_results):
                            extracted = "\n\n".join(ocr_results)
                            page_texts = ocr_page_texts
                            text = extracted
                        else:
                            if vl_model.is_available():
                                extracted = ""
                                for page_num in range(len(doc)):
                                    page = doc.load_page(page_num)
                                    pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
                                    img_bytes = pix.tobytes("png")
                                    description = vl_model.describe_pdf_page(img_bytes, page_num + 1)
                                    extracted += f"\n\n--- 第{page_num + 1}页 (VL分析) ---\n{description}"
                                page_texts = {i + 1: "" for i in range(len(doc))}
                                text = extracted
                            else:
                                text = "[No text detected in PDF even after OCR and VL not available]"
            finally:
                doc.close()
        except Exception as e:
            logger.error(f"PDF error: {e}", exc_info=True)
            text = safe_error_response("无法解析PDF文件，请确保文件未损坏。", log_error=e)

    elif ext in ['.docx', '.docm', '.dotx', '.dotm']:
        real_format = detect_word_format(file_bytes)
        if real_format == 'doc':
            text = extract_text_from_doc(file_bytes)
            if text:
                page_texts = {1: text}
            else:
                text = "[无法从 .doc 文件中提取文本。请转换为 .docx 格式后重试，或安装 antiword/catdoc。]"
        else:
            try:
                def extract_docx_text(byte_data):
                    doc = docx.Document(BytesIO(byte_data))
                    full_text = "\n".join([para.text for para in doc.paragraphs])
                    for table in doc.tables:
                        for row in table.rows:
                            row_text = "\t".join([cell.text for cell in row.cells])
                            full_text += "\n" + row_text
                    del doc
                    return full_text if full_text.strip() else "[No text in Word document]"

                text = extract_docx_text(file_bytes)
                page_texts = {1: text}
            except Exception as e:
                logger.error(f"DOCX parsing error: {e}")
                text = safe_error_response("无法解析Word文档，请转换为DOCX格式或检查文件。", log_error=e)

    elif ext in ['.xlsx', '.xlsm', '.xltx', '.xltm', '.xlsb']:
        wb = None
        try:
            wb = openpyxl.load_workbook(BytesIO(file_bytes), read_only=True, data_only=True)
            text_parts = []
            for sheet in wb.worksheets:
                sheet_text = []
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        sheet_text.append(row_text)
                if sheet_text:
                    text_parts.append(f"--- Sheet: {sheet.title} ---\n" + "\n".join(sheet_text))
            full_text = "\n\n".join(text_parts) if text_parts else "[No text in Excel]"
            text = full_text
            page_texts = {1: full_text}
        except Exception as e:
            logger.warning(f"openpyxl failed for {filename}: {e}. Trying fallback methods.")
            real_format = detect_excel_format(file_bytes)
            if real_format == 'xls':
                text = extract_text_from_xls(file_bytes)
                if text and not text.startswith("["):
                    page_texts = {1: text}
                else:
                    file_storage.seek(0)
                    md = MarkItDown()
                    result = md.convert(BytesIO(file_bytes), file_extension=original_ext.lstrip('.'))
                    text = result.text_content
                    if text and text.strip():
                        page_texts = {1: text}
                    else:
                        text = safe_error_response("无法解析Excel文档，请转换为xlsx格式或检查文件。", log_error=e)
            else:
                file_storage.seek(0)
                md = MarkItDown()
                result = md.convert(BytesIO(file_bytes), file_extension=original_ext.lstrip('.'))
                text = result.text_content
                if text and text.strip():
                    page_texts = {1: text}
                else:
                    text = safe_error_response("无法解析Excel文档，请转换为xlsx格式或检查文件。", log_error=e)
        finally:
            if wb:
                wb.close()

    elif ext in ['.pptx', '.pptm', '.potx', '.ppsx']:
        prs = None
        try:
            prs = Presentation(BytesIO(file_bytes))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            full_text = "\n".join(text_runs)
            text = full_text if full_text.strip() else "[No text in PowerPoint]"
            page_texts = {1: full_text}
        except Exception as e:
            logger.error(f"PPTX parsing error: {e}")
            text = safe_error_response("无法解析PowerPoint文件，请检查文件。", log_error=e)
        finally:
            if prs is not None:
                del prs

    elif ext == '.xls':
        xls = None
        try:
            xls = pd.ExcelFile(BytesIO(file_bytes), engine='xlrd')
            text_parts = []
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                sheet_text = df.to_string(index=False, header=True)
                if sheet_text.strip():
                    text_parts.append(f"--- Sheet: {sheet_name} ---\n{sheet_text}")
            full_text = "\n\n".join(text_parts) if text_parts else "[No text in Excel]"
            text = full_text
            page_texts = {1: full_text}
        except Exception as e:
            logger.error(f"XLS parsing error: {e}")
            text = f"[Excel parsing failed: {str(e)}]"
        finally:
            if xls is not None:
                try:
                    xls.close()
                except Exception:
                    pass
    elif ext == '.doc':
        text = extract_text_from_doc(file_bytes)
        if text:
            page_texts = {1: text}
        else:
            text = "[无法从 .doc 文件中提取文本。请转换为 .docx 格式后重试。]"
    elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
        if ocr_manager.is_available():
            try:
                image = Image.open(BytesIO(file_bytes))
                if image.mode != 'RGB':
                    image = image.convert('RGB')
                max_dim = 2000
                if max(image.size) > max_dim:
                    ratio = max_dim / max(image.size)
                    new_size = (int(image.size[0] * ratio), int(image.size[1] * ratio))
                    image = image.resize(new_size, Image.Resampling.LANCZOS)
                img_np = np.array(image)
                extracted_text = run_ocr(img_np)
                if extracted_text:
                    text = extracted_text
                    page_texts = {1: extracted_text}
            except Exception as e:
                logger.warning(f"OCR failed for image: {e}")
        if not text and vl_model.is_available():
            logger.info("Using VL model for image description")
            description = vl_model.describe_image(file_bytes)
            text = description
            page_texts = {1: description}
        if not text:
            text = safe_error_response("无法从图片提取文本，请确保图片清晰或使用其他格式。", log_error=e)
    else:
        try:
            file_storage.seek(0)
            md = MarkItDown()
            result = md.convert(BytesIO(file_bytes), file_extension=original_ext.lstrip('.'))
            text = result.text_content
            if text and text.strip():
                page_texts = {1: text}
            else:
                text = "[No text extracted by MarkItDown]"
        except Exception as e:
            logger.error(f"MarkItDown parsing failed for {original_ext}: {e}")
            text = safe_error_response(f"不支持的文件格式: {original_ext}", log_error=e)

    if not text or text.startswith("["):
        return text, page_texts

    analyze_images = session.get('analyze_images', True)
    if analyze_images:
        cached_desc = get_cached_image_description(file_hash)
        if cached_desc:
            image_desc = cached_desc
        else:
            image_desc = describe_images_in_file(file_bytes, filename, page_texts if page_texts else None)
            if image_desc:
                cache_image_description(file_hash, image_desc)
        if image_desc:
            text += "\n\n--- Image Descriptions ---\n" + image_desc

    def clean_report_headers(text):
        if not text:
            return text
        lines = text.split('\n')
        cleaned_lines = []
        skip = False
        for line in lines:
            stripped = line.strip()
            # Check if this line is a header to skip
            is_header = any(stripped.startswith(prefix) for prefix in [
                '--- Sheet:', '技术标规律性分析检查结果', '标段名称：', '投标单位个数：',
                '检查结果：', '检查规则：', '相似度计算说明：', '一、标书围串风险分析结果',
                '二、分析结果详情', '签字：', '日期：', '技术标规律性分析详情'
            ])
            if is_header:
                skip = True
                continue
            if skip:
                # Only stop skipping when we hit a blank line (section separator)
                if stripped == '':
                    skip = False
                continue
            cleaned_lines.append(line)
        return '\n'.join(cleaned_lines)

    if text:
        text = clean_report_headers(text)

    return text, page_texts

def get_or_extract_file_analysis(file_storage, file_type, user_id, thread_id=None, project_id=None):
    filename = file_storage.filename
    file_bytes = file_storage.read()
    file_storage.seek(0)
    file_hash = hashlib.sha256(file_bytes).hexdigest()
    file_size = len(file_bytes)

    # Composite key: hash + size (add mtime if needed, but size is enough to avoid most collisions)
    cache_key = f"{file_hash}_{file_size}"

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            # Check existing analysis
            cur.execute("""
                SELECT id, extracted_text FROM file_analysis
                WHERE file_hash = %s AND file_type = %s AND file_size = %s AND deleted_at IS NULL
            """, (file_hash, file_type, file_size))
            row = cur.fetchone()
            if row:
                analysis_id, extracted_text = row
                # If cached text is invalid marker, force re‑extraction
                if extracted_text == "[INVALID_EXTRACTION]":
                    logger.warning(f"Cached extraction for {filename} marked as invalid. Re‑extracting.")
                    extracted_text, _ = extract_text_from_file(file_storage)
                    if not extracted_text or extracted_text.startswith("["):
                        extracted_text = ""
                    # Validate again
                    if not is_valid_extracted_text(extracted_text):
                        extracted_text = "[INVALID_EXTRACTION]"
                    # Update cache
                    cur.execute("""
                        UPDATE file_analysis
                        SET extracted_text = %s, usage_count = usage_count + 1, last_used_at = NOW()
                        WHERE id = %s
                    """, (extracted_text, analysis_id))
                    conn.commit()
                else:
                    # Normal valid cache – update usage count
                    cur.execute("""
                        UPDATE file_analysis
                        SET usage_count = usage_count + 1, last_used_at = NOW()
                        WHERE id = %s
                    """, (analysis_id,))
                    conn.commit()
                return extracted_text if extracted_text != "[INVALID_EXTRACTION]" else ""
            else:
                # No cache – extract and validate
                extracted_text, _ = extract_text_from_file(file_storage)
                if not extracted_text or extracted_text.startswith("["):
                    extracted_text = ""
                if not is_valid_extracted_text(extracted_text):
                    extracted_text = "[INVALID_EXTRACTION]"
                # Store in cache
                cur.execute("""
                    INSERT INTO file_analysis (file_hash, file_type, file_size, original_filename, user_id, thread_id, project_id, extracted_text, usage_count, last_used_at)
                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s, 1, NOW())
                    RETURNING id
                """, (file_hash, file_type, file_size, filename, user_id, thread_id, project_id, extracted_text))
                analysis_id = cur.fetchone()[0]
                conn.commit()

                # Also update user_files.content if thread_id provided
                if thread_id:
                    cur.execute("""
                        SELECT id FROM user_files
                        WHERE user_id = %s AND file_hash = %s AND thread_id = %s
                        ORDER BY created_at DESC LIMIT 1
                    """, (user_id, file_hash, thread_id))
                    uf_row = cur.fetchone()
                    if uf_row:
                        cur.execute("""
                            UPDATE user_files SET content = %s WHERE id = %s
                        """, (extracted_text if extracted_text != "[INVALID_EXTRACTION]" else "", uf_row[0]))
                        conn.commit()

                return extracted_text if extracted_text != "[INVALID_EXTRACTION]" else ""

# Batch comparison helper functions
def _precompute_tfidf_for_files(file_data, template_text=None):
    texts = []
    for fd in file_data:
        clean = preprocess_text_for_similarity(fd['text'])
        if template_text:
            clean = remove_template_content(clean, template_text)
        texts.append(clean)
    vectorizer = TfidfVectorizer(stop_words=None, lowercase=True)
    tfidf_matrix = vectorizer.fit_transform(texts)
    return vectorizer, tfidf_matrix

def _compute_pair_similarity_from_matrix(tfidf_matrix, i, j):
    sim = cosine_similarity(tfidf_matrix[i:i + 1], tfidf_matrix[j:j + 1])[0][0]
    return sim

def store_batch_comparison_temp(data):
    fd, path = tempfile.mkstemp(suffix='.json', prefix='comp_', text=True)
    with os.fdopen(fd, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, default=str)
    return path

def load_batch_comparison_temp(path):
    with open(path, 'r', encoding='utf-8') as f:
        return json.load(f)

download_tokens = defaultdict(int)
download_tokens_lock = threading.Lock()

# File Cache Manager
class FileCacheManager:
    def __init__(self, max_cached_files=10, max_content_size=50 * 1024):
        self._lock = RLock()
        self.caches = {}
        self.recent = {}
        self.max_cached = max_cached_files
        self.max_size = max_content_size

    def add(self, thread_id, filename, content, user_id):
        if content is None:
            content = ''
        with self._lock:
            if len(content) > self.max_size:
                content = content[:self.max_size] + "\n[内容已截断，仅保留前50KB]"
            cache = self.caches.setdefault(thread_id, {})
            recent_list = self.recent.setdefault(thread_id, [])
            cache[filename] = content
            if filename in recent_list:
                recent_list.remove(filename)
            recent_list.insert(0, filename)
            while len(recent_list) > self.max_cached:
                old = recent_list.pop()
                del cache[old]

    def load_from_db(self, thread_id, user_id):
        with self._lock:
            if session.get('consent_value', 0) != 1:
                self.caches[thread_id] = {}
                self.recent[thread_id] = []
                return
            with get_db_connection() as conn:
                with conn.cursor() as cur:
                    cur.execute(
                        "SELECT filename, content FROM user_files WHERE thread_id = %s AND user_id = %s AND (expires_at IS NULL OR expires_at > NOW())",
                        (thread_id, user_id)
                    )
                    rows = cur.fetchall()
                    if rows:
                        cache = {}
                        recent_list = []
                        for filename, content in rows:
                            if content is None:
                                content = ''
                            cache[filename] = content
                            recent_list.append(filename)
                        self.caches[thread_id] = cache
                        self.recent[thread_id] = recent_list
                    else:
                        self.caches[thread_id] = {}
                        self.recent[thread_id] = []

    def get_recent_with_lock(self, thread_id):
        with self._lock:
            return self.recent.get(thread_id, []).copy()

    def get_content(self, thread_id, filename):
        with self._lock:
            return self.caches.get(thread_id, {}).get(filename)

    def clear_thread(self, thread_id):
        with self._lock:
            self.caches.pop(thread_id, None)
            self.recent.pop(thread_id, None)

    def evict_oldest(self, max_threads=20):
        """Evict oldest thread entries if cache exceeds max_threads."""
        with self._lock:
            while len(self.caches) > max_threads:
                oldest = list(self.caches.keys())[0]
                self.caches.pop(oldest, None)
                self.recent.pop(oldest, None)

    def add_thread(self, thread_id):
        """Track thread access order; auto-evict oldest if over limit."""
        with self._lock:
            self.evict_oldest()

file_cache_manager = FileCacheManager()

def add_to_cache(thread_id, filename, content, user_id):
    if content is None:
        content = ''
    file_cache_manager.add(thread_id, filename, content, user_id)

def load_cache_from_db(thread_id, user_id):
    file_cache_manager.load_from_db(thread_id, user_id)
    file_cache_manager.evict_oldest()

def get_user_total_storage_size(user_id):
    if session.get('consent_value', 0) != 1:
        return 0
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(
                "SELECT COALESCE(SUM(size_bytes), 0) FROM user_files WHERE user_id = %s AND (expires_at IS NULL OR expires_at > NOW())",
                (user_id,))
            return cur.fetchone()[0]

def record_file_usage(thread_id, filename, usage_type, question_text=None):
    if session.get('consent_value', 0) != 1:
        return
    user_id = get_user_id()
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(
                "INSERT INTO file_usage (user_id, thread_id, filename, usage_type, question) VALUES (%s, %s, %s, %s, %s)",
                (user_id, thread_id, filename, usage_type, question_text)
            )
            conn.commit()

# Flask app instance - injected by the factory via app.routes.register_all()
# All @app.route decorators reference this module-level variable.
app = None
csrf = None
ADMIN_PASSWORD_HASH = os.getenv("ADMIN_PASSWORD_HASH")
if not ADMIN_PASSWORD_HASH and os.getenv("ADMIN_PSWD"):
    ADMIN_PASSWORD_HASH = generate_password_hash(os.getenv("ADMIN_PSWD"))
    logger.warning("ADMIN_PASSWORD_HASH not set, using plaintext ADMIN_PSWD.")

def split_thinking_answer(text):
    patterns = [r'【思考】(.*?)【回答】', r'思考：(.*?)回答：', r'<思考>(.*?)</思考>']
    for pat in patterns:
        match = re.search(pat, text, re.DOTALL)
        if match:
            thinking = match.group(1).strip()
            answer = re.sub(pat, '', text, flags=re.DOTALL).strip()
            return thinking, answer
    return None, text

# Task locking
user_active_tasks = {}
user_task_lock = RLock()
TASK_TIMEOUT_SECONDS = 600

def cleanup_stale_tasks():
    with user_task_lock:
        now = datetime.now()
        stale = [uid for uid, info in user_active_tasks.items() if
                 (now - info['start_time']).total_seconds() > TASK_TIMEOUT_SECONDS]
        for uid in stale:
            logger.warning(f"Cleaning stale task lock for user {uid}")
            del user_active_tasks[uid]

def get_chat_short_name(thread_id):
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT title FROM chat_sessions WHERE thread_id = %s", (thread_id,))
            row = cur.fetchone()
            if row and row[0]:
                name = row[0]
                return name if len(name) <= 20 else name[:17] + '...'
            return "新对话"

def acquire_task_lock(user_id, thread_id, task_type):
    with user_task_lock:
        cleanup_stale_tasks()
        if user_id in user_active_tasks:
            busy = user_active_tasks[user_id]
            return False, busy['thread_id'], get_chat_short_name(busy['thread_id'])
        else:
            user_active_tasks[user_id] = {'thread_id': thread_id, 'task_type': task_type, 'start_time': datetime.now()}
            return True, None, None

def release_task_lock(user_id):
    with user_task_lock:
        if user_id in user_active_tasks:
            del user_active_tasks[user_id]

ALLOWED_EXTENSIONS = {'.txt', '.md', '.text', '.csv', '.pdf', '.docx', '.docm', '.dotx', '.dotm', '.doc',
                      '.xlsx', '.xlsm', '.xltx', '.xltm', '.xlsb', '.xls', '.pptx', '.pptm', '.potx', '.ppsx',
                      '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.wps', '.et', '.dps', '.webp',
                      '.html', '.htm', '.json'}

def allowed_file(filename):
    return os.path.splitext(filename)[1].lower() in ALLOWED_EXTENSIONS

# ---------- Routes ----------
@app.route('/')
def index():
    if 'consent_value' not in session:
        session['consent_value'] = 0
    if 'thread_id' not in session:
        session['thread_id'] = str(uuid.uuid4())
        get_or_create_session(session['thread_id'])
    if 'chat_history' not in session:
        session['chat_history'] = get_session_messages(session['thread_id'])
    user_id = get_user_id()
    load_cache_from_db(session['thread_id'], user_id)
    return render_template('index.html',
                           consent_given=(session.get('consent_value', 0) == 1),
                           chat_history=session['chat_history'],
                           recent_files=file_cache_manager.get_recent_with_lock(session['thread_id']))

@app.route('/get_csrf_token', methods=['GET'])
def get_csrf_token():
    return jsonify({'csrf_token': generate_csrf()})

@app.route('/consent', methods=['POST'])
def set_consent():
    data = request.get_json()
    choice = data.get('consent', False)
    session['consent_value'] = 1 if choice else 0
    if session['consent_value'] == 1:
        with get_db_connection() as conn:
            with conn.cursor() as cur:
                cur.execute(
                    "INSERT INTO consent (thread_id, consent_given, timestamp) VALUES (%s, %s, %s) ON CONFLICT (thread_id) DO UPDATE SET consent_given = EXCLUDED.consent_given, timestamp = EXCLUDED.timestamp",
                    (session['thread_id'], session['consent_value'], utc_now())
                )
                conn.commit()
    return jsonify({"status": "ok"})

@app.route('/logout', methods=['POST'])
def logout():
    session.clear()
    session['consent_value'] = 0
    session['thread_id'] = str(uuid.uuid4())
    get_or_create_session(session['thread_id'])
    return jsonify({"status": "ok"})

@app.route('/favicon.ico')
def favicon():
    favicon_path = os.path.join(os.getcwd(), 'static', 'favicon.ico')
    if os.path.isfile(favicon_path):
        return send_file(favicon_path, mimetype='image/vnd.microsoft.icon')
    # Return empty response to avoid 500 error
    return '', 204

@app.route('/send', methods=['POST'])
def send_message():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403

    user_msg = request.form.get('message', '').strip()
    message_id = request.form.get('message_id')
    if not user_msg and 'files' not in request.files:
        return jsonify({"error": "Empty message and no files"}), 400
    if not message_id:
        return jsonify({"error": "Missing message_id"}), 400

    thread_id = session['thread_id']
    user_id = get_user_id()
    get_or_create_session(thread_id)

    # Idempotency check
    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            try:
                cur.execute("""
                    INSERT INTO message_responses (message_id, thread_id, user_message, assistant_response, thinking)
                    VALUES (%s, %s, %s, '', '')
                    ON CONFLICT (message_id) DO NOTHING
                    RETURNING assistant_response, thinking
                """, (message_id, thread_id, user_msg))
                row = cur.fetchone()
                if row and row['assistant_response'] == '':
                    conn.commit()
                elif row:
                    return jsonify({
                        "user_message": user_msg,
                        "assistant_message": row['assistant_response'],
                        "thinking": row['thinking'],
                        "cached": True
                    })
            except Exception:
                conn.rollback()
                cur.execute("SELECT assistant_response, thinking FROM message_responses WHERE message_id = %s", (message_id,))
                row = cur.fetchone()
                if row:
                    return jsonify({
                        "user_message": user_msg,
                        "assistant_message": row['assistant_response'],
                        "thinking": row['thinking'],
                        "cached": True
                    })

    uploaded_files = request.files.getlist('files')
    has_files = len(uploaded_files) > 0 and uploaded_files[0].filename
    file_contents = []
    is_image = False
    image_analysis_used = True

    if has_files:
        for f in uploaded_files:
            if not allowed_file(f.filename):
                return jsonify({"error": f"不支持的文件类型: {f.filename}"}), 400
        success, busy_thread, busy_name = acquire_task_lock(user_id, thread_id, 'ocr_upload')
        if not success:
            return jsonify({
                "error": "resource_busy",
                "busy_chat": busy_name,
                "message": f"另一个资源密集型任务正在聊天“{busy_name}”中进行，请稍后再试。"
            }), 409
    else:
        success = True

    try:
        uploaded_filenames = []
        if has_files:
            for uploaded in uploaded_files:
                if not uploaded.filename:
                    continue
                uploaded_filenames.append(uploaded.filename)
                ext = os.path.splitext(uploaded.filename)[1].lower()
                if ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff']:
                    is_image = True

                if session.get('consent_value', 0) == 1:
                    ensure_user_exists(user_id)
                    file_bytes = uploaded.read()
                    file_hash = hashlib.sha256(file_bytes).hexdigest()
                    uploaded.seek(0)
                    with get_db_connection() as conn:
                        with conn.cursor() as cur:
                            cur.execute("SELECT id FROM user_files WHERE user_id = %s AND file_hash = %s", (user_id, file_hash))
                            existing = cur.fetchone()
                            if not existing:
                                unique_name = f"{file_hash}_{int(time.time())}{ext}"
                                original_dir = os.path.join(USER_FILES_ORIGINAL_ROOT, user_id)
                                os.makedirs(original_dir, exist_ok=True)
                                original_path = os.path.join(original_dir, unique_name)
                                uploaded.save(original_path)
                                cur.execute("""
                                    INSERT INTO user_files (user_id, thread_id, filename, size_bytes, expires_at,
                                                            original_stored_path, file_hash, original_expires_at, original_name, content)
                                    VALUES (%s, %s, %s, %s, NULL, %s, %s, NOW() + INTERVAL '3 days', %s, %s)
                                    ON CONFLICT (thread_id, filename) DO UPDATE SET
                                        size_bytes = EXCLUDED.size_bytes,
                                        original_stored_path = EXCLUDED.original_stored_path,
                                        file_hash = EXCLUDED.file_hash,
                                        original_expires_at = EXCLUDED.original_expires_at,
                                        original_name = EXCLUDED.original_name,
                                        content = EXCLUDED.content
                                """, (user_id, thread_id, uploaded.filename, len(file_bytes), original_path, file_hash, uploaded.filename, ""))
                                conn.commit()
                    # Use robust cache
                    file_content = get_or_extract_file_analysis(uploaded, 'chat', user_id, thread_id=thread_id)
                else:
                    file_content, _ = extract_text_from_file(uploaded)

                if file_content and not file_content.startswith("["):
                    # Also validate again (the robust cache already did, but double‑check)
                    if is_valid_extracted_text(file_content):
                        add_to_cache(thread_id, uploaded.filename, file_content, user_id)
                        record_file_usage(thread_id, uploaded.filename, 'chat', user_msg)
                        file_contents.append(file_content)
                    else:
                        logger.warning(f"Extracted text from {uploaded.filename} is invalid, skipping.")
                        file_contents.append(f"[文件 {uploaded.filename} 的内容无法读取，请检查文件格式。]")
                else:
                    # No text extracted – inform user
                    file_contents.append(f"[文件 {uploaded.filename} 的内容无法读取，请检查文件格式。]")

        # Build uploaded file text
        file_text = ""
        if has_files:
            if file_contents:
                combined = "\n\n".join(file_contents)
                if is_image:
                    file_text = f"The user uploaded an image. Extracted description:\n{combined}\n\nUser query: {user_msg}"
                else:
                    file_text = f"File content(s):\n{combined}\n\nUser query:\n{user_msg}"
            else:
                file_text = f"The user uploaded a file but no readable text could be extracted. The user's question is: {user_msg}"
        else:
            # Include recent batch compare files in context if available
            batch_files = session.get('batch_compare_files')
            if batch_files:
                batch_context = "The user previously ran a batch comparison on the following files:\n"
                for bf in batch_files:
                    batch_context += f"\n--- File: {bf['filename']} ---\n{bf['text']}\n"
                batch_context += f"\n\nUser query:\n{user_msg}"
                file_text = batch_context
            else:
                file_text = user_msg

        # Process knowledge files (same as before, but we also need to handle invalid extraction for knowledge files? For now, assume they are valid)
        knowledge_files_json = request.form.get('knowledge_files')
        knowledge_files = []
        if knowledge_files_json:
            try:
                knowledge_files = json.loads(knowledge_files_json)
            except Exception:
                pass

        knowledge_content = []
        for kf in knowledge_files:
            source = kf.get('source')
            fid = kf.get('id')
            filename = kf.get('filename')
            if source == 'user_file':
                with get_db_connection() as conn:
                    with conn.cursor() as cur:
                        cur.execute("SELECT content FROM user_files WHERE id = %s AND user_id = %s", (fid, user_id))
                        row = cur.fetchone()
                        if row and row[0]:
                            # Validate content
                            if is_valid_extracted_text(row[0]):
                                knowledge_content.append(f"--- 文件: {filename} ---\n{row[0]}")
                            else:
                                knowledge_content.append(f"--- 文件: {filename} (内容无效) ---\n无法读取文件内容")
            elif source == 'project_file':
                # similar validation
                with get_db_connection() as conn:
                    with conn.cursor() as cur:
                        cur.execute("""
                            SELECT pf.content, pf.original_name
                            FROM project_files pf
                            JOIN project_members pm ON pf.project_id = pm.project_id
                            WHERE pf.id = %s AND pm.user_id = %s
                        """, (fid, user_id))
                        row = cur.fetchone()
                        if row and row[0]:
                            if is_valid_extracted_text(row[0]):
                                knowledge_content.append(f"--- 文件: {row[1]} ---\n{row[0]}")
                            else:
                                knowledge_content.append(f"--- 文件: {row[1]} (内容无效) ---\n无法读取文件内容")
            elif source == 'knowledge_lab':
                with get_db_connection() as conn:
                    with conn.cursor() as cur:
                        cur.execute(
                            "SELECT content, original_name FROM knowledge_lab_files WHERE id = %s AND user_id = %s",
                            (fid, user_id))
                        row = cur.fetchone()
                        if row and row[0]:
                            if is_valid_extracted_text(row[0]):
                                knowledge_content.append(f"--- 知识库实验室文件: {row[1]} ---\n{row[0]}")
                            else:
                                knowledge_content.append(f"--- 知识库实验室文件: {row[1]} (内容无效) ---\n无法读取文件内容")
            elif source == 'company_kb':
                with get_db_connection() as conn:
                    with conn.cursor() as cur:
                        cur.execute("SELECT content, original_name FROM company_knowledge_base WHERE id = %s", (fid,))
                        row = cur.fetchone()
                        if row and row[0]:
                            if is_valid_extracted_text(row[0]):
                                knowledge_content.append(f"--- 公司知识库文件: {row[1]} ---\n{row[0]}")
                            else:
                                knowledge_content.append(f"--- 公司知识库文件: {row[1]} (内容无效) ---\n无法读取文件内容")

        # Build final query
        if knowledge_content:
            knowledge_text = "\n\n".join(knowledge_content)
            max_knowledge = 15000
            if len(knowledge_text) > max_knowledge:
                knowledge_text = knowledge_text[:max_knowledge] + "\n\n[知识库内容已截断...]"
            final_query = f"""你是一个基于知识库的助手。以下知识库内容具有最高优先级。请严格依据这些内容回答问题。如果知识库中没有相关信息，请明确告知用户。

=== 知识库内容 ===
{knowledge_text}

=== 用户上传的文件（如有） ===
{file_text}

=== 用户问题 ===
{user_msg}"""
        else:
            final_query = file_text

        # Store user message and get its ID
        user_msg_id = store_message(thread_id, 'user', user_msg)

        # ========== AGENT INVOCATION WITH ISOLATION ==========
        use_isolated_thread = (knowledge_files and len(knowledge_files) > 0) or has_files
        if use_isolated_thread:
            temp_thread_id = str(uuid.uuid4())
            config = {"configurable": {"thread_id": temp_thread_id}}
            from langgraph.checkpoint.memory import MemorySaver
            api_key = os.getenv("DEEPSEEK_API_KEY")
            if not api_key:
                api_key = os.getenv("QWEN_API_KEY") or os.getenv("DASHSCOPE_API_KEY")
            llm = ChatDeepSeek(
                model="deepseek-v4-pro",
                api_key=api_key,
                max_tokens=session.get('max_tokens', 1600),
                streaming=False,
                extra_body={"thinking": {"type": "disabled"}},
            )
            system_prompt = AGENT_SYSTEM_PROMPT
            checkpointer = MemorySaver()
            isolated_agent = create_agent(
                model=llm,
                tools=[get_date, bocha_search],
                system_prompt=system_prompt,
                checkpointer=checkpointer
            )
            try:
                response = isolated_agent.invoke(
                    {"messages": [{"role": "user", "content": final_query}]},
                    config
                )
            except Exception as e:
                logger.error(f"Isolated agent invoke failed: {e}", exc_info=True)
                return jsonify({"error": "AI 服务暂时不可用"}), 500
            finally:
                del isolated_agent
                del checkpointer
                del llm
        else:
            agent = get_agent()
            config = {"configurable": {"thread_id": thread_id}}
            try:
                response = agent.invoke({"messages": [{"role": "user", "content": final_query}]}, config)
            except Exception as e:
                logger.error(f"Agent invoke failed: {e}", exc_info=True)
                return jsonify({"error": "AI 服务暂时不可用"}), 500

        # Process response
        assistant_message = response["messages"][-1]
        raw_response = assistant_message.content
        reasoning = assistant_message.additional_kwargs.get('reasoning_content', '')
        if reasoning and reasoning.strip():
            thinking = reasoning.strip()
            answer = raw_response.strip() if raw_response else ''
        else:
            thinking, answer = split_thinking_answer(raw_response)

        assistant_msg_id = store_message(thread_id, 'assistant', answer, thinking)

        with get_db_connection() as conn:
            with conn.cursor() as cur:
                cur.execute("UPDATE message_responses SET assistant_response = %s, thinking = %s WHERE message_id = %s",
                            (answer, thinking, message_id))
                conn.commit()

        new_history = session.get('chat_history', [])
        new_history.append({"role": "user", "content": user_msg})
        new_history.append({"role": "assistant", "content": answer, "thinking": thinking})
        session['chat_history'] = new_history
        session['last_user_msg'] = user_msg
        session['last_assistant_msg'] = answer

        # Clear batch compare context after successful chat (one-time injection)
        session.pop('batch_compare_files', None)

        return jsonify({
            "assistant_message": answer,
            "thinking": thinking,
            "file_processed": len(uploaded_filenames) > 0,
            "ocr_attempted": is_image,
            "is_batch_report": False,
            "image_analysis_used": image_analysis_used,
            "assistant_message_id": assistant_msg_id,
            "user_message_id": user_msg_id
        })
    finally:
        if has_files:
            release_task_lock(user_id)

@app.route('/set_max_tokens', methods=['POST'])
def set_max_tokens():
    data = request.get_json()
    tokens = data.get('max_tokens', 4800)
    tokens = max(100, min(4800, tokens))
    session['max_tokens'] = tokens
    global _agent
    with _agent_lock:
        _agent = None
    return jsonify({"success": True, "max_tokens": tokens})


@app.route('/check_auth', methods=['GET'])
def check_auth():
    if session.get('consent_value', 0) != 1:
        return jsonify({"authenticated": False, "reason": "consent_not_given"})
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"authenticated": False, "reason": "no_user_id"})

    username = session.get('username')
    role = session.get('role')

    # If username not in session, try to fetch from database
    if not username:
        with get_db_connection() as conn:
            with conn.cursor() as cur:
                cur.execute("SELECT username, role FROM users WHERE user_id = %s", (user_id,))
                row = cur.fetchone()
                if row:
                    username = row[0]
                    role = row[1] or 'user'
                    # Restore to session for future requests
                    session['username'] = username
                    session['role'] = role
                else:
                    return jsonify({"authenticated": False, "reason": "user_deleted"})

    return jsonify({
        "authenticated": True,
        "username": username,
        "is_admin": role == 'admin',
        "user_id": user_id
    })

@app.route('/feedback', methods=['POST'])
def submit_feedback():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Feedback not allowed – no consent"}), 403
    data = request.get_json()
    rating = data.get('rating')
    comment = data.get('comment', '')
    user_message = data.get('user_message')
    assistant_response = data.get('assistant_response')
    if not user_message or not assistant_response:
        user_message = session.get('last_user_msg', '')
        assistant_response = session.get('last_assistant_msg', '')
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(
                "INSERT INTO feedback (thread_id, user_message, assistant_response, rating, comment, timestamp) VALUES (%s, %s, %s, %s, %s, %s)",
                (session['thread_id'], user_message, assistant_response, rating, comment, utc_now())
            )
            conn.commit()
    return jsonify({"status": "ok"})

@app.route('/get_recent_files', methods=['GET'])
def get_recent_files():
    thread_id = session.get('thread_id')
    if not thread_id:
        return jsonify({"recent_files": []})
    recent = file_cache_manager.get_recent_with_lock(thread_id)
    files_with_usage = []
    if session.get('consent_value', 0) == 1:
        with get_db_connection() as conn:
            with conn.cursor(cursor_factory=RealDictCursor) as cur:
                for filename in recent:
                    cur.execute(
                        """SELECT usage_type, question, timestamp
                           FROM file_usage
                           WHERE thread_id = %s
                             AND filename = %s
                           ORDER BY timestamp DESC
                           LIMIT 5""",
                        (thread_id, filename)
                    )
                    usage_records = []
                    for row in cur.fetchall():
                        ts_utc = row['timestamp']
                        if ts_utc:
                            ts_beijing = ts_utc.astimezone(BEIJING_TZ).strftime('%Y-%m-%d %H:%M:%S')
                        else:
                            ts_beijing = None
                        usage_records.append({
                            "type": row['usage_type'],
                            "question": row['question'],
                            "time": ts_beijing
                        })
                    files_with_usage.append({
                        "filename": filename,
                        "usage": usage_records
                    })
    else:
        for filename in recent:
            files_with_usage.append({"filename": filename, "usage": []})
    return jsonify({"recent_files": files_with_usage})

@app.route('/load_cached_file', methods=['POST'])
def load_cached_file():
    data = request.get_json()
    filename = data.get('filename')
    thread_id = session.get('thread_id')
    if not thread_id:
        return jsonify({"error": "Session expired"}), 401
    content = file_cache_manager.get_content(thread_id, filename)
    if content:
        return jsonify({"content": content})
    if session.get('consent_value', 0) != 1:
        user_id = get_user_id()
        temp_dir = get_anon_temp_dir(user_id)
        safe_name = re.sub(r'[^\w\-_\. ]', '_', filename) + '.txt'
        fpath = os.path.join(temp_dir, safe_name)
        if os.path.exists(fpath):
            with open(fpath, 'r', encoding='utf-8') as f:
                content = f.read()
            add_to_cache(thread_id, filename, content, user_id)
            return jsonify({"content": content})
        else:
            return jsonify({"error": "File not found"}), 404
    user_id = get_user_id()
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(
                "SELECT content FROM user_files WHERE user_id = %s AND filename = %s AND (expires_at IS NULL OR expires_at > NOW())",
                (user_id, filename)
            )
            row = cur.fetchone()
            if row:
                content = row[0] or ''
                add_to_cache(thread_id, filename, content, user_id)
                return jsonify({"content": content})
    return jsonify({"error": "File not found"}), 404

@app.route('/new_chat', methods=['POST'])
def new_chat():
    new_thread_id = str(uuid.uuid4())
    session['thread_id'] = new_thread_id
    session['chat_history'] = []
    get_or_create_session(new_thread_id)
    return jsonify({"thread_id": new_thread_id})

@app.route('/get_sessions', methods=['GET'])
def get_sessions():
    sessions = get_user_sessions()
    return jsonify({"sessions": sessions})

@app.route('/load_session/<thread_id>', methods=['GET'])
def load_session(thread_id):
    if session.get('consent_value', 0) != 1:
        messages = get_session_messages_anon(thread_id)
        session['thread_id'] = thread_id
        session['chat_history'] = messages
        return jsonify({"messages": messages, "thread_id": thread_id})

    user_sessions = get_user_sessions()
    if not any(s['thread_id'] == thread_id for s in user_sessions):
        return jsonify({"error": "Session not found"}), 404

    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            db_execute_readonly(cur)
            cur.execute(
                "SELECT id, role, content, thinking, timestamp FROM chat_messages WHERE thread_id = %s ORDER BY id ASC",
                (thread_id,)
            )
            rows = cur.fetchall()
            messages = []
            for row in rows:
                ts_utc = row['timestamp']
                ts_beijing = ts_utc.astimezone(BEIJING_TZ).strftime('%Y-%m-%d %H:%M:%S') if ts_utc else None
                messages.append({
                    "id": row['id'],
                    "role": row['role'],
                    "content": row['content'],
                    "thinking": row['thinking'],
                    "timestamp": ts_beijing
                })
            session['thread_id'] = thread_id
            session['chat_history'] = messages
            user_id = get_user_id()
            load_cache_from_db(thread_id, user_id)
            return jsonify({"messages": messages, "thread_id": thread_id})

@app.route('/delete_session/<thread_id>', methods=['POST'])
def delete_session_route(thread_id):
    user_sessions = get_user_sessions()
    if not any(s['thread_id'] == thread_id for s in user_sessions):
        return jsonify({"error": "Session not found"}), 404
    user_id = get_user_id()
    with user_task_lock:
        cleanup_stale_tasks()
        if user_id in user_active_tasks and user_active_tasks[user_id]['thread_id'] == thread_id:
            return jsonify({
                "error": "task_running",
                "message": "无法删除：该聊天正在进行资源密集型任务，请等待任务完成后再试。"
            }), 409
    try:
        archive_session(thread_id, user_id, reason="manual")
    except Exception as e:
        logger.error(f"Archive session failed for {thread_id}: {e}", exc_info=True)
    try:
        delete_session(thread_id)
        logger.info(f"Session {thread_id} deleted successfully for user {user_id}")
    except Exception as e:
        logger.error(f"Failed to delete session {thread_id}: {e}", exc_info=True)
        return jsonify({"error": "删除失败，请稍后重试"}), 500
    new_thread_id = None
    if session.get('thread_id') == thread_id:
        new_thread_id = str(uuid.uuid4())
        session['thread_id'] = new_thread_id
        session['chat_history'] = []
        get_or_create_session(new_thread_id)
        load_cache_from_db(new_thread_id, get_user_id())
    return jsonify({
        "status": "ok",
        "new_thread_id": new_thread_id,
        "messages": []
    })

@app.route('/archive_session/<thread_id>', methods=['POST'])
def archive_session_route(thread_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "Not logged in"}), 401
    try:
        archive_path = archive_session(thread_id, user_id, reason="manual")
        if archive_path:
            delete_session(thread_id)   # remove from active sessions
            return jsonify({"success": True})
        else:
            return jsonify({"error": "Archive failed"}), 500
    except Exception as e:
        logger.error(f"Archive session error: {e}", exc_info=True)
        return jsonify({"error": str(e)}), 500

@app.route('/regenerate', methods=['POST'])
def regenerate():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    data = request.get_json()
    user_message = data.get('user_message')
    if not user_message:
        return jsonify({"error": "Missing user_message"}), 400
    thread_id = session['thread_id']
    get_or_create_session(thread_id)
    agent = get_agent()
    config = {"configurable": {"thread_id": thread_id}}
    try:
        response = agent.invoke({"messages": [{"role": "user", "content": user_message}]}, config)
    except Exception as e:
        logger.error(f"Regenerate invoke failed: {e}", exc_info=True)
        return jsonify({"error": "AI 服务暂时不可用"}), 500
    assistant_message = response["messages"][-1]
    raw_response = assistant_message.content
    reasoning = assistant_message.additional_kwargs.get('reasoning_content', '')
    if reasoning and reasoning.strip():
        thinking = reasoning.strip()
        answer = raw_response.strip() if raw_response else ''
    else:
        thinking, answer = split_thinking_answer(raw_response)
    with get_db_connection() as conn:
        with db_transaction(conn):
            with conn.cursor() as cur:
                cur.execute("""
                            DELETE
                            FROM chat_messages
                            WHERE id IN (SELECT id
                                         FROM chat_messages
                                         WHERE thread_id = %s
                                         ORDER BY timestamp DESC
                                         LIMIT 2)
                            """, (thread_id,))
                conn.commit()
    store_message(thread_id, 'user', user_message)
    store_message(thread_id, 'assistant', answer if answer else raw_response, thinking if thinking else "")
    new_messages = get_session_messages(thread_id)
    session['chat_history'] = new_messages
    return jsonify({
        "assistant_message": answer if answer else raw_response,
        "thinking": thinking if thinking else ""
    })

@app.route('/check_storage', methods=['GET'])
def check_storage():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = get_user_id()
    total_bytes = get_user_total_storage_size(user_id)
    total_mb = total_bytes / (1024 * 1024)
    warning = total_mb > 300
    return jsonify({
        "total_mb": round(total_mb, 2),
        "warning": warning,
        "message": f"已使用 {total_mb:.2f} MB / 300 MB" if warning else None
    })

@app.route('/cleanup_now', methods=['POST'])
def cleanup_now():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    cleanup_old_sessions(days=15)
    return jsonify({"status": "ok", "message": "Cleanup completed"})

@app.route('/cleanup_anon_temp', methods=['POST'])
def cleanup_anon_temp():
    if session.get('consent_value', 0) != 1:
        user_id = get_user_id()
        temp_dir = get_anon_temp_dir(user_id)
        if os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)
            logger.info(f"Cleaned up anonymous temp directory for user {user_id}")
    return jsonify({"status": "ok"})

# Account routes
@app.route('/create_account', methods=['POST'])
def create_account():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403

    data = request.get_json()
    username = data.get('username', '').strip()
    pin = data.get('pin', '').strip()
    pin_length = data.get('pin_length', 6)

    if not username or not pin:
        return jsonify({"error": "用户名和PIN不能为空"}), 400
    if len(username) < 5 or len(username) > 18:
        return jsonify({"error": "用户名长度应为5-18个字符"}), 400
    if pin_length not in [4, 6] or len(pin) != pin_length:
        return jsonify({"error": f"PIN必须是{pin_length}位数字"}), 400
    if not pin.isdigit():
        return jsonify({"error": "PIN只能包含数字"}), 400

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT 1 FROM users WHERE username = %s", (username,))
            if cur.fetchone():
                return jsonify({"error": "用户名已存在"}), 409

            user_id = session.get('user_id')
            if not user_id:
                user_id = str(uuid.uuid4())
                session['user_id'] = user_id

            pin_hash = generate_password_hash(pin)

            cur.execute("""
                INSERT INTO users (user_id, username, pin_hash, pin_length, created_at, role)
                VALUES (%s, %s, %s, %s, NOW(), 'user')
                ON CONFLICT (user_id) DO UPDATE SET
                    username = EXCLUDED.username,
                    pin_hash = EXCLUDED.pin_hash,
                    pin_length = EXCLUDED.pin_length,
                    role = 'user'
                RETURNING user_id
            """, (user_id, username, pin_hash, pin_length))

            conn.commit()
            session['username'] = username
            session['role'] = 'user'
            session.modified = True

            return jsonify({"success": True, "username": username})

def cleanup_orphan_users():
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("""
                DELETE FROM users
                WHERE (username IS NULL OR username = '')
                  AND created_at < NOW() - INTERVAL '1 day'
                  AND NOT EXISTS (SELECT 1 FROM chat_sessions WHERE user_id = users.user_id)
                  AND NOT EXISTS (SELECT 1 FROM user_files WHERE user_id = users.user_id)
                  AND NOT EXISTS (SELECT 1 FROM project_members WHERE user_id = users.user_id)
                  AND NOT EXISTS (SELECT 1 FROM projects WHERE created_by = users.user_id)
            """)
            conn.commit()
            logger.info(f"Deleted {cur.rowcount} orphan empty users")

@app.route('/login', methods=['POST'])
def login():
    data = request.get_json()
    username = data.get('username', '').strip()
    pin = data.get('pin', '').strip()
    if not username or not pin:
        return jsonify({"error": "用户名和PIN不能为空"}), 400

    if username == "admin":
        if not ADMIN_PASSWORD_HASH:
            logger.error("ADMIN_PASSWORD_HASH not set in environment")
            return jsonify({"error": "管理员账户未配置"}), 500
        if not check_password_hash(ADMIN_PASSWORD_HASH, pin):
            logger.warning(f"Admin login failed for {username}")
            return jsonify({"error": "用户名或PIN错误"}), 401

        with get_db_connection() as conn:
            with conn.cursor() as cur:
                cur.execute("SELECT user_id FROM users WHERE username = 'admin'")
                admin_row = cur.fetchone()
                if admin_row:
                    user_id = admin_row[0]
                else:
                    user_id = str(uuid.uuid4())
                    cur.execute(
                        "INSERT INTO users (user_id, username, role) VALUES (%s, %s, %s)",
                        (user_id, 'admin', 'admin')
                    )
                conn.commit()

        session['user_id'] = user_id
        session['consent_value'] = 1
        session['username'] = 'admin'
        session['role'] = 'admin'
        session.permanent = True

        with get_db_connection() as conn:
            with conn.cursor() as cur:
                cur.execute(
                    "INSERT INTO consent (thread_id, consent_given, timestamp) VALUES (%s, %s, NOW()) ON CONFLICT (thread_id) DO UPDATE SET consent_given = EXCLUDED.consent_given, timestamp = EXCLUDED.timestamp",
                    (session.get('thread_id', str(uuid.uuid4())), 1)
                )
                conn.commit()

        logger.info(f"Admin logged in: {user_id}")
        return jsonify({
            "success": True,
            "username": "admin",
            "is_admin": True,
            "user_id": session['user_id']
        })

    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute(
                "SELECT user_id, pin_hash, pin_length, role FROM users WHERE username = %s AND is_active = TRUE AND pin_hash IS NOT NULL",
                (username,)
            )
            user = cur.fetchone()
            if not user or not check_password_hash(user['pin_hash'], pin):
                return jsonify({"error": "用户名或PIN错误"}), 401

            session['user_id'] = user['user_id']
            session['consent_value'] = 1
            session['username'] = username
            session['role'] = user.get('role', 'user')
            session.permanent = True

            with conn.cursor() as cur2:
                cur2.execute(
                    "INSERT INTO consent (thread_id, consent_given, timestamp) VALUES (%s, %s, NOW()) ON CONFLICT (thread_id) DO UPDATE SET consent_given = EXCLUDED.consent_given, timestamp = EXCLUDED.timestamp",
                    (session.get('thread_id', str(uuid.uuid4())), 1)
                )
            conn.commit()
            return jsonify({
                "success": True,
                "username": username,
                "is_admin": session['role'] == 'admin',
                "user_id": session['user_id']
            })

@app.route('/update_account', methods=['POST'])
def update_account():
    if session.get('consent_value', 0) != 1 or not session.get('user_id'):
        return jsonify({"error": "未登录"}), 401

    data = request.get_json()
    new_username = data.get('new_username', '').strip()
    new_pin = data.get('new_pin', '').strip()
    pin_length = int(data.get('pin_length', 6))
    current_pin = data.get('current_pin', '').strip()

    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("SELECT pin_hash FROM users WHERE user_id = %s", (session['user_id'],))
            user = cur.fetchone()
            if not user:
                return jsonify({"error": "用户不存在"}), 404

            existing_hash = user['pin_hash']

            if existing_hash is not None:
                if not current_pin or not check_password_hash(existing_hash, current_pin):
                    return jsonify({"error": "当前PIN错误"}), 401

            updates = []
            params = []

            if new_username:
                if len(new_username) < 5 or len(new_username) > 18:
                    return jsonify({"error": "用户名长度应为5-18个字符"}), 400
                cur.execute("SELECT 1 FROM users WHERE username = %s AND user_id != %s",
                            (new_username, session['user_id']))
                if cur.fetchone():
                    return jsonify({"error": "用户名已存在"}), 409
                updates.append("username = %s")
                params.append(new_username)
                session['username'] = new_username

            if new_pin:
                if pin_length not in (4, 6) or len(new_pin) != pin_length:
                    return jsonify({"error": f"PIN必须是{pin_length}位数字"}), 400
                if not new_pin.isdigit():
                    return jsonify({"error": "PIN只能包含数字"}), 400
                updates.append("pin_hash = %s")
                params.append(generate_password_hash(new_pin))
                updates.append("pin_length = %s")
                params.append(pin_length)

            if updates:
                params.append(session['user_id'])
                cur.execute(f"UPDATE users SET {', '.join(updates)} WHERE user_id = %s", params)
                conn.commit()
            return jsonify({"success": True})

@app.route('/delete_account', methods=['POST'])
def delete_account():
    if session.get('consent_value', 0) != 1 or not session.get('user_id'):
        return jsonify({"error": "未登录"}), 401

    data = request.get_json()
    pin = data.get('pin', '').strip()
    user_id = session['user_id']

    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("SELECT pin_hash FROM users WHERE user_id = %s", (user_id,))
            user = cur.fetchone()
            if not user:
                return jsonify({"error": "用户不存在"}), 404

            pin_hash = user['pin_hash']

            if pin_hash is None:
                return jsonify({"error": "您的账户没有设置PIN，无法删除。请先通过账户设置设置PIN后再试。"}), 400

            if not check_password_hash(pin_hash, pin):
                return jsonify({"error": "PIN错误"}), 401

            with db_transaction(conn):
                cur.execute("""
                    SELECT DISTINCT p.id, p.name, p.created_by
                    FROM projects p
                    LEFT JOIN project_members pm ON p.id = pm.project_id
                    WHERE pm.user_id = %s OR p.created_by = %s
                """, (user_id, user_id))
                projects = cur.fetchall()
                for proj in projects:
                    proj_id = proj['id']
                    proj_name = proj['name']
                    cur.execute("""
                        INSERT INTO task_deposit_items (original_user_id, original_username, project_id,
                                                        project_name, item_type, item_data)
                        VALUES (%s, %s, %s, %s, 'project', %s)
                    """, (user_id, user.get('username', 'unknown'), proj_id, proj_name,
                          json.dumps({'project_id': proj_id, 'name': proj_name})))
                    cur.execute("""
                        SELECT id, original_name, stored_path, uploaded_by, folder_id, filename, version, file_size
                        FROM project_files WHERE project_id = %s
                    """, (proj_id,))
                    files = cur.fetchall()
                    for f in files:
                        cur.execute("""
                            INSERT INTO task_deposit_items (original_user_id, original_username, project_id,
                                                            project_name, item_type, item_data, stored_path)
                            VALUES (%s, %s, %s, %s, 'file', %s, %s)
                        """, (user_id, user.get('username', 'unknown'), proj_id, proj_name,
                              json.dumps(dict(f)), f['stored_path']))
                    cur.execute("""
                        SELECT id, name, parent_folder_id, created_by
                        FROM project_folders WHERE project_id = %s
                    """, (proj_id,))
                    folders = cur.fetchall()
                    for fold in folders:
                        cur.execute("""
                            INSERT INTO task_deposit_items (original_user_id, original_username, project_id,
                                                            project_name, item_type, item_data)
                            VALUES (%s, %s, %s, %s, 'folder', %s)
                        """, (user_id, user.get('username', 'unknown'), proj_id, proj_name,
                              json.dumps(dict(fold))))
                    cur.execute("""
                        SELECT id, file_id, user_id, comment, created_at
                        FROM project_file_comments
                        WHERE file_id IN (SELECT id FROM project_files WHERE project_id = %s)
                    """, (proj_id,))
                    comments = cur.fetchall()
                    for comm in comments:
                        cur.execute("""
                            INSERT INTO task_deposit_items (original_user_id, original_username, project_id,
                                                            project_name, item_type, item_data)
                            VALUES (%s, %s, %s, %s, 'comment', %s)
                        """, (user_id, user.get('username', 'unknown'), proj_id, proj_name,
                              json.dumps(dict(comm))))

                cur.execute("SELECT user_id FROM users WHERE role = 'admin' AND user_id != %s LIMIT 1", (user_id,))
                admin_row = cur.fetchone()
                admin_id = admin_row['user_id'] if admin_row else None

                if admin_id:
                    cur.execute("UPDATE projects SET created_by = %s WHERE created_by = %s", (admin_id, user_id))
                else:
                    cur.execute("UPDATE projects SET created_by = NULL WHERE created_by = %s", (user_id,))

                if admin_id:
                    cur.execute("UPDATE project_folders SET created_by = %s WHERE created_by = %s", (admin_id, user_id))
                else:
                    cur.execute("UPDATE project_folders SET created_by = NULL WHERE created_by = %s", (user_id,))

                if admin_id:
                    cur.execute("UPDATE project_files SET uploaded_by = %s WHERE uploaded_by = %s", (admin_id, user_id))
                else:
                    cur.execute("UPDATE project_files SET uploaded_by = NULL WHERE uploaded_by = %s", (user_id,))

                if admin_id:
                    cur.execute("""
                        UPDATE project_file_versions SET uploaded_by = %s
                        WHERE uploaded_by = %s
                    """, (admin_id, user_id))
                else:
                    cur.execute("UPDATE project_file_versions SET uploaded_by = NULL WHERE uploaded_by = %s", (user_id,))

                if admin_id:
                    cur.execute("UPDATE project_members SET added_by = %s WHERE added_by = %s", (admin_id, user_id))
                else:
                    cur.execute("UPDATE project_members SET added_by = NULL WHERE added_by = %s", (user_id,))

                cur.execute("UPDATE project_file_comments SET user_id = NULL WHERE user_id = %s", (user_id,))
                cur.execute("UPDATE project_folder_comments SET user_id = NULL WHERE user_id = %s", (user_id,))
                cur.execute("UPDATE task_deposit_items SET transferred_to_user_id = NULL WHERE transferred_to_user_id = %s", (user_id,))
                cur.execute("DELETE FROM recycle_bin WHERE user_id = %s", (user_id,))
                cur.execute("UPDATE project_recycle_bin SET uploaded_by = NULL WHERE uploaded_by = %s", (user_id,))
                cur.execute("UPDATE task_deposit_items SET original_user_id = NULL WHERE original_user_id = %s", (user_id,))
                cur.execute("UPDATE task_deposit_permissions SET manager_id = NULL WHERE manager_id = %s", (user_id,))
                cur.execute("UPDATE task_deposit_permissions SET granted_by = NULL WHERE granted_by = %s", (user_id,))
                cur.execute("DELETE FROM project_members WHERE user_id = %s", (user_id,))
                cur.execute("DELETE FROM chat_messages WHERE thread_id IN (SELECT thread_id FROM chat_sessions WHERE user_id = %s)", (user_id,))
                cur.execute("DELETE FROM user_files WHERE user_id = %s", (user_id,))
                cur.execute("DELETE FROM file_usage WHERE user_id = %s", (user_id,))
                cur.execute("DELETE FROM feedback WHERE thread_id IN (SELECT thread_id FROM chat_sessions WHERE user_id = %s)", (user_id,))
                cur.execute("DELETE FROM consent WHERE thread_id IN (SELECT thread_id FROM chat_sessions WHERE user_id = %s)", (user_id,))
                cur.execute("DELETE FROM chat_sessions WHERE user_id = %s", (user_id,))
                cur.execute("DELETE FROM users WHERE user_id = %s", (user_id,))

                conn.commit()

            session.clear()
            session['consent_value'] = 0
            session['thread_id'] = str(uuid.uuid4())
            get_or_create_session(session['thread_id'])
            return jsonify({"success": True})
# Check if column_name exists in table_name
def validate_table_column(table_name, column_name):
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT column_name 
                FROM information_schema.columns 
                WHERE table_name = %s AND column_name = %s
            """, (table_name, column_name))
            return cur.fetchone() is not None

# Task deposit endpoints
@app.route('/admin/task_deposit', methods=['GET'])
def get_task_deposit():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "Not logged in"}), 401
    is_admin_user = session.get('role') == 'admin'
    if not is_admin_user:
        return jsonify({"error": "Access denied"}), 403
    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                        SELECT id,
                               original_user_id,
                               original_username,
                               project_id,
                               project_name,
                               item_type,
                               item_data,
                               stored_path,
                               transferred_to_user_id,
                               transferred_at,
                               created_at
                        FROM task_deposit_items
                        WHERE deleted_at IS NULL
                        ORDER BY created_at DESC
                        """)
            items = cur.fetchall()
            return jsonify({"items": items})

@app.route('/admin/task_deposit/transfer/<int:item_id>', methods=['POST'])
def transfer_task_deposit_item(item_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if session.get('role') != 'admin':
        return jsonify({"error": "Only admin can transfer deposit items"}), 403
    data = request.get_json()
    target_user_id = data.get('target_user_id')
    if not target_user_id:
        return jsonify({"error": "Missing target_user_id"}), 400
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT 1 FROM users WHERE user_id = %s", (target_user_id,))
            if not cur.fetchone():
                return jsonify({"error": "Target user not found"}), 404
            cur.execute("""
                        UPDATE task_deposit_items
                        SET transferred_to_user_id = %s,
                            transferred_at         = NOW()
                        WHERE id = %s
                          AND deleted_at IS NULL
                        RETURNING id, item_type, item_data, stored_path
                        """, (target_user_id, item_id))
            item = cur.fetchone()
            if not item:
                return jsonify({"error": "Item not found or already deleted"}), 404
            conn.commit()
            return jsonify({"success": True, "item": dict(item)})

# Permission helpers for projects
def is_admin():
    return session.get('role') == 'admin'

def admin_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if not is_admin():
            return jsonify({"error": "Admin access required"}), 403
        return f(*args, **kwargs)
    return decorated_function

def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if session.get('consent_value', 0) != 1:
            return jsonify({"error": "Consent not given"}), 403
        if not session.get('user_id'):
            return jsonify({"error": "Not logged in"}), 401
        return f(*args, **kwargs)
    return decorated_function

def get_user_role_in_project(project_id, user_id):
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT role FROM project_members WHERE project_id = %s AND user_id = %s",
                        (project_id, user_id))
            row = cur.fetchone()
            return row[0] if row else None


@app.route('/user_project_files', methods=['GET'])
def get_user_project_files():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "Not logged in"}), 401

    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                        SELECT pf.id, pf.original_name, pf.file_size, p.name as project_name
                        FROM project_files pf
                                 JOIN projects p ON pf.project_id = p.id
                                 JOIN project_members pm ON p.id = pm.project_id
                        WHERE pm.user_id = %s
                        ORDER BY p.name, pf.uploaded_at DESC
                        """, (user_id,))
            files = cur.fetchall()
            return jsonify({"files": files})

def can_manage_files(project_id, user_id):
    if is_admin():
        return True
    role = get_user_role_in_project(project_id, user_id)
    return role in ('admin', 'manager')

def can_edit_file(project_id, file_id, user_id):
    if is_admin():
        return True
    role = get_user_role_in_project(project_id, user_id)
    if role == 'manager':
        return True
    if role == 'member':
        with get_db_connection() as conn:
            with conn.cursor() as cur:
                cur.execute("SELECT uploaded_by FROM project_files WHERE id = %s AND project_id = %s",
                            (file_id, project_id))
                row = cur.fetchone()
                return row and row[0] == user_id
    return False

def can_move_file(project_id, file_id, user_id):
    if is_admin():
        return True
    role = get_user_role_in_project(project_id, user_id)
    if role in ('admin', 'manager', 'member'):
        return True
    return False

def can_edit_folder(project_id, folder_id, user_id):
    if is_admin():
        return True
    role = get_user_role_in_project(project_id, user_id)
    if role == 'manager':
        return True
    if role == 'member':
        with get_db_connection() as conn:
            with conn.cursor() as cur:
                cur.execute("SELECT created_by FROM project_folders WHERE id = %s AND project_id = %s",
                            (folder_id, project_id))
                row = cur.fetchone()
                return row and row[0] == user_id
    return False

def can_manage_members(project_id, user_id):
    if is_admin():
        return True
    role = get_user_role_in_project(project_id, user_id)
    return role == 'manager'

def can_access_project(project_id, user_id):
    if is_admin():
        return True
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT 1 FROM project_members WHERE project_id = %s AND user_id = %s", (project_id, user_id))
            return cur.fetchone() is not None

def user_has_any_project(user_id):
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT 1 FROM project_members WHERE user_id = %s LIMIT 1", (user_id,))
            return cur.fetchone() is not None

# ---------- Project management routes ----------
@app.route('/admin/projects', methods=['POST'])
@admin_required
def create_project():
    data = request.get_json()
    name = data.get('name', '').strip()
    description = data.get('description', '').strip()
    if not name:
        return jsonify({"error": "Project name required"}), 400
    user_id = session.get('user_id')
    with get_db_connection() as conn:
        with db_transaction(conn):
            with conn.cursor() as cur:
                cur.execute(
                    "INSERT INTO projects (name, description, created_by, status) VALUES (%s, %s, %s, 'active') RETURNING id",
                    (name, description, user_id))
                project_id = cur.fetchone()[0]
                cur.execute(
                    "INSERT INTO project_members (project_id, user_id, role, added_by) VALUES (%s, %s, 'admin', %s)",
                    (project_id, user_id, user_id))
                cur.execute(
                    "INSERT INTO project_folders (project_id, parent_folder_id, name, created_by) VALUES (%s, NULL, %s, %s)",
                    (project_id, name, user_id))
                conn.commit()
                return jsonify({"success": True, "id": project_id})

@app.route('/admin/projects', methods=['GET'])
def get_projects():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"projects": [], "has_projects": False})
    if is_admin():
        with get_db_connection() as conn:
            with conn.cursor(cursor_factory=RealDictCursor) as cur:
                cur.execute(
                    "SELECT id, name, description, created_at, updated_at, status, archived_at, deletion_scheduled_at FROM projects ORDER BY CASE status WHEN 'active' THEN 1 WHEN 'archived' THEN 2 WHEN 'aborted' THEN 3 END, created_at DESC")
                projects = cur.fetchall()
                has_projects = len(projects) > 0
                return jsonify({"projects": projects, "has_projects": has_projects})
    else:
        with get_db_connection() as conn:
            with conn.cursor(cursor_factory=RealDictCursor) as cur:
                cur.execute("""
                            SELECT p.id,
                                   p.name,
                                   p.description,
                                   p.created_at,
                                   p.updated_at,
                                   p.status,
                                   p.archived_at,
                                   p.deletion_scheduled_at
                            FROM projects p
                                     JOIN project_members pm ON p.id = pm.project_id
                            WHERE pm.user_id = %s
                            ORDER BY CASE p.status
                                         WHEN 'active' THEN 1
                                         WHEN 'archived' THEN 2
                                         WHEN 'aborted' THEN 3 END, p.created_at DESC
                            """, (user_id,))
                projects = cur.fetchall()
                has_projects = len(projects) > 0
                return jsonify({"projects": projects, "has_projects": has_projects})

@app.route('/admin/projects/<int:project_id>', methods=['PUT'])
def update_project(project_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "Not logged in"}), 401

    if not is_admin() and not can_manage_files(project_id, user_id):
        return jsonify({"error": "Permission denied"}), 403

    data = request.get_json()
    name = data.get('name', '').strip()
    description = data.get('description', '').strip()
    if not name:
        return jsonify({"error": "Project name required"}), 400

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("""
                UPDATE projects
                SET name = %s, description = %s, updated_at = NOW()
                WHERE id = %s
                RETURNING id
            """, (name, description, project_id))
            if cur.fetchone():
                conn.commit()
                return jsonify({"success": True})
            else:
                return jsonify({"error": "Project not found"}), 404

@app.route('/admin/projects/<int:project_id>', methods=['DELETE'])
@admin_required
def delete_project(project_id):
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT status FROM projects WHERE id = %s", (project_id,))
            row = cur.fetchone()
            if not row:
                return jsonify({"error": "Project not found"}), 404
            status = row[0]
            if status not in ('archived', 'aborted'):
                return jsonify({"error": "Only archived or aborted projects can be deleted"}), 400

            cur.execute("SELECT stored_path FROM project_files WHERE project_id = %s", (project_id,))
            for (stored_path,) in cur.fetchall():
                if stored_path and os.path.exists(stored_path):
                    try:
                        os.remove(stored_path)
                    except Exception as e:
                        logger.warning(f"Could not delete project file {stored_path}: {e}")

            cur.execute("DELETE FROM projects WHERE id = %s", (project_id,))
            conn.commit()
            return jsonify({"success": True})

@app.route('/admin/projects/<int:project_id>/files/<int:file_id>', methods=['DELETE'])
def delete_project_file(project_id, file_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not can_edit_file(project_id, file_id, user_id):
        return jsonify({"error": "Permission denied"}), 403

    with get_db_connection() as conn:
        with db_transaction(conn):
            with conn.cursor(cursor_factory=RealDictCursor) as cur:
                cur.execute("""
                    SELECT id, original_name, file_size, stored_path, uploaded_by, folder_id, filename, version, file_hash, project_id
                    FROM project_files
                    WHERE id = %s AND project_id = %s
                """, (file_id, project_id))
                file_record = cur.fetchone()
                if not file_record:
                    return jsonify({"error": "File not found"}), 404

                cur.execute("""
                    INSERT INTO project_recycle_bin 
                    (original_table, original_id, project_id, folder_id, file_name, original_name, file_size, stored_path, file_hash, version, uploaded_by, deleted_at, expires_at)
                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, NOW(), NOW() + INTERVAL '3 days')
                """, (
                    'project_files', file_record['id'], file_record['project_id'], file_record['folder_id'],
                    file_record['original_name'],
                    file_record['original_name'],
                    file_record['file_size'], file_record['stored_path'], file_record['file_hash'],
                    file_record['version'], file_record['uploaded_by']
                ))

                cur.execute("DELETE FROM project_files WHERE id = %s AND project_id = %s", (file_id, project_id))
                conn.commit()
                return jsonify({"success": True, "moved_to_recycle_bin": True})

@app.route('/admin/projects/<int:project_id>/abort', methods=['POST'])
@admin_required
def abort_project(project_id):
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("UPDATE projects SET status = 'aborted', archived_at = NOW() WHERE id = %s RETURNING id",
                        (project_id,))
            if cur.fetchone():
                conn.commit()
                return jsonify({"success": True})
            return jsonify({"error": "Project not found"}), 404

@app.route('/admin/projects/<int:project_id>/finish', methods=['POST'])
def finish_project(project_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not can_manage_files(project_id, user_id):
        return jsonify({"error": "Only admin or project manager can finish a project"}), 403
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT name FROM projects WHERE id = %s AND status = 'active'", (project_id,))
            project = cur.fetchone()
            if not project:
                return jsonify({"error": "Project not found or already finished/aborted"}), 404
            project_name = project[0]
            cur.execute("SELECT stored_path, original_name FROM project_files WHERE project_id = %s", (project_id,))
            files = cur.fetchall()
            if not files:
                return jsonify({"error": "No files to archive"}), 400
            zip_dir = os.path.join(PROJECT_FILES_ROOT, 'archives')
            os.makedirs(zip_dir, exist_ok=True)
            safe_name = re.sub(r'[^\w\-_\.]', '_', project_name)
            zip_filename = f"project_{project_id}_{safe_name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip"
            zip_path = os.path.join(zip_dir, zip_filename)
            with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for stored_path, original_name in files:
                    zipf.write(stored_path, original_name)
            cur.execute("UPDATE projects SET status = 'archived', archived_at = NOW() WHERE id = %s", (project_id,))
            conn.commit()
            return jsonify({
                "success": True,
                "download_url": f"/admin/projects/{project_id}/download_archive/{zip_filename}",
                "zip_filename": zip_filename
            })

@app.route('/admin/projects/<int:project_id>/download_archive/<zip_filename>', methods=['GET'])
def download_archive(project_id, zip_filename):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    zip_dir = os.path.join(PROJECT_FILES_ROOT, 'archives')
    zip_path = os.path.join(zip_dir, zip_filename)
    if not os.path.exists(zip_path):
        return jsonify({"error": "Archive not found"}), 404
    return send_file(zip_path, as_attachment=True, download_name=zip_filename)

# Project members routes
@app.route('/admin/projects/<int:project_id>/members', methods=['GET'])
def get_project_members(project_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not is_admin() and not can_access_project(project_id, user_id):
        return jsonify({"error": "Access denied"}), 403
    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                SELECT u.user_id, u.username, pm.role, pm.added_at
                FROM project_members pm
                JOIN users u ON pm.user_id = u.user_id
                WHERE pm.project_id = %s AND u.username IS NOT NULL AND u.username != ''
                ORDER BY pm.role, u.username
            """, (project_id,))
            members = cur.fetchall()
            return jsonify({"members": members})

@app.route('/admin/projects/<int:project_id>/members/search', methods=['GET'])
def search_users_to_add(project_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not is_admin() and not can_manage_members(project_id, user_id):
        return jsonify({"error": "Access denied"}), 403

    query = request.args.get('q', '').strip()
    if len(query) < 2:
        return jsonify({"users": []})
    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                SELECT user_id, username
                FROM users
                WHERE username ILIKE %s
                  AND user_id NOT IN (SELECT user_id FROM project_members WHERE project_id = %s)
                  AND user_id != %s
                  AND role != 'admin'
                LIMIT 20
            """, (f'%{query}%', project_id, user_id))
            users = cur.fetchall()
            return jsonify({"users": users})

@app.route('/admin/projects/<int:project_id>/all_users', methods=['GET'])
def get_all_users_for_project(project_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    current_user_id = session.get('user_id')
    if not current_user_id:
        return jsonify({"error": "Not logged in"}), 401
    if not is_admin() and not can_manage_members(project_id, current_user_id):
        return jsonify({"error": "Access denied"}), 403

    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                SELECT user_id, username
                FROM users
                WHERE username IS NOT NULL AND username != ''
                  AND user_id NOT IN (SELECT user_id FROM project_members WHERE project_id = %s)
                  AND user_id != %s
                  AND role != 'admin'
                ORDER BY username
                LIMIT 100
            """, (project_id, current_user_id))
            users = cur.fetchall()
            return jsonify({"users": users})

@app.route('/admin/projects/<int:project_id>/members', methods=['POST'])
def add_project_member(project_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not can_manage_members(project_id, user_id):
        return jsonify({"error": "Only admin or project manager can add members"}), 403

    data = request.get_json()
    new_user_id = data.get('user_id')
    role = data.get('role', 'member')
    if role == 'manager' and not is_admin():
        return jsonify({"error": "Only admin can add managers"}), 403
    if role not in ('member', 'manager'):
        return jsonify({"error": "Invalid role"}), 400

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT role FROM users WHERE user_id = %s", (new_user_id,))
            row = cur.fetchone()
            if row and row[0] == 'admin':
                return jsonify({"error": "Cannot add a global admin as a project member"}), 403

            if not row:
                return jsonify({"error": "User not found"}), 404

            cur.execute("SELECT 1 FROM project_members WHERE project_id = %s AND user_id = %s",
                        (project_id, new_user_id))
            if cur.fetchone():
                return jsonify({"error": "User already a member"}), 409

            cur.execute("""
                INSERT INTO project_members (project_id, user_id, role, added_by)
                VALUES (%s, %s, %s, %s)
            """, (project_id, new_user_id, role, user_id))
            conn.commit()
            return jsonify({"success": True})

@app.route('/admin/projects/<int:project_id>/members/<user_id>', methods=['PUT'])
@admin_required
def update_member_role(project_id, user_id):
    data = request.get_json()
    new_role = data.get('role')
    if new_role not in ('member', 'manager'):
        return jsonify({"error": "Invalid role"}), 400

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT role FROM users WHERE user_id = %s", (user_id,))
            row = cur.fetchone()
            if row and row[0] == 'admin':
                return jsonify({"error": "Cannot modify global admin's role"}), 403

            cur.execute("""
                UPDATE project_members
                SET role = %s
                WHERE project_id = %s AND user_id = %s
                RETURNING user_id
            """, (new_role, project_id, user_id))
            if cur.rowcount == 0:
                return jsonify({"error": "Member not found"}), 404
            conn.commit()
            return jsonify({"success": True})

@app.route('/admin/projects/<int:project_id>/members/<user_id>', methods=['DELETE'])
def remove_project_member(project_id, user_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    current_user_id = session.get('user_id')
    if not can_manage_members(project_id, current_user_id):
        return jsonify({"error": "Only admin or project manager can remove members"}), 403

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT role FROM users WHERE user_id = %s", (user_id,))
            row = cur.fetchone()
            if row and row[0] == 'admin':
                return jsonify({"error": "Cannot remove a global admin"}), 403

            cur.execute("""
                SELECT role FROM project_members
                WHERE project_id = %s AND user_id = %s
            """, (project_id, user_id))
            target_member = cur.fetchone()
            if not target_member:
                return jsonify({"error": "Member not found"}), 404

            target_role = target_member[0]
            if target_role == 'admin':
                return jsonify({"error": "Cannot remove the project admin"}), 403
            if target_role == 'manager' and not is_admin():
                return jsonify({"error": "Only admin can remove managers"}), 403

            cur.execute("DELETE FROM project_members WHERE project_id = %s AND user_id = %s",
                        (project_id, user_id))
            if cur.rowcount == 0:
                return jsonify({"error": "Member not found"}), 404
            conn.commit()
            return jsonify({"success": True})

@app.route('/admin/projects/<int:project_id>/transfer_manager/<user_id>', methods=['POST'])
def transfer_manager_role(project_id, user_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    current_user_id = session.get('user_id')
    current_role = get_user_role_in_project(project_id, current_user_id)
    if current_role != 'manager':
        return jsonify({"error": "Only a manager can transfer manager rights"}), 403
    target_role = get_user_role_in_project(project_id, user_id)
    if target_role != 'member':
        return jsonify({"error": "Target user must be a member"}), 400
    with get_db_connection() as conn:
        with db_transaction(conn):
            with conn.cursor() as cur:
                cur.execute("UPDATE project_members SET role = 'member' WHERE project_id = %s AND user_id = %s",
                            (project_id, current_user_id))
                cur.execute("UPDATE project_members SET role = 'manager' WHERE project_id = %s AND user_id = %s",
                            (project_id, user_id))
                conn.commit()
    return jsonify({"success": True})

# Project folders and files
def ensure_root_folder(project_id):
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT id FROM project_folders WHERE project_id = %s AND parent_folder_id IS NULL",
                        (project_id,))
            if not cur.fetchone():
                cur.execute("SELECT name FROM projects WHERE id = %s", (project_id,))
                row = cur.fetchone()
                if row:
                    project_name = row[0]
                    cur.execute(
                        "INSERT INTO project_folders (project_id, parent_folder_id, name, created_by) VALUES (%s, NULL, %s, %s)",
                        (project_id, project_name, session.get('user_id')))
                    conn.commit()
                    logger.info(f"Created missing root folder for project {project_id}")

def build_folder_path(folder_id, folder_dict):
    parts = []
    current_id = folder_id
    while current_id:
        folder = folder_dict.get(current_id)
        if not folder:
            break
        parts.insert(0, folder['name'])
        current_id = folder['parent_folder_id']
    return '/' + '/'.join(parts) if parts else '/'

@app.route('/admin/projects/<int:project_id>/folders', methods=['GET'])
def get_folders(project_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not is_admin() and not can_access_project(project_id, user_id):
        return jsonify({"error": "Access denied"}), 403
    ensure_root_folder(project_id)
    try:
        with get_db_connection() as conn:
            with conn.cursor(cursor_factory=RealDictCursor) as cur:
                cur.execute(
                    "SELECT id, parent_folder_id, name FROM project_folders WHERE project_id = %s ORDER BY parent_folder_id, name",
                    (project_id,))
                folders = cur.fetchall()
                if not folders:
                    return jsonify({"folders": []})
                folder_dict = {f['id']: f for f in folders}
                for f in folder_dict.values():
                    f['children'] = []
                root_folders = []
                for f in folder_dict.values():
                    if f['parent_folder_id'] is None:
                        root_folders.append(f)
                    else:
                        parent = folder_dict.get(f['parent_folder_id'])
                        if parent:
                            parent['children'].append(f)
                        else:
                            root_folders.append(f)
                for f in folder_dict.values():
                    f['path'] = build_folder_path(f['id'], folder_dict)
                return jsonify({"folders": root_folders})
    except Exception as e:
        logger.error(f"Error in get_folders: {e}", exc_info=True)
        return jsonify({"error": "Internal server error"}), 500

@app.route('/admin/projects/<int:project_id>/folders', methods=['POST'])
def create_folder(project_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not is_admin() and not can_access_project(project_id, user_id):
        return jsonify({"error": "Access denied"}), 403
    data = request.get_json()
    name = data.get('name', '').strip()
    parent_folder_id = data.get('parent_folder_id')
    if not name:
        return jsonify({"error": "Folder name required"}), 400
    if parent_folder_id is None:
        return jsonify({"error": "Cannot create root folder. Only one root folder exists per project."}), 400
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT id FROM project_folders WHERE id = %s AND project_id = %s",
                        (parent_folder_id, project_id))
            if not cur.fetchone():
                return jsonify({"error": "Parent folder not found"}), 404
            cur.execute(
                "INSERT INTO project_folders (project_id, parent_folder_id, name, created_by) VALUES (%s, %s, %s, %s) RETURNING id",
                (project_id, parent_folder_id, name, user_id))
            new_id = cur.fetchone()[0]
            conn.commit()
            return jsonify({"success": True, "id": new_id})

@app.route('/admin/projects/<int:project_id>/folders/<int:folder_id>', methods=['DELETE'])
def delete_folder(project_id, folder_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not can_edit_folder(project_id, folder_id, user_id):
        return jsonify({"error": "Permission denied"}), 403

    with get_db_connection() as conn:
        with db_transaction(conn):
            with conn.cursor(cursor_factory=RealDictCursor) as cur:
                cur.execute("""
                    WITH RECURSIVE folder_tree AS (
                        SELECT id, name, parent_folder_id, created_at, created_by
                        FROM project_folders
                        WHERE id = %s AND project_id = %s
                        UNION ALL
                        SELECT pf.id, pf.name, pf.parent_folder_id, pf.created_at, pf.created_by
                        FROM project_folders pf
                        INNER JOIN folder_tree ft ON pf.parent_folder_id = ft.id
                    )
                    SELECT * FROM folder_tree
                """, (folder_id, project_id))
                folders = cur.fetchall()

                folder_ids = [f['id'] for f in folders]
                for f in folders:
                    cur.execute("""
                        INSERT INTO project_folders_recycle_bin
                        (original_id, project_id, name, parent_folder_id, original_parent_id, created_at, created_by, deleted_at, expires_at)
                        VALUES (%s, %s, %s, %s, %s, %s, %s, NOW(), NOW() + INTERVAL '3 days')
                    """, (
                        f['id'], project_id, f['name'], f['parent_folder_id'], f['parent_folder_id'],
                        f['created_at'], f['created_by']
                    ))

                if folder_ids:
                    placeholders = ','.join(['%s'] * len(folder_ids))
                    cur.execute(f"""
                        SELECT id, original_name, file_size, stored_path, file_hash, version, uploaded_by, folder_id
                        FROM project_files
                        WHERE project_id = %s AND folder_id IN ({placeholders})
                    """, [project_id] + folder_ids)
                    files = cur.fetchall()
                    for f in files:
                        cur.execute("""
                            INSERT INTO project_recycle_bin 
                            (original_table, original_id, project_id, folder_id, file_name, original_name, file_size, stored_path, file_hash, version, uploaded_by, deleted_at, expires_at)
                            VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, NOW(), NOW() + INTERVAL '3 days')
                        """, (
                            'project_files', f['id'], project_id, f['folder_id'],
                            f['original_name'], f['original_name'], f['file_size'],
                            f['stored_path'], f['file_hash'], f['version'],
                            f['uploaded_by']
                        ))
                    cur.execute(f"""
                        DELETE FROM project_files
                        WHERE project_id = %s AND folder_id IN ({placeholders})
                    """, [project_id] + folder_ids)

                cur.execute(f"""
                    DELETE FROM project_folders
                    WHERE project_id = %s AND id IN ({','.join(['%s']*len(folder_ids))})
                """, [project_id] + folder_ids)

                conn.commit()
                return jsonify({
                    "success": True,
                    "folders_moved": len(folders),
                    "files_moved": len(files) if files else 0
                })

@app.route('/admin/projects/<int:project_id>/folders/<int:folder_id>/rename', methods=['PUT'])
def rename_folder(project_id, folder_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not can_edit_folder(project_id, folder_id, user_id):
        return jsonify({"error": "Permission denied"}), 403
    data = request.get_json()
    new_name = data.get('name', '').strip()
    if not new_name:
        return jsonify({"error": "Folder name required"}), 400
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT parent_folder_id FROM project_folders WHERE id = %s AND project_id = %s",
                        (folder_id, project_id))
            row = cur.fetchone()
            if not row:
                return jsonify({"error": "Folder not found"}), 404
            parent_id = row[0]
            cur.execute(
                "SELECT id FROM project_folders WHERE project_id = %s AND parent_folder_id = %s AND name = %s AND id != %s",
                (project_id, parent_id, new_name, folder_id))
            if cur.fetchone():
                return jsonify({"error": "A folder with this name already exists in this location"}), 400
            cur.execute("UPDATE project_folders SET name = %s WHERE id = %s", (new_name, folder_id))
            conn.commit()
            return jsonify({"success": True})

# Project files management
os.makedirs(PROJECT_FILES_ROOT, exist_ok=True)

def get_project_file_path(project_id, unique_filename):
    project_dir = os.path.join(PROJECT_FILES_ROOT, str(project_id))
    os.makedirs(project_dir, exist_ok=True)
    return os.path.join(project_dir, unique_filename)

@app.route('/admin/projects/<int:project_id>/folders/<int:folder_id>/upload', methods=['POST'])
def upload_project_file(project_id, folder_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not is_admin() and not can_access_project(project_id, user_id):
        return jsonify({"error": "Access denied"}), 403

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT status FROM projects WHERE id = %s", (project_id,))
            row = cur.fetchone()
            if not row or row[0] != 'active':
                return jsonify({"error": "Project is not active. Cannot upload."}), 400

    if 'file' not in request.files:
        return jsonify({"error": "No file"}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "Empty filename"}), 400

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT id FROM project_folders WHERE id = %s AND project_id = %s", (folder_id, project_id))
            if not cur.fetchone():
                return jsonify({"error": "Folder not found"}), 404

    original_name = file.filename
    file_bytes = file.read()
    file_hash = hashlib.sha256(file_bytes).hexdigest()
    file.seek(0)

    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("SELECT id, original_name, stored_path, version, folder_id FROM project_files WHERE project_id = %s AND file_hash = %s", (project_id, file_hash))
            duplicate = cur.fetchone()
            if duplicate:
                return jsonify({
                    "duplicate": True,
                    "existing_file": {
                        "id": duplicate['id'],
                        "original_name": duplicate['original_name'],
                        "folder_id": duplicate['folder_id'],
                        "version": duplicate['version']
                    },
                    "new_filename": original_name
                })

    ext = os.path.splitext(original_name)[1]
    unique_name = f"{uuid.uuid4().hex}{ext}"
    stored_path = get_project_file_path(project_id, unique_name)
    # Save the binary file
    file.save(stored_path)
    file_size = os.path.getsize(stored_path)

    # Extract text content from the file for knowledge base and search
    fake_file = FileStorage(BytesIO(file_bytes), filename=original_name)
    text_content, _ = extract_text_from_file(fake_file)
    if not text_content or text_content.startswith("["):
        text_content = ""  # fallback

    with get_db_connection() as conn:
        with db_transaction(conn):
            with conn.cursor() as cur:
                cur.execute("""
                    INSERT INTO project_files (project_id, folder_id, filename, original_name, file_size,
                                               stored_path, uploaded_by, file_hash, content)
                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s)
                    RETURNING id
                """, (project_id, folder_id, unique_name, original_name, file_size, stored_path, user_id, file_hash, text_content))
                file_id = cur.fetchone()[0]
                conn.commit()
                return jsonify({"success": True, "file_id": file_id, "original_name": original_name, "version": 1})

@app.route('/admin/projects/<int:project_id>/files/<int:file_id>/new_version', methods=['POST'])
def new_file_version(project_id, file_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not is_admin() and not can_access_project(project_id, user_id):
        return jsonify({"error": "Access denied"}), 403
    if 'file' not in request.files:
        return jsonify({"error": "No file"}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "Empty filename"}), 400

    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                        SELECT stored_path, version, original_name, folder_id, filename
                        FROM project_files
                        WHERE id = %s
                          AND project_id = %s
                        """, (file_id, project_id))
            existing = cur.fetchone()
            if not existing:
                return jsonify({"error": "File not found"}), 404

            original_name = file.filename
            file_bytes = file.read()
            file_hash = hashlib.sha256(file_bytes).hexdigest()
            file.seek(0)

            # Only extract text for supported text/office files; for images, just store.
            ext = os.path.splitext(original_name)[1].lower()
            text_extensions = {'.txt', '.md', '.text', '.csv', '.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt'}
            if ext in text_extensions:
                # Only attempt text extraction for office documents, not for images
                if file.filename.lower().endswith(('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp')):
                    file_content = "[Image file – no text extracted]"
                else:
                    file_content, _ = extract_text_from_file(file)
                    if not file_content or file_content.startswith("["):
                        return jsonify({"error": "Could not extract text from new version"}), 400
            else:
                # For images, audio, video, etc., just store empty content
                file_content = ""

            ext = os.path.splitext(original_name)[1]
            unique_name = f"{uuid.uuid4().hex}{ext}"
            stored_path = get_project_file_path(project_id, unique_name)
            file.save(stored_path)
            file_size = os.path.getsize(stored_path)
            new_version = existing['version'] + 1

            cur.execute("""
                        INSERT INTO project_file_versions (file_id, version, stored_path, file_size, uploaded_by)
                        VALUES (%s, %s, %s, %s, %s)
                        """, (file_id, existing['version'], existing['stored_path'], file_size, user_id))

            cur.execute("""
                        UPDATE project_files
                        SET version       = %s,
                            stored_path   = %s,
                            file_size     = %s,
                            uploaded_at   = NOW(),
                            uploaded_by   = %s,
                            file_hash     = %s,
                            original_name = %s,
                            content       = %s
                        WHERE id = %s
                        """,
                        (new_version, stored_path, file_size, user_id, file_hash, original_name, file_content, file_id))

            cur.execute("""
                        INSERT INTO project_file_usage (file_id, user_id, action, details)
                        VALUES (%s, %s, 'new_version', %s)
                        """, (file_id, user_id, json.dumps({'version': new_version, 'size': file_size})))

            conn.commit()
            return jsonify({"success": True, "file_id": file_id, "original_name": original_name, "version": new_version})

@app.route('/admin/projects/<int:project_id>/folders/<int:folder_id>/files', methods=['GET'])
def list_project_files(project_id, folder_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not is_admin() and not can_access_project(project_id, user_id):
        return jsonify({"error": "Access denied"}), 403
    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                        SELECT f.id,
                               f.original_name,
                               f.file_size,
                               f.version,
                               f.uploaded_at,
                               f.uploaded_by,
                               (SELECT username FROM users WHERE user_id = f.uploaded_by)        as uploaded_by_name,
                               (SELECT COUNT(*) FROM project_file_versions WHERE file_id = f.id) as version_count
                        FROM project_files f
                        WHERE f.project_id = %s
                          AND f.folder_id = %s
                        ORDER BY f.uploaded_at DESC
                        """, (project_id, folder_id))
            files = cur.fetchall()
            result = []
            for f in files:
                f['has_versions'] = f['version_count'] > 0
                f['can_move'] = can_move_file(project_id, f['id'], user_id)
                f['can_delete'] = can_edit_file(project_id, f['id'], user_id)
                f['can_rename'] = can_edit_file(project_id, f['id'], user_id)
                f['can_download'] = True
                f['file_size_kb'] = round(f['file_size'] / 1024, 1)
                f['uploaded_at_str'] = f['uploaded_at'].strftime('%Y-%m-%d %H:%M:%S')
                result.append(f)
            return jsonify({"files": result})

@app.route('/admin/projects/<int:project_id>/files', methods=['GET'])
def list_root_files(project_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not is_admin() and not can_access_project(project_id, user_id):
        return jsonify({"error": "Access denied"}), 403
    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                        SELECT f.id,
                               f.original_name,
                               f.file_size,
                               f.version,
                               f.uploaded_at,
                               f.uploaded_by,
                               (SELECT username FROM users WHERE user_id = f.uploaded_by)        as uploaded_by_name,
                               (SELECT COUNT(*) FROM project_file_versions WHERE file_id = f.id) as version_count
                        FROM project_files f
                        WHERE f.project_id = %s
                          AND f.folder_id IS NULL
                        ORDER BY f.uploaded_at DESC
                        """, (project_id,))
            files = cur.fetchall()
            result = []
            for f in files:
                f['has_versions'] = f['version_count'] > 0
                f['can_move'] = can_move_file(project_id, f['id'], user_id)
                f['can_delete'] = can_edit_file(project_id, f['id'], user_id)
                f['can_rename'] = can_edit_file(project_id, f['id'], user_id)
                f['can_download'] = True
                f['file_size_kb'] = round(f['file_size'] / 1024, 1)
                f['uploaded_at_str'] = f['uploaded_at'].strftime('%Y-%m-%d %H:%M:%S')
                result.append(f)
            return jsonify({"files": result})

@app.route('/admin/projects/<int:project_id>/files/<int:file_id>/versions', methods=['GET'])
def get_file_versions(project_id, file_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not is_admin() and not can_access_project(project_id, user_id):
        return jsonify({"error": "Access denied"}), 403
    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                        SELECT version,
                               file_size,
                               uploaded_at,
                               uploaded_by,
                               (SELECT username FROM users WHERE user_id = fv.uploaded_by) as uploaded_by_name
                        FROM project_file_versions fv
                        WHERE file_id = %s
                        ORDER BY version DESC
                        """, (file_id,))
            versions = cur.fetchall()
            return jsonify({"versions": versions})

@app.route('/admin/projects/<int:project_id>/files/<int:file_id>/download', methods=['GET'])
def download_project_file(project_id, file_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not is_admin() and not can_access_project(project_id, user_id):
        return jsonify({"error": "Access denied"}), 403
    version = request.args.get('version', type=int)
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            if version:
                cur.execute(
                    "SELECT stored_path, original_name FROM project_file_versions WHERE file_id = %s AND version = %s",
                    (file_id, version))
            else:
                cur.execute("SELECT stored_path, original_name FROM project_files WHERE id = %s", (file_id,))
            row = cur.fetchone()
            if not row:
                return jsonify({"error": "File not found"}), 404
            stored_path, original_name = row
    return send_file(stored_path, as_attachment=True, download_name=original_name)

@app.route('/admin/projects/<int:project_id>/files/<int:file_id>/comments', methods=['GET'])
def get_file_comments(project_id, file_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not is_admin() and not can_access_project(project_id, user_id):
        return jsonify({"error": "Access denied"}), 403
    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                        SELECT c.id, c.comment, c.created_at, u.username
                        FROM project_file_comments c
                                 JOIN users u ON c.user_id = u.user_id
                        WHERE c.file_id = %s
                        ORDER BY c.created_at ASC
                        """, (file_id,))
            comments = cur.fetchall()
            return jsonify({"comments": comments})

@app.route('/admin/projects/<int:project_id>/files/<int:file_id>/comments', methods=['POST'])
def add_file_comment(project_id, file_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not is_admin() and not can_access_project(project_id, user_id):
        return jsonify({"error": "Access denied"}), 403
    data = request.get_json()
    comment = data.get('comment', '').strip()
    if not comment:
        return jsonify({"error": "Comment cannot be empty"}), 400
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("INSERT INTO project_file_comments (file_id, user_id, comment) VALUES (%s, %s, %s)",
                        (file_id, user_id, comment))
            conn.commit()
            return jsonify({"success": True})

@app.route('/admin/projects/<int:project_id>/files/<int:file_id>/move', methods=['POST'])
def move_file(project_id, file_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not can_move_file(project_id, file_id, user_id):
        return jsonify({"error": "Permission denied"}), 403
    data = request.get_json()
    target_folder_id = data.get('folder_id')
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            if target_folder_id:
                cur.execute("SELECT id FROM project_folders WHERE id = %s AND project_id = %s", (target_folder_id, project_id))
                if not cur.fetchone():
                    return jsonify({"error": "Target folder not found in this project"}), 404
            cur.execute("UPDATE project_files SET folder_id = %s WHERE id = %s AND project_id = %s", (target_folder_id, file_id, project_id))
            if cur.rowcount == 0:
                return jsonify({"error": "File not found"}), 404
            conn.commit()
            return jsonify({"success": True})

@app.route('/admin/projects/<int:project_id>/files/batch_move', methods=['POST'])
def batch_move_files(project_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "Not logged in"}), 401

    data = request.get_json()
    file_ids = data.get('file_ids', [])
    target_folder_id = data.get('folder_id')
    if not file_ids:
        return jsonify({"error": "No files selected"}), 400
    if not target_folder_id:
        return jsonify({"error": "Target folder required"}), 400

    role = get_user_role_in_project(project_id, user_id)
    if not role and not is_admin():
        return jsonify({"error": "You are not a member of this project"}), 403

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT id FROM project_folders WHERE id = %s AND project_id = %s", (target_folder_id, project_id))
            if not cur.fetchone():
                return jsonify({"error": "Target folder not found in this project"}), 404

            placeholders = ','.join(['%s'] * len(file_ids))
            cur.execute(f"""
                SELECT id FROM project_files 
                WHERE id IN ({placeholders}) AND project_id = %s
            """, file_ids + [project_id])
            found = cur.fetchall()
            if len(found) != len(file_ids):
                return jsonify({"error": "Some files not found in this project"}), 404

            cur.execute(f"""
                UPDATE project_files SET folder_id = %s 
                WHERE id IN ({placeholders}) AND project_id = %s
            """, [target_folder_id] + file_ids + [project_id])
            conn.commit()
            return jsonify({"success": True, "moved_count": len(file_ids)})

@app.route('/admin/projects/<int:project_id>/batch_download', methods=['POST'])
def batch_download_files(project_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not is_admin() and not can_access_project(project_id, user_id):
        return jsonify({"error": "Access denied"}), 403
    data = request.get_json()
    file_ids = data.get('file_ids', [])
    if not file_ids:
        return jsonify({"error": "No files selected"}), 400
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            placeholders = ','.join(['%s'] * len(file_ids))
            cur.execute(
                f"SELECT stored_path, original_name FROM project_files WHERE id IN ({placeholders}) AND project_id = %s",
                file_ids + [project_id])
            files = cur.fetchall()
            if not files:
                return jsonify({"error": "No valid files found"}), 404
            zip_buffer = BytesIO()
            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for stored_path, original_name in files:
                    zipf.write(stored_path, original_name)
            zip_buffer.seek(0)
            return send_file(zip_buffer, as_attachment=True, download_name=f"project_{project_id}_files.zip",
                             mimetype='application/zip')

@app.route('/admin/projects/<int:project_id>/files/search', methods=['GET'])
def search_project_files(project_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not is_admin() and not can_access_project(project_id, user_id):
        return jsonify({"error": "Access denied"}), 403
    query = request.args.get('q', '').strip()
    if len(query) < 2:
        return jsonify({"files": []})
    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                        SELECT f.id, f.original_name, f.file_size, f.uploaded_at, fo.name as folder_name
                        FROM project_files f
                                 LEFT JOIN project_folders fo ON f.folder_id = fo.id
                        WHERE f.project_id = %s
                          AND f.original_name ILIKE %s
                        ORDER BY f.uploaded_at DESC
                        LIMIT 50
                        """, (project_id, f'%{query}%'))
            files = cur.fetchall()
            for f in files:
                f['file_size_kb'] = round(f['file_size'] / 1024, 1)
            return jsonify({"files": files})

@app.route('/admin/projects/<int:project_id>/folders/<int:folder_id>/comments', methods=['GET'])
def get_folder_comments(project_id, folder_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not is_admin() and not can_access_project(project_id, user_id):
        return jsonify({"error": "Access denied"}), 403
    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                        SELECT c.id, c.comment, c.created_at, u.username
                        FROM project_folder_comments c
                                 JOIN users u ON c.user_id = u.user_id
                        WHERE c.folder_id = %s
                        ORDER BY c.created_at ASC
                        """, (folder_id,))
            comments = cur.fetchall()
            return jsonify({"comments": comments})

@app.route('/admin/projects/<int:project_id>/folders/<int:folder_id>/comments', methods=['POST'])
def add_folder_comment(project_id, folder_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not is_admin() and not can_access_project(project_id, user_id):
        return jsonify({"error": "Access denied"}), 403
    data = request.get_json()
    comment = data.get('comment', '').strip()
    if not comment:
        return jsonify({"error": "Comment cannot be empty"}), 400
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("INSERT INTO project_folder_comments (folder_id, user_id, comment) VALUES (%s, %s, %s)",
                        (folder_id, user_id, comment))
            conn.commit()
            return jsonify({"success": True})

@app.route('/admin/projects/<int:project_id>/files/<int:file_id>/rename', methods=['PUT'])
def rename_project_file(project_id, file_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not can_edit_file(project_id, file_id, user_id):
        return jsonify({"error": "Permission denied"}), 403
    data = request.get_json()
    new_name = data.get('original_name', '').strip()
    if not new_name:
        return jsonify({"error": "New name required"}), 400
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("UPDATE project_files SET original_name = %s WHERE id = %s AND project_id = %s",
                        (new_name, file_id, project_id))
            if cur.rowcount == 0:
                return jsonify({"error": "File not found"}), 404
            conn.commit()
    return jsonify({"success": True})


# ========== Knowledge Lab Routes ==========
@app.route('/knowledge_lab/upload', methods=['POST'])
@login_required  # you need to define this decorator or just check session
def upload_knowledge_lab_file():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "Not logged in"}), 401

    if 'file' not in request.files:
        return jsonify({"error": "No file"}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "Empty filename"}), 400

    # Extract text
    file_bytes = file.read()
    file.seek(0)
    file_hash = hashlib.sha256(file_bytes).hexdigest()

    from io import BytesIO
    from werkzeug.datastructures import FileStorage
    fake_file = FileStorage(BytesIO(file_bytes), filename=file.filename)
    text_content, _ = extract_text_from_file(fake_file)
    if not text_content or text_content.startswith("["):
        text_content = ""

    # Save file permanently
    KNOWLEDGE_LAB_DIR = os.path.join(os.getcwd(), 'knowledge_lab_files')
    os.makedirs(KNOWLEDGE_LAB_DIR, exist_ok=True)
    unique_name = f"{file_hash}_{int(time.time())}_{file.filename}"
    stored_path = os.path.join(KNOWLEDGE_LAB_DIR, unique_name)
    with open(stored_path, 'wb') as f:
        f.write(file_bytes)

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            # Check for duplicate by hash
            cur.execute("SELECT id FROM knowledge_lab_files WHERE file_hash = %s", (file_hash,))
            existing = cur.fetchone()
            if existing:
                return jsonify({"error": "File already exists in knowledge lab", "file_id": existing[0]}), 409
            cur.execute("""
                        INSERT INTO knowledge_lab_files (user_id, filename, original_name, file_size, content,
                                                         file_hash, stored_path)
                        VALUES (%s, %s, %s, %s, %s, %s, %s)
                        RETURNING id
                        """,
                        (user_id, file.filename, file.filename, len(file_bytes), text_content, file_hash, stored_path))
            new_id = cur.fetchone()[0]
            conn.commit()
            return jsonify({
                "success": True,
                "file_id": new_id,
                "filename": file.filename,
                "file_size": len(file_bytes),
                "uploaded_at": datetime.now(timezone.utc).isoformat()
            })

@app.route('/knowledge_lab/list', methods=['GET'])
def list_knowledge_lab_files():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "Not logged in"}), 401

    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                        SELECT id, filename, original_name, file_size, uploaded_at
                        FROM knowledge_lab_files
                        WHERE user_id = %s
                        ORDER BY uploaded_at DESC
                        """, (user_id,))
            files = cur.fetchall()
            return jsonify({"files": files})

@app.route('/knowledge_lab/content/<int:file_id>', methods=['GET'])
def get_knowledge_lab_content(file_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "Not logged in"}), 401
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT content, original_name FROM knowledge_lab_files WHERE id = %s AND user_id = %s", (file_id, user_id))
            row = cur.fetchone()
            if not row:
                return jsonify({"error": "File not found"}), 404
            return jsonify({"content": row[0], "filename": row[1]})

@app.route('/knowledge_lab/delete/<int:file_id>', methods=['POST'])
def delete_knowledge_lab_file(file_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "Not logged in"}), 401

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT stored_path FROM knowledge_lab_files WHERE id = %s AND user_id = %s",
                        (file_id, user_id))
            row = cur.fetchone()
            if not row:
                return jsonify({"error": "File not found"}), 404
            stored_path = row[0]
            if stored_path and os.path.exists(stored_path):
                os.remove(stored_path)
            cur.execute("DELETE FROM knowledge_lab_files WHERE id = %s AND user_id = %s", (file_id, user_id))
            conn.commit()
            return jsonify({"success": True})

@app.route('/company_kb/upload', methods=['POST'])
def upload_company_kb_file():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not user_id or session.get('role') != 'admin':
        return jsonify({"error": "Admin access required"}), 403

    if 'file' not in request.files:
        return jsonify({"error": "No file"}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "Empty filename"}), 400

    category = request.form.get('category', '').strip()
    if not category:
        return jsonify({"error": "Category is required"}), 400

    # Read file bytes
    file_bytes = file.read()
    file_hash = hashlib.sha256(file_bytes).hexdigest()

    # Extract text
    from io import BytesIO
    from werkzeug.datastructures import FileStorage
    fake_file = FileStorage(BytesIO(file_bytes), filename=file.filename)
    text_content, _ = extract_text_from_file(fake_file)
    if not text_content or text_content.startswith("["):
        text_content = ""

    # Permanent storage
    COMPANY_KB_DIR = os.path.join(os.getcwd(), 'company_kb_files')
    os.makedirs(COMPANY_KB_DIR, exist_ok=True)
    unique_name = f"{file_hash}_{int(time.time())}_{file.filename}"
    stored_path = os.path.join(COMPANY_KB_DIR, unique_name)
    with open(stored_path, 'wb') as f:
        f.write(file_bytes)

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            # Check for existing file by hash
            cur.execute("SELECT id, stored_path FROM company_knowledge_base WHERE file_hash = %s", (file_hash,))
            existing = cur.fetchone()
            if existing:
                # Overwrite: delete old physical file, update record
                old_path = existing[1]
                if old_path and os.path.exists(old_path):
                    try:
                        os.remove(old_path)
                    except Exception as e:
                        logger.warning(f"Could not delete old file: {e}")
                cur.execute("""
                    UPDATE company_knowledge_base
                    SET filename = %s, original_name = %s, file_size = %s, content = %s,
                        stored_path = %s, category = %s, uploaded_by = %s, updated_at = NOW()
                    WHERE id = %s
                """, (file.filename, file.filename, len(file_bytes), text_content, stored_path, category, user_id, existing[0]))
                conn.commit()
                return jsonify({"success": True, "file_id": existing[0], "filename": file.filename, "category": category, "updated": True})
            else:
                # New file
                cur.execute("""
                    INSERT INTO company_knowledge_base (filename, original_name, file_size, content, file_hash, stored_path, category, uploaded_by)
                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
                    RETURNING id
                """, (file.filename, file.filename, len(file_bytes), text_content, file_hash, stored_path, category, user_id))
                new_id = cur.fetchone()[0]
                conn.commit()
                return jsonify({"success": True, "file_id": new_id, "filename": file.filename, "category": category})

@app.route('/company_kb/list', methods=['GET'])
def list_company_kb_files():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "Not logged in"}), 401

    category = request.args.get('category', '')
    search = request.args.get('search', '').strip()
    page = int(request.args.get('page', 1))
    per_page = int(request.args.get('per_page', 50))
    offset = (page - 1) * per_page

    ts_config = 'simple'

    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            # Build where clause and base parameters (for count query)
            where_parts = []
            count_params = []
            if category:
                where_parts.append("category = %s")
                count_params.append(category)
            if search:
                where_parts.append(f"to_tsvector('{ts_config}', content) @@ plainto_tsquery('{ts_config}', %s)")
                count_params.append(search)

            where_clause = "WHERE " + " AND ".join(where_parts) if where_parts else ""

            # Count total
            count_query = f"SELECT COUNT(*) as total FROM company_knowledge_base {where_clause}"
            cur.execute(count_query, count_params)
            total = cur.fetchone()['total']

            # Build main query parameters
            if search:
                # We need two copies of the search term: one for the rank function, one for the where clause
                main_params = []
                if category:
                    main_params.append(category)
                main_params.append(search)   # for rank
                main_params.append(search)   # for where
                # Add pagination
                main_params.extend([per_page, offset])

                query = f"""
                    SELECT id, original_name as filename, file_size, category, uploaded_at,
                           (SELECT username FROM users WHERE user_id = uploaded_by) as uploaded_by_name,
                           ts_rank(to_tsvector('{ts_config}', content), plainto_tsquery('{ts_config}', %s)) as rank
                    FROM company_knowledge_base
                    {where_clause}
                    ORDER BY rank DESC, uploaded_at DESC
                    LIMIT %s OFFSET %s
                """
            else:
                main_params = count_params.copy()
                main_params.extend([per_page, offset])
                query = f"""
                    SELECT id, original_name as filename, file_size, category, uploaded_at,
                           (SELECT username FROM users WHERE user_id = uploaded_by) as uploaded_by_name
                    FROM company_knowledge_base
                    {where_clause}
                    ORDER BY uploaded_at DESC
                    LIMIT %s OFFSET %s
                """

            cur.execute(query, main_params)
            files = cur.fetchall()

            return jsonify({
                "files": files,
                "total": total,
                "page": page,
                "per_page": per_page,
                "has_next": offset + per_page < total
            })

@app.route('/company_kb/search', methods=['GET'])
def search_company_kb():
    query = request.args.get('q', '').strip()
    if len(query) < 2:
        return jsonify({"results": []})

    ts_config = 'simple'
    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute(f"""
                SELECT id, original_name as filename, category,
                       ts_headline('{ts_config}', content, plainto_tsquery('{ts_config}', %s), 'MaxWords=30, MinWords=15') as snippet
                FROM company_knowledge_base
                WHERE to_tsvector('{ts_config}', content) @@ plainto_tsquery('{ts_config}', %s)
                ORDER BY ts_rank(to_tsvector('{ts_config}', content), plainto_tsquery('{ts_config}', %s)) DESC
                LIMIT 20
            """, (query, query, query))
            results = cur.fetchall()
            return jsonify({"results": results})

@app.route('/company_kb/content/<int:file_id>', methods=['GET'])
def get_company_kb_content(file_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "Not logged in"}), 401

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT content, original_name FROM company_knowledge_base WHERE id = %s", (file_id,))
            row = cur.fetchone()
            if not row:
                return jsonify({"error": "File not found"}), 404
            return jsonify({"content": row[0], "filename": row[1]})

@app.route('/company_kb/delete/<int:file_id>', methods=['POST'])
def delete_company_kb_file(file_id):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not user_id or session.get('role') != 'admin':
        return jsonify({"error": "Admin access required"}), 403

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT stored_path FROM company_knowledge_base WHERE id = %s", (file_id,))
            row = cur.fetchone()
            if not row:
                return jsonify({"error": "File not found"}), 404
            stored_path = row[0]
            if stored_path and os.path.exists(stored_path):
                os.remove(stored_path)
            cur.execute("DELETE FROM company_knowledge_base WHERE id = %s", (file_id,))
            conn.commit()
            return jsonify({"success": True})

@app.route('/company_kb/categories', methods=['GET'])
def get_company_kb_categories():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "Not logged in"}), 401

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT DISTINCT category FROM company_knowledge_base WHERE category IS NOT NULL AND category != '' ORDER BY category")
            rows = cur.fetchall()
            categories = [row[0] for row in rows]
            return jsonify({"categories": categories})

# File station routes
@app.route('/upload_file', methods=['POST'])
def upload_file():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403

    if 'file' not in request.files:
        return jsonify({"error": "No file"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "Empty filename"}), 400

    if not allowed_file(file.filename):
        return jsonify({"error": f"不支持的文件类型: {file.filename}"}), 400

    user_id = get_user_id()
    thread_id = session.get('thread_id')
    if not thread_id:
        thread_id = str(uuid.uuid4())
        session['thread_id'] = thread_id
        get_or_create_session(thread_id)

    # Read file bytes once
    file_bytes = file.read()
    file_hash = hashlib.sha256(file_bytes).hexdigest()

    # Extract text from file (for storage in user_files.content)
    from io import BytesIO
    from werkzeug.datastructures import FileStorage
    fake_file = FileStorage(BytesIO(file_bytes), filename=file.filename)
    extracted_text, _ = extract_text_from_file(fake_file)
    if not extracted_text or extracted_text.startswith("["):
        extracted_text = ""   # fallback

    if session.get('consent_value', 0) != 1:
        # Anonymous: store only as temp file (no DB)
        anon_id = user_id
        temp_dir = get_anon_temp_dir(anon_id)
        existing_files = [f for f in os.listdir(temp_dir) if f.endswith('.txt')]
        if len(existing_files) >= 5:
            return jsonify({"error": "Anonymous users can only store up to 5 files."}), 400
        if len(file_bytes) > 5 * 1024 * 1024:
            return jsonify({"error": "File exceeds 5MB limit for anonymous users."}), 400

        # For anonymous, we still want to cache the content for the current session
        add_to_cache(thread_id, file.filename, extracted_text, user_id)
        # Also store as .txt file in temp dir
        safe_name = re.sub(r'[^\w\-_\. ]', '_', file.filename) + '.txt'
        file_path = os.path.join(temp_dir, safe_name)
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(extracted_text)
        return jsonify({"success": True, "filename": file.filename})

    # Registered user: check for existing file by hash
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT id, filename, original_stored_path FROM user_files WHERE user_id = %s AND file_hash = %s", (user_id, file_hash))
            existing = cur.fetchone()
            if existing and request.form.get('force') != 'true':
                return jsonify({
                    "exists": True,
                    "file_id": existing[0],
                    "filename": existing[1],
                    "original_path": existing[2] if existing[2] else None
                })

    ext = os.path.splitext(file.filename)[1]
    unique_name = f"{file_hash}_{int(time.time())}{ext}"
    original_dir = os.path.join(USER_FILES_ORIGINAL_ROOT, user_id)
    os.makedirs(original_dir, exist_ok=True)
    original_path = os.path.join(original_dir, unique_name)
    # Save original binary file
    with open(original_path, 'wb') as f:
        f.write(file_bytes)

    # Add to in‑memory cache
    add_to_cache(thread_id, file.filename, extracted_text, user_id)
    record_file_usage(thread_id, file.filename, 'standalone_upload', "上传文件供日后使用")

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            if existing and request.form.get('force') == 'true':
                old_path = existing[2]
                if old_path and os.path.exists(old_path):
                    try:
                        os.remove(old_path)
                    except Exception:
                        pass
                cur.execute("""
                    UPDATE user_files
                    SET filename = %s,
                        size_bytes = %s,
                        original_stored_path = %s,
                        file_hash = %s,
                        expires_at = NULL,
                        original_expires_at = NOW() + INTERVAL '3 days',
                        original_name = %s,
                        content = %s
                    WHERE id = %s
                """, (file.filename, len(file_bytes), original_path, file_hash, file.filename, extracted_text, existing[0]))
            else:
                ensure_user_exists(user_id)
                cur.execute("""
                    INSERT INTO user_files (user_id, thread_id, filename, size_bytes, expires_at,
                                            original_stored_path, file_hash, original_expires_at, original_name, content)
                    VALUES (%s, %s, %s, %s, NULL, %s, %s, NOW() + INTERVAL '3 days', %s, %s)
                    ON CONFLICT (thread_id, filename) DO UPDATE SET
                        size_bytes = EXCLUDED.size_bytes,
                        original_stored_path = EXCLUDED.original_stored_path,
                        file_hash = EXCLUDED.file_hash,
                        original_expires_at = EXCLUDED.original_expires_at,
                        original_name = EXCLUDED.original_name,
                        content = EXCLUDED.content
                """, (user_id, thread_id, file.filename, len(file_bytes), original_path, file_hash, file.filename, extracted_text))
            conn.commit()

    return jsonify({"success": True, "filename": file.filename})

@app.route('/download_original_file', methods=['POST'])
def download_original_file():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403

    data = request.get_json()
    filename = data.get('filename')
    if not filename:
        return jsonify({"error": "Missing filename"}), 400

    user_id = get_user_id()
    thread_id = session.get('thread_id')
    if not thread_id:
        return jsonify({"error": "No active session"}), 400

    if session.get('consent_value', 0) != 1:
        return jsonify({
            "error": "anonymous_not_allowed",
            "message": "匿名用户无法下载原文件。请注册或登录账户后使用此功能。"
        }), 403

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT original_stored_path
                FROM user_files
                WHERE user_id = %s AND thread_id = %s AND filename = %s
                  AND (original_expires_at IS NULL OR original_expires_at > NOW())
            """, (user_id, thread_id, filename))
            row = cur.fetchone()
            if not row or not row[0]:
                return jsonify({"error": "Original file not found or expired"}), 404
            original_path = row[0]
            if not os.path.exists(original_path):
                return jsonify({"error": "File missing on server"}), 404
            return send_file(original_path, as_attachment=True, download_name=filename)

@app.route('/delete_file_station', methods=['POST'])
def delete_file_station():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403

    data = request.get_json()
    file_id = data.get('file_id')
    if not file_id:
        return jsonify({"error": "Missing file_id"}), 400

    user_id = get_user_id()

    if session.get('consent_value', 0) != 1:
        anon_id = user_id
        temp_dir = get_anon_temp_dir(anon_id)
        fpath = os.path.join(temp_dir, file_id)
        if os.path.exists(fpath):
            os.remove(fpath)
            return jsonify({"success": True})
        else:
            return jsonify({"error": "File not found"}), 404

    with get_db_connection() as conn:
        with db_transaction(conn):
            with conn.cursor(cursor_factory=RealDictCursor) as cur:
                cur.execute("""
                    SELECT id, filename, original_name, content, size_bytes, original_stored_path, file_hash, thread_id, user_id
                    FROM user_files
                    WHERE id = %s AND user_id = %s
                """, (file_id, user_id))
                file_record = cur.fetchone()
                if not file_record:
                    return jsonify({"error": "File not found or not owned"}), 404

                cur.execute("""
                            INSERT INTO recycle_bin
                            (original_table, original_id, user_id, file_name, file_content, file_size,
                             original_stored_path, file_hash, thread_id, deleted_at, expires_at,
                             uploaded_by, deleted_by)
                            VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, NOW(), NOW() + INTERVAL '3 days',
                                    %s, %s)
                            """, (
                                'user_files', file_record['id'], user_id, file_record['original_name'],
                                file_record['content'], file_record['size_bytes'], file_record['original_stored_path'],
                                file_record['file_hash'], file_record['thread_id'],
                                file_record['user_id'],
                                user_id
                            ))

                cur.execute("DELETE FROM user_files WHERE id = %s AND user_id = %s", (file_id, user_id))

                conn.commit()
                return jsonify({"success": True, "moved_to_recycle_bin": True})

@app.route('/get_file_station', methods=['GET'])
def get_file_station():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = get_user_id()
    is_admin = session.get('role') == 'admin'

    if session.get('consent_value', 0) != 1:
        temp_dir = get_anon_temp_dir(user_id)
        files = []
        try:
            for fname in os.listdir(temp_dir):
                if fname.endswith('.txt'):
                    fpath = os.path.join(temp_dir, fname)
                    stat = os.stat(fpath)
                    original_name = fname[:-4] if fname.endswith('.txt') else fname
                    files.append({
                        "id": fname,
                        "filename": original_name,
                        "size_bytes": stat.st_size,
                        "created_at": datetime.fromtimestamp(stat.st_ctime).isoformat(),
                        "expires_at": None,
                        "usage": [],
                        "source": "user_file"
                    })
        except Exception as e:
            logger.error(f"Failed to list anon files: {e}")
            return jsonify({"error": "无法读取临时文件"}), 500
        return jsonify({"files": files, "is_admin": False})

    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            # User files
            cur.execute("""
                SELECT 
                    'user_file' as source,
                    uf.id::text as id,
                    uf.original_name as filename,
                    uf.size_bytes,
                    uf.created_at,
                    uf.expires_at,
                    uf.meta_data,
                    uf.user_id as owner_id,
                    (uf.user_id = %s) as can_delete,
                    NULL as project_name,
                    NULL as project_id,
                    NULL as folder_path,
                    (SELECT json_agg(
                        json_build_object(
                            'usage_type', fu.usage_type,
                            'question', fu.question,
                            'timestamp', fu.timestamp,
                            'thread_id', fu.thread_id
                        ) ORDER BY fu.timestamp DESC
                    ) FROM file_usage fu WHERE fu.user_id = uf.user_id AND fu.filename = uf.original_name LIMIT 10) as usage
                FROM user_files uf
                WHERE uf.user_id = %s AND (uf.expires_at IS NULL OR uf.expires_at > NOW())
                ORDER BY uf.created_at DESC
            """, (user_id, user_id))
            user_files = cur.fetchall()

            # Project files
            if is_admin:
                cur.execute("""
                    SELECT 
                        'project_file' as source,
                        pf.id::text as id,
                        pf.original_name as filename,
                        pf.file_size as size_bytes,
                        pf.uploaded_at as created_at,
                        NULL as expires_at,
                        p.name as project_name,
                        p.id as project_id,
                        (SELECT string_agg(f.name, '/') FROM project_folders f WHERE f.id = pf.folder_id) as folder_path,
                        NULL as usage
                    FROM project_files pf
                    JOIN projects p ON pf.project_id = p.id
                    ORDER BY pf.uploaded_at DESC
                """)
            else:
                cur.execute("""
                    SELECT 
                        'project_file' as source,
                        pf.id::text as id,
                        pf.original_name as filename,
                        pf.file_size as size_bytes,
                        pf.uploaded_at as created_at,
                        NULL as expires_at,
                        p.name as project_name,
                        p.id as project_id,
                        (SELECT string_agg(f.name, '/') FROM project_folders f WHERE f.id = pf.folder_id) as folder_path,
                        NULL as usage
                    FROM project_files pf
                    JOIN projects p ON pf.project_id = p.id
                    JOIN project_members pm ON p.id = pm.project_id
                    WHERE pm.user_id = %s
                    ORDER BY pf.uploaded_at DESC
                """, (user_id,))
            project_files = cur.fetchall()

    all_files = user_files + project_files
    return jsonify({"files": all_files, "is_admin": is_admin})

@app.route('/load_project_file', methods=['POST'])
def load_project_file():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = get_user_id()
    data = request.get_json()
    project_id = data.get('project_id')
    file_id = data.get('file_id')
    if not project_id or not file_id:
        return jsonify({"error": "Missing project_id or file_id"}), 400

    if not is_admin() and not can_access_project(project_id, user_id):
        return jsonify({"error": "Access denied"}), 403

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT stored_path, original_name
                FROM project_files
                WHERE id = %s AND project_id = %s
            """, (file_id, project_id))
            row = cur.fetchone()
            if not row:
                return jsonify({"error": "File not found"}), 404
            stored_path, original_name = row
            if not os.path.exists(stored_path):
                return jsonify({"error": "File missing on server"}), 404

            with open(stored_path, 'rb') as f:
                file_bytes = f.read()

            from io import BytesIO
            from werkzeug.datastructures import FileStorage
            fake_file = FileStorage(BytesIO(file_bytes), filename=original_name)
            text, _ = extract_text_from_file(fake_file)
            if not text or text.startswith("["):
                return jsonify({"error": "Could not extract text from file"}), 400

            return jsonify({"content": text, "filename": original_name})

# ---------- Batch compare endpoints ----------
@app.route('/compare_batch', methods=['POST'])
def compare_batch():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    if 'files' not in request.files:
        return jsonify({"error": "No files uploaded"}), 400
    files = request.files.getlist('files')
    if len(files) < 2:
        return jsonify({"error": "Need at least 2 files for comparison"}), 400
    if len(files) > 10:
        return jsonify({"error": "Maximum 10 files allowed"}), 400
    for f in files:
        if not allowed_file(f.filename):
            return jsonify({"error": f"不支持的文件类型: {f.filename}"}), 400

    user_id = get_user_id()
    thread_id = session['thread_id']
    success, busy_thread, busy_name = acquire_task_lock(user_id, thread_id, 'batch_compare')
    get_or_create_session(thread_id)
    if not success:
        return jsonify({
            "error": "resource_busy",
            "busy_chat": busy_name,
            "message": f"另一个资源密集型任务正在聊天“{busy_name}”中进行，请稍后再试。"
        }), 409

    try:
        template_file = request.files.get('template_file')
        template_text = None
        if template_file and template_file.filename:
            if not allowed_file(template_file.filename):
                return jsonify({"error": f"不支持的文件类型: {template_file.filename}"}), 400
            if session.get('consent_value', 0) == 1:
                template_text = get_or_extract_file_analysis(template_file, 'chat', user_id, thread_id=thread_id)
            else:
                template_text, _ = extract_text_from_file(template_file)
            if template_text and not template_text.startswith("["):
                if not is_valid_extracted_text(template_text):
                    template_text = None   # invalid template, ignore
                else:
                    record_file_usage(thread_id, template_file.filename, 'template_upload', "上传模板文件用于对比")

        check_items_json = request.form.get('check_items', '{}')
        try:
            check_items = json.loads(check_items_json)
        except Exception:
            check_items = {}

        defaults = {
            'text_sim': True,
            'key_info': True,
            'file_attr': True,
            'image_sim': True,
            'semantic': False
        }
        for k, v in defaults.items():
            if k not in check_items:
                check_items[k] = v

        if len(files) > 10 and check_items.get('semantic'):
            check_items['semantic'] = False
            logger.info("Semantic analysis disabled because number of files exceeds 10.")

        file_data = []
        for f in files:
            if not f.filename:
                continue
            if session.get('consent_value', 0) == 1:
                f.seek(0)
                file_bytes = f.read()
                file_hash = hashlib.sha256(file_bytes).hexdigest()
                file_size = len(file_bytes)
                f.seek(0)
                with get_db_connection() as conn:
                    with conn.cursor() as cur:
                        cur.execute("SELECT id FROM user_files WHERE user_id = %s AND file_hash = %s", (user_id, file_hash))
                        existing = cur.fetchone()
                        if not existing:
                            ext = os.path.splitext(f.filename)[1]
                            unique_name = f"{file_hash}_{int(time.time())}{ext}"
                            original_dir = os.path.join(USER_FILES_ORIGINAL_ROOT, user_id)
                            os.makedirs(original_dir, exist_ok=True)
                            original_path = os.path.join(original_dir, unique_name)
                            f.seek(0)
                            f.save(original_path)
                            ensure_user_exists(user_id)
                            cur.execute("""
                                INSERT INTO user_files (user_id, thread_id, filename, size_bytes, expires_at, original_stored_path, file_hash, original_expires_at, original_name, content)
                                VALUES (%s, %s, %s, %s, NULL, %s, %s, NOW() + INTERVAL '3 days', %s, '')
                                ON CONFLICT (thread_id, filename) DO UPDATE SET
                                    size_bytes = EXCLUDED.size_bytes,
                                    original_stored_path = EXCLUDED.original_stored_path,
                                    file_hash = EXCLUDED.file_hash,
                                    original_expires_at = EXCLUDED.original_expires_at,
                                    original_name = EXCLUDED.original_name,
                                    content = EXCLUDED.content
                            """, (user_id, thread_id, f.filename, file_size, original_path, file_hash, f.filename))
                            conn.commit()
                text = get_or_extract_file_analysis(f, 'chat', user_id, thread_id=thread_id)
            else:
                text, _ = extract_text_from_file(f)

            if text and not text.startswith("[") and is_valid_extracted_text(text):
                record_file_usage(thread_id, f.filename, 'compare_batch', "批量对比")
                # Store in memory cache so files are immediately available in chat file station
                add_to_cache(thread_id, f.filename, text, user_id)
                f.seek(0)
                meta = extract_metadata(f)
                images = extract_images_from_file(f)
                file_data.append({
                    'filename': f.filename,
                    'text': text,
                    'metadata': meta,
                    'images': images
                })
            else:
                logger.warning(f"Skipping file {f.filename}: extraction failed or invalid (text='{text}')")
                continue

        if len(file_data) < 2:
            return jsonify({"error": "Could not extract valid text from at least two files"}), 400

        # Store batch compare file texts in session so chat agent can reference them in follow-up questions
        session['batch_compare_files'] = [
            {'filename': fd['filename'], 'text': fd['text'][:3000]}  # truncated to avoid huge session
            for fd in file_data
        ]

        n = len(file_data)
        if check_items.get('text_sim', True) or check_items.get('key_info', True):
            vectorizer, tfidf_matrix = _precompute_tfidf_for_files(file_data, template_text)
        else:
            vectorizer = tfidf_matrix = None

        semantic_sim_matrix = None
        if check_items.get('semantic', False):
            all_texts = [fd['text'] for fd in file_data]
            semantic_sim_matrix = compute_batch_semantic_similarity(all_texts)
            logger.info("Semantic similarity matrix computed.")
        if semantic_sim_matrix is None and check_items.get('semantic', False):
            check_items['semantic'] = False
            logger.warning("Semantic analysis disabled due to model load failure.")

        pairs = []
        risk_matrix = [[0] * n for _ in range(n)]
        for i in range(n):
            for j in range(i + 1, n):
                text1 = file_data[i]['text']
                text2 = file_data[j]['text']
                meta1 = file_data[i]['metadata']
                meta2 = file_data[j]['metadata']
                images1 = file_data[i]['images']
                images2 = file_data[j]['images']
                img_sim = image_similarity(images1, images2) if check_items.get('image_sim', True) else 0.0

                if check_items.get('text_sim', True) and tfidf_matrix is not None:
                    sim = _compute_pair_similarity_from_matrix(tfidf_matrix, i, j)
                else:
                    sim = 0.0

                if check_items.get('key_info', True):
                    t1 = preprocess_text_for_similarity(text1)
                    t2 = preprocess_text_for_similarity(text2)
                    if template_text:
                        t1 = remove_template_content(t1, template_text)
                        t2 = remove_template_content(t2, template_text)
                    key_sim = keyword_overlap_similarity(t1, t2)
                else:
                    key_sim = 0.0

                if check_items.get('file_attr', True) and meta1 and meta2:
                    attr_sim = file_attr_similarity(meta1, meta2)
                else:
                    attr_sim = 0.0

                text_sim_val = sim * 100
                key_info_val = key_sim * 100
                file_attr_val = attr_sim
                img_sim_val = img_sim

                risk = 0.3 * key_info_val + 0.3 * file_attr_val + 0.2 * text_sim_val + 0.2 * img_sim_val

                _, html1, html2, blocks = compute_similarity_with_numbers(text1, text2, template_text)

                pair_info = {
                    'i': i, 'j': j,
                    'name1': file_data[i]['filename'],
                    'name2': file_data[j]['filename'],
                    'text1': text1,
                    'text2': text2,
                    'sim': sim * 100,
                    'risk': risk,
                    'blocks': blocks,
                    'html1': html1,
                    'html2': html2,
                    'used_weights': {},
                    'attr_same': 1 if meta1.get('author') and meta1['author'] == meta2.get('author') else 0
                }
                pairs.append(pair_info)
                risk_matrix[i][j] = risk
                risk_matrix[j][i] = risk

        key_info_matches = []
        for p in pairs:
            kw1 = set(extract_keywords(p['text1'], 20))
            kw2 = set(extract_keywords(p['text2'], 20))
            common = kw1 & kw2
            key_info_matches.append({
                'name1': p['name1'],
                'name2': p['name2'],
                'common_keywords': list(common)[:10]
            })

        attr_details = []
        for fd in file_data:
            meta = fd['metadata']
            attr_details.append({
                'filename': fd['filename'],
                'author': meta.get('author', ''),
                'creation_date': meta.get('creationDate', ''),
                'creator': meta.get('creator', ''),
                'producer': meta.get('producer', '')
            })

        batch_data = {
            'file_data': [{'filename': fd['filename'], 'metadata': fd['metadata']} for fd in file_data],
            'pairs': pairs,
            'check_items': check_items,
            'timestamp': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            'key_info_matches': key_info_matches,
            'attr_details': attr_details,
            'semantic_sim_matrix': semantic_sim_matrix,
        }

        temp_path = store_batch_comparison_temp(batch_data)
        session['batch_comparison_path'] = temp_path

        high_risk_files = []
        strong_alert_files = []
        for i in range(n):
            for j in range(i + 1, n):
                if risk_matrix[i][j] > 20:
                    strong_alert_files.extend([file_data[i]['filename'], file_data[j]['filename']])
                elif risk_matrix[i][j] > 10:
                    high_risk_files.extend([file_data[i]['filename'], file_data[j]['filename']])
        strong_alert_files = list(set(strong_alert_files))
        high_risk_files = list(set(high_risk_files) - set(strong_alert_files))

        short_names = [truncate_filename(fd['filename'], 20) for fd in file_data]

        # ----- Build summary HTML (escaped) -----
        summary_html = '<details style="margin-bottom:4px; border-radius:6px; padding:6px;"><summary style="cursor:pointer; font-weight:bold; font-size:0.9rem;">📋 对比摘要 (点击展开)</summary><div style="margin-top:12px; border-left:8px solid #2c3e50; padding-left:8px;">'
        for fd in file_data:
            preview = html.escape(fd['text'][:200].replace('\n', ' ')) + '…'
            safe_filename = html.escape(fd['filename'])
            summary_html += f'<div style="margin-bottom:15px;"><strong>📄 {safe_filename}</strong><br><span style="color:#666; font-size:0.85rem;">{preview}</span></div>'
        if strong_alert_files:
            safe_strong = ', '.join(html.escape(f) for f in strong_alert_files)
            summary_html += f'<p style="color:#d9534f; font-weight:bold;">🚨 强烈警告：以下文件风险度超过20：{safe_strong}</p>'
        elif high_risk_files:
            safe_high = ', '.join(html.escape(f) for f in high_risk_files)
            summary_html += f'<p style="color:#f0ad4e; font-weight:bold;">⚠️ 可疑文件：以下文件风险度超过10：{safe_high}</p>'
        else:
            summary_html += '<p style="color:#5cb85c;">✅ 未发现高风险文件（风险度均≤10）</p>'
        summary_html += '</div></details>'

        if n == 2:
            p = pairs[0]
            if p['blocks']:
                detail_rows = ""
                for b in p['blocks']:
                    detail_rows += f'''
                    <tr>
                        <td style="border:1px solid #ccc; padding:8px; text-align:center;">{b["id"]}</td>
                        <td style="border:1px solid #ccc; padding:8px; text-align:center;">{b["size"]}</td>
                        <td style="border:1px solid #ccc; padding:8px; word-break:break-word; max-width:300px;">{html.escape(b["text1_snippet"])}</td>
                        <td style="border:1px solid #ccc; padding:8px; word-break:break-word; max-width:300px;">{html.escape(b["text2_snippet"])}</td>
                    </tr>
                    '''
                detailed_report = f'<details><summary style="cursor:pointer; font-weight:bold;">📋 详细相似度明细报告（共 {len(p["blocks"])} 个匹配块）</summary><div style="margin-top:12px;"><p><strong>总匹配字符数：</strong>{sum(b["size"] for b in p["blocks"])} 字符 &nbsp;|&nbsp;<strong>平均匹配块长度：</strong>{round(sum(b["size"] for b in p["blocks"]) / len(p["blocks"]), 1)} 字符</p><div style="overflow-x:auto;"><table style="width:100%; border-collapse:collapse; margin-top:10px;"><thead><tr style="background:#f0f0f0;"><th style="border:1px solid #ccc; padding:8px;">块序号</th><th style="border:1px solid #ccc; padding:8px;">匹配字符数</th><th style="border:1px solid #ccc; padding:8px;">文档A片段</th><th style="border:1px solid #ccc; padding:8px;">文档B片段</th></tr></thead><tbody>{detail_rows}</tbody></table></div></div></details>'
                main_report = detailed_report
            else:
                main_report = "<p>未检测到显著匹配块。</p>"
        else:
            matrix_html = '<details><summary style="cursor:pointer; font-weight:bold;">📊 风险度矩阵 (点击展开/折叠)</summary><div style="overflow-x:auto; margin-top:12px;"><table style="border-collapse:collapse; font-size:0.85rem; min-width:400px; width:100%;"><thead><tr><th style="padding:8px; border:1px solid #ddd;"></th>' + ''.join(f'<th style="padding:8px; border:1px solid #ddd; word-break:break-word;">{html.escape(short_names[i])}</th>' for i in range(n)) + '</tr></thead><tbody>'
            for i in range(n):
                matrix_html += f'<tr><td style="border:1px solid #ddd; padding:8px; font-weight:bold;">{html.escape(short_names[i])}</td>'
                for j in range(n):
                    if i == j:
                        val = '--'
                        bg = ''
                    else:
                        val = f'{risk_matrix[i][j]:.2f}'
                        if risk_matrix[i][j] > 20:
                            bg = ' style="background:#d9534f; color:white; font-weight:bold;"'
                        elif risk_matrix[i][j] > 10:
                            bg = ' style="background:#f0ad4e;"'
                        else:
                            bg = ''
                    matrix_html += f'<td style="border:1px solid #ddd; padding:8px; text-align:center;"{bg}>{html.escape(val)}</td>'
                matrix_html += '</tr>'
            matrix_html += '</tbody></table></div><p style="font-size:0.7rem; color:#666; margin-top:8px;">风险度矩阵（值越高风险越大）</p></details>'
            main_report = matrix_html

        download_token = secrets.token_urlsafe(32)
        download_tokens[download_token] = 20
        session[f'download_path_{download_token}'] = temp_path
        download_link = url_for('export_batch_excel_download', token=download_token, _external=True)
        export_html = f'<p><a href="{download_link}" target="_blank" style="background:#27ae60; color:white; text-decoration:none; border-radius:8px; padding:8px 16px; display:inline-block; margin-top:12px;">📊 下载Excel报告 (可下载20次)</a></p>'
        full_message = f"<!--COMPARE_REPORT--><div style='font-family: -apple-system, BlinkMacSystemFont, \"Segoe UI\", Roboto, sans-serif; line-height:1.5; max-width:100%; overflow-x:auto;'><h4>📁 批量对比结果（{len(file_data)}个文件）</h4>{summary_html}{main_report}{export_html}</div>"

        if session.get('consent_value', 0) == 1:
            ensure_user_exists(user_id)
            report_filename = f"批量对比_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
            report_bytes = full_message.encode('utf-8')
            report_hash = hashlib.sha256(report_bytes).hexdigest()
            ext = '.html'
            unique_name = f"{report_hash}_{int(time.time())}{ext}"
            original_dir = os.path.join(USER_FILES_ORIGINAL_ROOT, user_id)
            os.makedirs(original_dir, exist_ok=True)
            report_path = os.path.join(original_dir, unique_name)
            with open(report_path, 'wb') as f:
                f.write(report_bytes)
            with get_db_connection() as conn:
                with conn.cursor() as cur:
                    cur.execute("""
                        INSERT INTO user_files (user_id, thread_id, filename, size_bytes, expires_at, original_stored_path, file_hash, original_expires_at, original_name)
                        VALUES (%s, %s, %s, %s, NULL, %s, %s, NOW() + INTERVAL '3 days', %s)
                        ON CONFLICT (thread_id, filename) DO NOTHING
                    """, (user_id, thread_id, report_filename, len(report_bytes), report_path, report_hash, report_filename))
                    conn.commit()
            record_file_usage(thread_id, report_filename, 'compare_batch_report', "批量对比生成的报告")

        store_message(thread_id, 'assistant', full_message, thinking="")
        session['chat_history'].append({
            "role": "assistant",
            "content": full_message,
            "thinking": ""
        })
        return jsonify({"success": True, "pair_count": len(pairs)})
    finally:
        release_task_lock(user_id)

@app.route('/export_batch_excel_download/<token>', methods=['GET'])
def export_batch_excel_download(token):
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    remaining = download_tokens.get(token, 0)
    if remaining <= 0:
        return jsonify({"error": "Download link has expired or already used the maximum number of times."}), 410
    temp_path = session.get(f'download_path_{token}')
    if not temp_path or not os.path.exists(temp_path):
        return jsonify({"error": "Comparison data not found."}), 404
    try:
        batch_data = load_batch_comparison_temp(temp_path)
    except Exception as e:
        logger.error(f"Failed to load batch data: {e}")
        return jsonify({"error": "Comparison data corrupted."}), 400

    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment
    from openpyxl.utils import get_column_letter

    file_data = batch_data['file_data']
    pairs = batch_data['pairs']
    timestamp = batch_data['timestamp']
    check_items = batch_data.get('check_items', {})
    key_info_matches = batch_data.get('key_info_matches', [])
    attr_details = batch_data.get('attr_details', [])

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "规律性分析结果"
    ws1.merge_cells('A1:H1')
    ws1['A1'] = "技术标规律性分析检查结果"
    ws1['A1'].font = Font(bold=True, size=14)
    ws1['A2'] = "标段名称：用户自定义"
    ws1['A3'] = f"投标单位个数：{len(file_data)}"
    ws1['A4'] = f"创建时间：{timestamp}"
    max_risk = max(p['risk'] for p in pairs) if pairs else 0
    max_sim = max(p['sim'] for p in pairs) if pairs else 0
    result_parts = []
    if check_items.get('text_sim', True):
        result_parts.append("文本相似度检查存在异常" if max_sim > 0 else "文本相似度检查无异常")
    if check_items.get('key_info', True):
        result_parts.append("重点信息无异常")
    if check_items.get('file_attr', True):
        result_parts.append("文档属性检查无异常")
    if check_items.get('image_sim', True):
        result_parts.append("图片相似度检查无异常")
    result_str = "；".join(result_parts) if result_parts else "无异常"
    ws1['A5'] = f"检查结果：{result_str}"
    ws1['A6'] = "检查规则：检查相似度≥80%的段落，文本中重点信息，相似图片，相同作者；忽略与招标文件相同内容，忽略标点符号及小于6个字的内容，忽略目录，忽略文件中的技术标准，忽略【公司/组织、地名/地址、项目、人员、奖项、身份证号码、电话号码、统一社会信用代码、证书编号】"
    ws1['A7'] = "相似度计算说明：风险度=0.3×重点信息雷同风险+0.3×文件属性雷同风险+0.2×文本相似度×100+0.2×图片相似度×100\n*若某项不参与检查，则其余项按照比例进行折算"
    ws1.merge_cells('A6:H6')
    ws1.merge_cells('A7:H7')
    row = 10
    ws1[f'A{row}'] = "一、标书围串风险分析结果"
    ws1[f'A{row}'].font = Font(bold=True)
    row += 1
    headers = ["投标单位"] + [fd['filename'] for fd in file_data]
    for col, h in enumerate(headers, 1):
        cell = ws1.cell(row=row, column=col, value=h)
        cell.font = Font(bold=True)
        cell.alignment = Alignment(horizontal='center')
    row += 1
    for i in range(len(file_data)):
        ws1.cell(row=row, column=1, value=file_data[i]['filename'])
        for j in range(len(file_data)):
            if i == j:
                val = "--"
            else:
                for p in pairs:
                    if (p['i'] == i and p['j'] == j) or (p['i'] == j and p['j'] == i):
                        val = p['risk']
                        break
                else:
                    val = 0
            ws1.cell(row=row, column=j+2, value=val)
        row += 1
    row += 2
    ws1[f'A{row}'] = "二、分析结果详情"
    ws1[f'A{row}'].font = Font(bold=True)
    row += 1
    detail_headers = ["序号", "投标单位1", "投标单位2", "风险度", "文本相似度（%）", "语义相似度（%）", "图片相似度（%）", "文件属性雷同", "重点信息雷同（项）"]
    for col, h in enumerate(detail_headers, 1):
        cell = ws1.cell(row=row, column=col, value=h)
        cell.font = Font(bold=True)
        cell.alignment = Alignment(horizontal='center')
    row += 1
    for idx, p in enumerate(pairs, 1):
        ws1.cell(row=row, column=1, value=idx)
        ws1.cell(row=row, column=2, value=p['name1'])
        ws1.cell(row=row, column=3, value=p['name2'])
        ws1.cell(row=row, column=4, value=p['risk'])
        ws1.cell(row=row, column=5, value=p['sim'])
        ws1.cell(row=row, column=6, value=p.get('semantic_sim', 0))
        ws1.cell(row=row, column=7, value=0)
        ws1.cell(row=row, column=8, value="是" if p['attr_same'] else "否")
        ki_match = next((k for k in key_info_matches if k['name1'] == p['name1'] and k['name2'] == p['name2']), None)
        ki_count = len(ki_match['common_keywords']) if ki_match else 0
        ws1.cell(row=row, column=9, value=ki_count)
        row += 1
    for col in range(1, 9):
        ws1.column_dimensions[get_column_letter(col)].width = 20

    ws2 = wb.create_sheet("规律性分析详情（文本）")
    ws2['A1'] = "技术标规律性分析详情（文本）"
    ws2.merge_cells('A1:I1')
    ws2['A2'] = "标段名称：用户自定义"
    ws2['A3'] = "检查规则：检查相似度≥80%的段落，忽略与招标文件相同内容，忽略标点符号及小于6个字的内容，忽略目录，忽略文件中的技术标准"
    ws2['A5'] = "序号"
    ws2['B5'] = "目标单位"
    ws2['C5'] = "目标单位对应文档"
    ws2['D5'] = "页码"
    ws2['E5'] = "目标单位内容"
    ws2['F5'] = "对比单位"
    ws2['G5'] = "对比单位对应文档"
    ws2['H5'] = "页码"
    ws2['I5'] = "对比单位相似内容"
    for col in range(1, 10):
        cell = ws2.cell(row=5, column=col)
        cell.font = Font(bold=True)
        cell.alignment = Alignment(horizontal='center')
    r = 6
    seq = 1
    for p in pairs:
        if not p['blocks']:
            continue
        for block in p['blocks']:
            ws2.cell(row=r, column=1, value=seq)
            ws2.cell(row=r, column=2, value=p['name1'])
            ws2.cell(row=r, column=3, value=p['name1'] + ".pdf")
            ws2.cell(row=r, column=4, value=block.get('page1', ''))
            ws2.cell(row=r, column=5, value=block['text1_snippet'])
            ws2.cell(row=r, column=6, value=p['name2'])
            ws2.cell(row=r, column=7, value=p['name2'] + ".pdf")
            ws2.cell(row=r, column=8, value=block.get('page2', ''))
            ws2.cell(row=r, column=9, value=block['text2_snippet'])
            r += 1
            seq += 1
    for col in range(1, 10):
        ws2.column_dimensions[get_column_letter(col)].width = 30

    ws3 = wb.create_sheet("规律性分析详情（重点信息）")
    ws3['A1'] = "技术标规律性分析详情（重点信息）"
    ws3.merge_cells('A1:I1')
    ws3['A2'] = "标段名称：用户自定义"
    ws3['A3'] = "检查规则：检查文本中重点信息；忽略【公司/组织、地名/地址、项目、人员、奖项、身份证号码、电话号码、统一社会信用代码、证书编号】"
    ws3['A5'] = "序号"
    ws3['B5'] = "AI识别类型"
    ws3['C5'] = "内容"
    ws3['D5'] = "目标单位"
    ws3['E5'] = "目标单位对应文档"
    ws3['F5'] = "页码"
    ws3['G5'] = "对比单位"
    ws3['H5'] = "对比单位对应文档"
    ws3['I5'] = "页码"
    for col in range(1, 10):
        cell = ws3.cell(row=5, column=col)
        cell.font = Font(bold=True)
        cell.alignment = Alignment(horizontal='center')
    r = 6
    seq = 1
    for ki in key_info_matches:
        for kw in ki['common_keywords']:
            ws3.cell(row=r, column=1, value=seq)
            ws3.cell(row=r, column=2, value="关键词")
            ws3.cell(row=r, column=3, value=kw)
            ws3.cell(row=r, column=4, value=ki['name1'])
            ws3.cell(row=r, column=5, value=ki['name1'] + ".pdf")
            ws3.cell(row=r, column=7, value=ki['name2'])
            ws3.cell(row=r, column=8, value=ki['name2'] + ".pdf")
            r += 1
            seq += 1
    for col in range(1, 10):
        ws3.column_dimensions[get_column_letter(col)].width = 20

    ws4 = wb.create_sheet("技术标规律性分析详情（文件属性-汇总）")
    ws4['A1'] = "技术标规律性分析详情（文件属性）"
    ws4.merge_cells('A1:E1')
    ws4['A2'] = "标段名称：用户自定义"
    ws4['A3'] = "检查规则：相同作者"
    ws4['A5'] = "序号"
    ws4['B5'] = "单位名称"
    ws4['C5'] = "作者"
    ws4['D5'] = "属性相同单位数量"
    ws4['E5'] = "属性相同单位名称"
    for col in range(1, 6):
        cell = ws4.cell(row=5, column=col)
        cell.font = Font(bold=True)
        cell.alignment = Alignment(horizontal='center')
    author_map = {}
    for fd in file_data:
        name = fd['filename']
        author = fd['metadata'].get('author', '')
        author_map.setdefault(author, []).append(name)
    r = 6
    seq = 1
    for fd in file_data:
        name = fd['filename']
        author = fd['metadata'].get('author', '')
        same_authors = author_map.get(author, [])
        same_count = len(same_authors) - 1
        same_names = ", ".join([n for n in same_authors if n != name]) if same_count > 0 else ""
        ws4.cell(row=r, column=1, value=seq)
        ws4.cell(row=r, column=2, value=name)
        ws4.cell(row=r, column=3, value=author)
        ws4.cell(row=r, column=4, value=same_count)
        ws4.cell(row=r, column=5, value=same_names)
        r += 1
        seq += 1
    for col in range(1, 6):
        ws4.column_dimensions[get_column_letter(col)].width = 30

    ws5 = wb.create_sheet("技术标规律性分析详情（文件属性-详情）")
    ws5['A1'] = "技术标规律性分析详情（文件属性日志信息）"
    ws5.merge_cells('A1:G1')
    ws5['A2'] = "标段名称：用户自定义"
    ws5['A3'] = "检查规则：相同作者"
    ws5['A5'] = "序号"
    ws5['B5'] = "单位名称"
    ws5['C5'] = "文档名称"
    ws5['D5'] = "作者"
    ws5['E5'] = "属性相同单位数量"
    ws5['F5'] = "属性相同单位名称"
    ws5['G5'] = "属性相同文档名称"
    for col in range(1, 8):
        cell = ws5.cell(row=5, column=col)
        cell.font = Font(bold=True)
        cell.alignment = Alignment(horizontal='center')
    r = 6
    seq = 1
    for fd in file_data:
        name = fd['filename']
        author = fd['metadata'].get('author', '')
        same_authors = author_map.get(author, [])
        same_count = len(same_authors) - 1
        same_names = ", ".join([n for n in same_authors if n != name]) if same_count > 0 else "无重复单位"
        same_docs = ", ".join([n for n in same_authors if n != name]) if same_count > 0 else "无重复文档"
        ws5.cell(row=r, column=1, value=seq)
        ws5.cell(row=r, column=2, value=name)
        ws5.cell(row=r, column=3, value=name + ".pdf")
        ws5.cell(row=r, column=4, value=author)
        ws5.cell(row=r, column=5, value=same_count)
        ws5.cell(row=r, column=6, value=same_names)
        ws5.cell(row=r, column=7, value=same_docs)
        r += 1
        seq += 1
    for col in range(1, 8):
        ws5.column_dimensions[get_column_letter(col)].width = 30

    output = BytesIO()
    wb.save(output)
    output.seek(0)
    filename = f"清标分析结果_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
    with download_tokens_lock:
        download_tokens[token] -= 1
        if download_tokens[token] <= 0:
            del download_tokens[token]
            session.pop(f'download_path_{token}', None)
            if temp_path and os.path.exists(temp_path):
                os.unlink(temp_path)
    return send_file(output, as_attachment=True, download_name=filename, mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')

@app.route('/export_batch_excel', methods=['POST'])
def export_batch_excel():
    return jsonify({"error": "Please use the download link from the batch comparison report."}), 400

# ---------- Recycle Bin routes ----------
@app.route('/get_recycle_bin', methods=['GET'])
def get_recycle_bin():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "Not logged in"}), 401

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("""
                DO $$ 
                BEGIN
                    BEGIN
                        ALTER TABLE recycle_bin ADD COLUMN deletion_reason TEXT DEFAULT 'manual';
                    EXCEPTION WHEN duplicate_column THEN NULL;
                    END;
                END $$;
            """)
            conn.commit()

    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                        SELECT rb.id,
                               rb.original_table,
                               rb.original_id,
                               rb.file_name,
                               rb.file_size,
                               rb.deleted_at,
                               rb.expires_at,
                               rb.deletion_reason,
                               'chat'              as source,
                               u_uploader.username as uploaded_by_name,
                               u_deleter.username  as deleted_by_name
                        FROM recycle_bin rb
                                 LEFT JOIN users u_uploader ON rb.uploaded_by = u_uploader.user_id
                                 LEFT JOIN users u_deleter ON rb.deleted_by = u_deleter.user_id
                        WHERE rb.user_id = %s
                          AND rb.expires_at > NOW()
                        ORDER BY rb.deleted_at DESC
                        """, (user_id,))
            chat_items = cur.fetchall()

            cur.execute("""
                        SELECT prb.id,
                               prb.original_table,
                               prb.original_id,
                               prb.file_name,
                               prb.file_size,
                               prb.deleted_at,
                               prb.expires_at,
                               p.name              as project_name,
                               'project'           as source,
                               u_uploader.username as uploaded_by_name,
                               u_deleter.username  as deleted_by_name
                        FROM project_recycle_bin prb
                                 JOIN projects p ON prb.project_id = p.id
                                 LEFT JOIN users u_uploader ON prb.uploaded_by = u_uploader.user_id
                                 LEFT JOIN users u_deleter ON prb.deleted_by = u_deleter.user_id
                        WHERE prb.expires_at > NOW()
                        ORDER BY prb.deleted_at DESC
                        """)
            project_items = cur.fetchall()

            cur.execute("""
                SELECT pfrb.id, pfrb.original_id, pfrb.name, pfrb.original_parent_id, pfrb.deleted_at, pfrb.expires_at, 
                       p.name as project_name, 'folder' as source
                FROM project_folders_recycle_bin pfrb
                JOIN projects p ON pfrb.project_id = p.id
                WHERE pfrb.expires_at > NOW()
                ORDER BY pfrb.deleted_at DESC
            """)
            folder_items = cur.fetchall()

            return jsonify({
                "chat_items": chat_items,
                "project_items": project_items,
                "folder_items": folder_items
            })

@app.route('/restore_from_recycle_bin', methods=['POST'])
def restore_from_recycle_bin():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "Not logged in"}), 401

    data = request.get_json()
    item_id = data.get('item_id')
    source = data.get('source')
    section = data.get('section')
    restore_all = data.get('restore_all', False)

    with get_db_connection() as conn:
        with db_transaction(conn):
            with conn.cursor(cursor_factory=RealDictCursor) as cur:
                if restore_all:
                    restored_count = 0
                    if section == 'chat':
                        cur.execute("SELECT * FROM recycle_bin WHERE user_id = %s AND expires_at > NOW()", (user_id,))
                        items = cur.fetchall()
                        for item in items:
                            meta_data = {}
                            if item.get('deletion_reason') == 'chat_deleted':
                                meta_data['restored_from'] = 'chat_deletion'
                                meta_data['original_thread_id'] = item.get('original_thread_id')
                            meta_data_json = json.dumps(meta_data)
                            cur.execute("""
                                INSERT INTO user_files (user_id, thread_id, filename, content, size_bytes, expires_at,
                                                        original_stored_path, file_hash, meta_data, original_name)
                                VALUES (%s, %s, %s, %s, %s, NOW() + INTERVAL '3 days', %s, %s, %s::jsonb, %s)
                            """, (user_id, None, item['file_name'], item['file_content'], item['file_size'],
                                  item['original_stored_path'], item['file_hash'], meta_data_json, item['file_name']))
                            cur.execute("DELETE FROM recycle_bin WHERE id = %s", (item['id'],))
                            restored_count += 1
                    elif section == 'project_files':
                        cur.execute("SELECT * FROM project_recycle_bin WHERE expires_at > NOW()")
                        items = cur.fetchall()
                        for item in items:
                            folder_id = item['folder_id']
                            if folder_id:
                                cur.execute("SELECT id FROM project_folders WHERE id = %s AND project_id = %s",
                                            (folder_id, item['project_id']))
                                if not cur.fetchone():
                                    restore_folder_path_for_file(item, conn, cur)
                            cur.execute("""
                                INSERT INTO project_files (project_id, folder_id, filename, original_name, file_size,
                                                           stored_path, version, uploaded_by, file_hash)
                                VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s)
                            """, (item['project_id'], item['folder_id'], item['file_name'], item['original_name'],
                                  item['file_size'], item['stored_path'], item['version'],
                                  item['uploaded_by'], item['file_hash']))
                            cur.execute("DELETE FROM project_recycle_bin WHERE id = %s", (item['id'],))
                            restored_count += 1
                    elif section == 'project_folders':
                        cur.execute("SELECT * FROM project_folders_recycle_bin WHERE expires_at > NOW()")
                        folders = cur.fetchall()
                        for folder in folders:
                            restore_folder_recursive(folder, conn, cur)
                            restored_count += 1
                    else:
                        return jsonify({"error": "Invalid section"}), 400
                    return jsonify({"success": True, "restored_count": restored_count})

                if source == 'chat':
                    cur.execute("SELECT * FROM recycle_bin WHERE id = %s AND user_id = %s AND expires_at > NOW()",
                                (item_id, user_id))
                    item = cur.fetchone()
                    if not item:
                        return jsonify({"error": "Item not found or expired"}), 404
                    meta_data = {}
                    if item.get('deletion_reason') == 'chat_deleted':
                        meta_data['restored_from'] = 'chat_deletion'
                        meta_data['original_thread_id'] = item.get('original_thread_id')
                    meta_data_json = json.dumps(meta_data)
                    cur.execute("""
                        INSERT INTO user_files (user_id, thread_id, filename, content, size_bytes, expires_at,
                                                original_stored_path, file_hash, meta_data, original_name)
                        VALUES (%s, %s, %s, %s, %s, NOW() + INTERVAL '3 days', %s, %s, %s::jsonb, %s)
                    """, (user_id, None, item['file_name'], item['file_content'], item['file_size'],
                          item['original_stored_path'], item['file_hash'], meta_data_json, item['file_name']))
                    cur.execute("DELETE FROM recycle_bin WHERE id = %s", (item_id,))
                elif source == 'project':
                    cur.execute("SELECT * FROM project_recycle_bin WHERE id = %s", (item_id,))
                    item = cur.fetchone()
                    if not item:
                        return jsonify({"error": "Item not found"}), 404
                    folder_id = item['folder_id']
                    if folder_id:
                        cur.execute("SELECT id FROM project_folders WHERE id = %s AND project_id = %s",
                                    (folder_id, item['project_id']))
                        if not cur.fetchone():
                            restore_folder_path_for_file(item, conn, cur)
                    cur.execute("""
                        INSERT INTO project_files (project_id, folder_id, filename, original_name, file_size,
                                                   stored_path, version, uploaded_by, file_hash)
                        VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s)
                    """, (item['project_id'], item['folder_id'], item['file_name'], item['original_name'],
                          item['file_size'], item['stored_path'], item['version'],
                          item['uploaded_by'], item['file_hash']))
                    cur.execute("DELETE FROM project_recycle_bin WHERE id = %s", (item_id,))
                elif source == 'folder':
                    cur.execute("SELECT * FROM project_folders_recycle_bin WHERE id = %s", (item_id,))
                    folder = cur.fetchone()
                    if not folder:
                        return jsonify({"error": "Folder not found"}), 404
                    restore_folder_recursive(folder, conn, cur)
                else:
                    return jsonify({"error": "Invalid source"}), 400
                conn.commit()
                return jsonify({"success": True})

@app.route('/delete_recycle_item', methods=['POST'])
def delete_recycle_item():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "Not logged in"}), 401

    data = request.get_json()
    item_id = data.get('item_id')
    source = data.get('source')

    if not item_id or not source:
        return jsonify({"error": "Missing item_id or source"}), 400

    with get_db_connection() as conn:
        with db_transaction(conn):
            with conn.cursor() as cur:
                if source == 'chat':
                    cur.execute("SELECT original_stored_path FROM recycle_bin WHERE id = %s AND user_id = %s", (item_id, user_id))
                    row = cur.fetchone()
                    if row and row[0] and os.path.exists(row[0]):
                        try:
                            os.remove(row[0])
                        except Exception as e:
                            logger.warning(f"Failed to delete physical file {row[0]}: {e}")
                    cur.execute("DELETE FROM recycle_bin WHERE id = %s AND user_id = %s", (item_id, user_id))
                elif source == 'project':
                    cur.execute("SELECT stored_path FROM project_recycle_bin WHERE id = %s", (item_id,))
                    row = cur.fetchone()
                    if row and row[0] and os.path.exists(row[0]):
                        try:
                            os.remove(row[0])
                        except Exception as e:
                            logger.warning(f"Failed to delete project file {row[0]}: {e}")
                    cur.execute("DELETE FROM project_recycle_bin WHERE id = %s", (item_id,))
                elif source == 'folder':
                    cur.execute("SELECT project_id, original_id FROM project_folders_recycle_bin WHERE id = %s", (item_id,))
                    folder = cur.fetchone()
                    if folder:
                        project_id = folder[0]
                        original_folder_id = folder[1]
                        cur.execute("SELECT stored_path FROM project_recycle_bin WHERE project_id = %s AND folder_id = %s", (project_id, original_folder_id))
                        for (stored_path,) in cur.fetchall():
                            if stored_path and os.path.exists(stored_path):
                                try:
                                    os.remove(stored_path)
                                except Exception as e:
                                    logger.warning(f"Failed to delete file {stored_path}: {e}")
                        cur.execute("DELETE FROM project_recycle_bin WHERE project_id = %s AND folder_id = %s", (project_id, original_folder_id))
                        cur.execute("""
                            WITH RECURSIVE folder_tree AS (
                                SELECT id, original_id, project_id, original_parent_id
                                FROM project_folders_recycle_bin
                                WHERE id = %s
                                UNION ALL
                                SELECT pf.id, pf.original_id, pf.project_id, pf.original_parent_id
                                FROM project_folders_recycle_bin pf
                                INNER JOIN folder_tree ft ON pf.original_parent_id = ft.original_id AND pf.project_id = ft.project_id
                            )
                            SELECT id, original_id FROM folder_tree
                        """, (item_id,))
                        subfolders = cur.fetchall()
                        for (sf_id, sf_orig_id) in subfolders:
                            cur.execute("SELECT stored_path FROM project_recycle_bin WHERE project_id = %s AND folder_id = %s", (project_id, sf_orig_id))
                            for (sp,) in cur.fetchall():
                                if sp and os.path.exists(sp):
                                    try:
                                        os.remove(sp)
                                    except Exception:
                                        pass
                            cur.execute("DELETE FROM project_recycle_bin WHERE project_id = %s AND folder_id = %s", (project_id, sf_orig_id))
                            cur.execute("DELETE FROM project_folders_recycle_bin WHERE id = %s", (sf_id,))
                    else:
                        cur.execute("DELETE FROM project_folders_recycle_bin WHERE id = %s", (item_id,))
                else:
                    return jsonify({"error": "Invalid source"}), 400
                conn.commit()
                return jsonify({"success": True})

@app.route('/empty_recycle_bin', methods=['POST'])
def empty_recycle_bin():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "Not logged in"}), 401

    data = request.get_json()
    source = data.get('source')

    with get_db_connection() as conn:
        with db_transaction(conn):
            with conn.cursor() as cur:
                if source == 'chat' or source == 'all':
                    cur.execute("SELECT original_stored_path FROM recycle_bin WHERE user_id = %s", (user_id,))
                    paths = cur.fetchall()
                    for row in paths:
                        if row[0] and os.path.exists(row[0]):
                            try:
                                os.remove(row[0])
                            except Exception:
                                pass
                    cur.execute("DELETE FROM recycle_bin WHERE user_id = %s", (user_id,))
                if source == 'project_files' or source == 'all':
                    cur.execute("SELECT stored_path FROM project_recycle_bin")
                    paths = cur.fetchall()
                    for row in paths:
                        if row[0] and os.path.exists(row[0]):
                            try:
                                os.remove(row[0])
                            except Exception:
                                pass
                    cur.execute("DELETE FROM project_recycle_bin")
                if source == 'project_folders' or source == 'all':
                    cur.execute("DELETE FROM project_folders_recycle_bin")
                conn.commit()
                return jsonify({"success": True})

@app.route('/set_image_analysis', methods=['POST'])
def set_image_analysis():
    data = request.get_json()
    enabled = data.get('enabled', True)
    session['analyze_images'] = enabled
    return jsonify({"success": True})

@app.route('/search_chat', methods=['GET'])
def search_chat():
    if session.get('consent_value', 0) != 1:
        return jsonify({"error": "Consent not given"}), 403
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "Not logged in"}), 401

    q = request.args.get('q', '').strip()
    if len(q) < 2:
        return jsonify({"error": "Search query must be at least 2 characters"}), 400

    start_date = request.args.get('start_date')
    end_date = request.args.get('end_date')
    fuzzy = request.args.get('fuzzy', 'false').lower() == 'true'
    role = request.args.get('role', 'assistant')

    if fuzzy:
        search_pattern = f"%{q}%"
    else:
        search_pattern = q

    date_condition = ""
    params = [user_id, search_pattern]
    if start_date:
        date_condition += " AND cm.timestamp >= %s"
        params.append(start_date)
    if end_date:
        date_condition += " AND cm.timestamp <= %s"
        params.append(end_date)

    if role == 'assistant':
        role_condition = " AND cm.role = 'assistant'"
    elif role == 'user':
        role_condition = " AND cm.role = 'user'"
    else:
        role_condition = ""

    query = f"""
        SELECT cs.thread_id, cs.title, cm.role, cm.content, cm.timestamp, cm.id as message_id,
               SUBSTRING(cm.content, 1, 200) as snippet
        FROM chat_messages cm
        JOIN chat_sessions cs ON cm.thread_id = cs.thread_id
        WHERE cs.user_id = %s
          AND cm.content ILIKE %s
          {role_condition}
          {date_condition}
        ORDER BY cm.timestamp DESC
        LIMIT 100
    """

    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute(query, params)
            results = cur.fetchall()
            formatted_results = []
            for row in results:
                ts_utc = row['timestamp']
                ts_beijing = ts_utc.astimezone(BEIJING_TZ).strftime('%Y-%m-%d %H:%M:%S') if ts_utc else None
                snippet = row['snippet'] or ""
                if fuzzy:
                    import re
                    escaped = re.escape(q)
                    highlighted = re.sub(f"({escaped})", r'<mark>\1</mark>', snippet, flags=re.IGNORECASE)
                else:
                    highlighted = snippet
                formatted_results.append({
                    'thread_id': row['thread_id'],
                    'title': row['title'],
                    'role': row['role'],
                    'snippet': snippet,
                    'timestamp_str': ts_beijing,
                    'highlighted_snippet': highlighted,
                    'message_id': row['message_id']
                })
            return jsonify({"results": formatted_results})

# ---------- Admin database browser ----------
@app.route('/admin/db_tables', methods=['GET'])
@admin_required
def admin_db_tables():
    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                        SELECT tablename
                        FROM pg_tables
                        WHERE schemaname = 'public'
                        ORDER BY tablename
                        """)
            tables = [row['tablename'] for row in cur.fetchall()]
            return jsonify({"tables": tables})

@app.route('/admin/db_table_data', methods=['POST'])
@admin_required
def admin_db_table_data():
    data = request.get_json()
    table = data.get('table')
    page = int(data.get('page', 1))
    per_page = int(data.get('per_page', 50))
    search = data.get('search', '').strip()
    search_column = data.get('search_column', '')

    if not table:
        logger.warning("admin_db_table_data: no table name provided")
        return jsonify({"error": "Table name required"}), 400

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT tablename FROM pg_tables WHERE schemaname = 'public' AND tablename = %s", (table,))
            if not cur.fetchone():
                logger.warning(f"admin_db_table_data: table '{table}' not found in pg_tables")
                return jsonify({"error": "Invalid table name"}), 400

    offset = (page - 1) * per_page
    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                        SELECT column_name, data_type
                        FROM information_schema.columns
                        WHERE table_name = %s
                        ORDER BY ordinal_position
                        """, (table,))
            columns = cur.fetchall()
            col_names = [col['column_name'] for col in columns]

            if 'id' in col_names:
                order_col = 'id'
            else:
                order_col = col_names[0] if col_names else '1'

            where_clause = sql.SQL("")
            params = []
            if search and search_column and search_column in col_names:
                where_clause = sql.SQL(" WHERE {col}::text ILIKE %s").format(col=sql.Identifier(search_column))
                params.append(f"%{search}%")
            elif search:
                text_cols = [col['column_name'] for col in columns if
                             col['data_type'] in ('text', 'varchar', 'character varying', 'char', 'name')]
                if text_cols:
                    conditions = sql.SQL(" OR ").join([
                        sql.SQL("{c}::text ILIKE %s").format(c=sql.Identifier(c)) for c in text_cols
                    ])
                    where_clause = sql.SQL(" WHERE {conditions}").format(conditions=conditions)
                    params.extend([f"%{search}%"] * len(text_cols))

            table_ident = sql.Identifier(table)
            order_ident = sql.Identifier(order_col)
            count_query = sql.SQL("SELECT COUNT(*) as total FROM {table} {where}").format(
                table=table_ident, where=where_clause
            ).as_string(conn)
            cur.execute(count_query, params)
            total = cur.fetchone()['total']

            query = sql.SQL("SELECT * FROM {table} {where} ORDER BY {order_col} DESC LIMIT %s OFFSET %s").format(
                table=table_ident, where=where_clause, order_col=order_ident
            ).as_string(conn)
            cur.execute(query, params + [per_page, offset])
            rows = cur.fetchall()

            return jsonify({
                "columns": col_names,
                "rows": rows,
                "total": total,
                "page": page,
                "per_page": per_page
            })

@app.route('/admin/db_update_row', methods=['POST'])
@admin_required
@admin_rate_limiter
def admin_db_update_row():
    data = request.get_json()
    table = data.get('table')
    row_id = data.get('row_id')
    column = data.get('column')
    new_value = data.get('value')
    pin = data.get('pin', '').strip()
    admin_user_id = session.get('user_id')
    admin_username = session.get('username', 'admin')

    if not ADMIN_PASSWORD_HASH:
        logger.error("Admin password hash not configured")
        log_admin_action(admin_user_id, admin_username, 'UPDATE', table, row_id,
                         column, None, new_value, success=False,
                         error_message="Admin password hash not configured")
        return jsonify({"error": "Admin authentication not configured"}), 500

    if not check_password_hash(ADMIN_PASSWORD_HASH, pin):
        log_admin_action(admin_user_id, admin_username, 'UPDATE', table, row_id,
                         column, None, new_value, success=False,
                         error_message="Invalid admin PIN")
        return jsonify({"error": "Invalid admin PIN"}), 403

    # Validate table exists
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT tablename FROM pg_tables WHERE schemaname = 'public' AND tablename = %s", (table,))
            if not cur.fetchone():
                log_admin_action(admin_user_id, admin_username, 'UPDATE', table, row_id,
                                 column, None, new_value, success=False,
                                 error_message=f"Invalid table name: {table}")
                return jsonify({"error": "Invalid table name"}), 400

    # Validate column exists
    if not validate_table_column(table, column):
        log_admin_action(admin_user_id, admin_username, 'UPDATE', table, row_id,
                         column, None, new_value, success=False,
                         error_message=f"Invalid column: {column}")
        return jsonify({"error": f"Column '{column}' does not exist"}), 400

    # Determine primary key column
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT column_name 
                FROM information_schema.columns 
                WHERE table_name = %s AND column_name IN ('id', 'thread_id')
            """, (table,))
            pk_col = cur.fetchone()
            if not pk_col:
                log_admin_action(admin_user_id, admin_username, 'UPDATE', table, row_id,
                                 column, None, new_value, success=False,
                                 error_message="No primary key column found")
                return jsonify({"error": f"Table '{table}' has no known primary key"}), 400
            pk_col = pk_col[0]

    if pk_col == 'id':
        try:
            row_id_val = int(row_id)
        except ValueError:
            log_admin_action(admin_user_id, admin_username, 'UPDATE', table, row_id,
                             column, None, new_value, success=False,
                             error_message="Invalid row_id (not an integer)")
            return jsonify({"error": "Row ID must be an integer"}), 400
    else:
        row_id_val = row_id

    old_value = None
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            select_q = sql.SQL("SELECT {col} FROM {table} WHERE {pk} = %s").format(
                col=sql.Identifier(column), table=sql.Identifier(table), pk=sql.Identifier(pk_col)
            ).as_string(conn)
            cur.execute(select_q, (row_id_val,))
            row = cur.fetchone()
            if row:
                old_value = str(row[0]) if row[0] is not None else None
            else:
                log_admin_action(admin_user_id, admin_username, 'UPDATE', table, str(row_id_val),
                                 column, None, new_value, success=False,
                                 error_message="Row not found")
                return jsonify({"error": "Row not found"}), 404

    # Execute update
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            update_q = sql.SQL("UPDATE {table} SET {col} = %s WHERE {pk} = %s").format(
                table=sql.Identifier(table), col=sql.Identifier(column), pk=sql.Identifier(pk_col)
            ).as_string(conn)
            cur.execute(update_q, (new_value, row_id_val))
            conn.commit()

    # Log successful update
    log_admin_action(admin_user_id, admin_username, 'UPDATE', table, str(row_id_val),
                     column, old_value, new_value, success=True)

    return jsonify({"success": True})

@app.route('/admin/db_delete_row', methods=['POST'])
@admin_required
@admin_rate_limiter
def admin_db_delete_row():
    data = request.get_json()
    table = data.get('table')
    row_id = data.get('row_id')
    pin = data.get('pin', '').strip()
    admin_user_id = session.get('user_id')
    admin_username = session.get('username', 'admin')

    if not table or not row_id or row_id == 'undefined':
        log_admin_action(admin_user_id, admin_username, 'DELETE', table, row_id,
                         success=False, error_message="Missing table or row_id")
        return jsonify({"error": "Missing table or valid row_id"}), 400

    if not ADMIN_PASSWORD_HASH:
        log_admin_action(admin_user_id, admin_username, 'DELETE', table, row_id,
                         success=False, error_message="Admin password hash not configured")
        return jsonify({"error": "Admin authentication not configured"}), 500

    if not check_password_hash(ADMIN_PASSWORD_HASH, pin):
        log_admin_action(admin_user_id, admin_username, 'DELETE', table, row_id,
                         success=False, error_message="Invalid admin PIN")
        return jsonify({"error": "Invalid admin PIN"}), 403

    # Validate table exists
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT tablename FROM pg_tables WHERE schemaname = 'public' AND tablename = %s", (table,))
            if not cur.fetchone():
                log_admin_action(admin_user_id, admin_username, 'DELETE', table, row_id,
                                 success=False, error_message=f"Invalid table: {table}")
                return jsonify({"error": "Invalid table name"}), 400

    # Determine primary key column
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT column_name 
                FROM information_schema.columns 
                WHERE table_name = %s AND column_name IN ('id', 'thread_id')
            """, (table,))
            pk_col = cur.fetchone()
            if not pk_col:
                log_admin_action(admin_user_id, admin_username, 'DELETE', table, row_id,
                                 success=False, error_message="No primary key column found")
                return jsonify({"error": f"Table '{table}' has no known primary key"}), 400
            pk_col = pk_col[0]

    if pk_col == 'id':
        try:
            row_id_val = int(row_id)
        except ValueError:
            log_admin_action(admin_user_id, admin_username, 'DELETE', table, row_id,
                             success=False, error_message="Invalid row_id (not integer)")
            return jsonify({"error": "Row ID must be an integer"}), 400
    else:
        row_id_val = row_id

    row_snapshot = None
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            select_q = sql.SQL("SELECT * FROM {table} WHERE {pk} = %s").format(
                table=sql.Identifier(table), pk=sql.Identifier(pk_col)
            ).as_string(conn)
            cur.execute(select_q, (row_id_val,))
            row = cur.fetchone()
            if row:
                # Convert to dict for logging
                columns = [desc[0] for desc in cur.description]
                row_snapshot = dict(zip(columns, row))
            else:
                log_admin_action(admin_user_id, admin_username, 'DELETE', table, str(row_id_val),
                                 success=False, error_message="Row not found")
                return jsonify({"error": "Row not found"}), 404

    # Execute deletion
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            delete_q = sql.SQL("DELETE FROM {table} WHERE {pk} = %s").format(
                table=sql.Identifier(table), pk=sql.Identifier(pk_col)
            ).as_string(conn)
            cur.execute(delete_q, (row_id_val,))
            conn.commit()

    # Log successful deletion (store row snapshot as old_value JSON)
    import json as json_module
    log_admin_action(admin_user_id, admin_username, 'DELETE', table, str(row_id_val),
                     old_value=json_module.dumps(row_snapshot, default=str) if row_snapshot else None,
                     success=True)

    return jsonify({"success": True})

@app.route('/admin/clear_file_cache', methods=['POST'])
@admin_required
def clear_file_cache():
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("TRUNCATE TABLE file_text_cache")
            conn.commit()
    return jsonify({"success": True})

@app.route('/admin/audit_log', methods=['GET'])
@admin_required
def admin_audit_log():
    page = request.args.get('page', 1, type=int)
    per_page = 50
    offset = (page - 1) * per_page
    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                SELECT * FROM admin_audit_log
                ORDER BY created_at DESC
                LIMIT %s OFFSET %s
            """, (per_page, offset))
            logs = cur.fetchall()
            cur.execute("SELECT COUNT(*) as total FROM admin_audit_log")
            total = cur.fetchone()['total']
    return jsonify({
        "logs": logs,
        "total": total,
        "page": page,
        "per_page": per_page
    })

# ---------- Helper functions for recycle bin folder restoration ----------
def restore_folder_recursive(folder_item, conn, cur, target_parent_id=None):
    parent_id = target_parent_id if target_parent_id is not None else folder_item['original_parent_id']
    cur.execute("""
        INSERT INTO project_folders (id, project_id, parent_folder_id, name, created_at, created_by)
        VALUES (%s, %s, %s, %s, %s, %s)
        ON CONFLICT (id) DO NOTHING
    """, (folder_item['original_id'], folder_item['project_id'], parent_id,
          folder_item['name'], folder_item['created_at'], folder_item['created_by']))
    cur.execute("""
        SELECT * FROM project_recycle_bin
        WHERE project_id = %s AND folder_id = %s
    """, (folder_item['project_id'], folder_item['original_id']))
    files = cur.fetchall()
    for f in files:
        cur.execute("""
            INSERT INTO project_files (project_id, folder_id, filename, original_name, file_size, stored_path, version, uploaded_by, file_hash)
            VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s)
        """, (f['project_id'], folder_item['original_id'], f['file_name'], f['original_name'],
              f['file_size'], f['stored_path'], f['version'], f['uploaded_by'], f['file_hash']))
        cur.execute("DELETE FROM project_recycle_bin WHERE id = %s", (f['id'],))
    cur.execute("""
        SELECT * FROM project_folders_recycle_bin
        WHERE project_id = %s AND original_parent_id = %s
    """, (folder_item['project_id'], folder_item['original_id']))
    subfolders = cur.fetchall()
    for sf in subfolders:
        restore_folder_recursive(sf, conn, cur, target_parent_id=folder_item['original_id'])
    cur.execute("DELETE FROM project_folders_recycle_bin WHERE id = %s", (folder_item['id'],))

def restore_folder_path_for_file(file_item, conn, cur):
    folder_id = file_item['folder_id']
    if folder_id is None:
        return
    cur.execute("SELECT id FROM project_folders WHERE id = %s AND project_id = %s", (folder_id, file_item['project_id']))
    if cur.fetchone():
        return
    cur.execute("SELECT * FROM project_folders_recycle_bin WHERE original_id = %s AND project_id = %s", (folder_id, file_item['project_id']))
    folder = cur.fetchone()
    if not folder:
        return
    if folder['original_parent_id']:
        cur.execute("SELECT * FROM project_folders_recycle_bin WHERE original_id = %s AND project_id = %s", (folder['original_parent_id'], file_item['project_id']))
        parent = cur.fetchone()
        if parent:
            restore_folder_path_for_file(parent, conn, cur)
    cur.execute("""
        INSERT INTO project_folders (id, project_id, parent_folder_id, name, created_at, created_by)
        VALUES (%s, %s, %s, %s, %s, %s)
        ON CONFLICT (id) DO NOTHING
    """, (folder['original_id'], folder['project_id'], folder['original_parent_id'],
          folder['name'], folder['created_at'], folder['created_by']))
    cur.execute("DELETE FROM project_folders_recycle_bin WHERE id = %s", (folder['id'],))

# ---------- Scheduled jobs ----------
def delete_expired_original_files():
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT id, original_stored_path
                FROM user_files
                WHERE original_expires_at IS NOT NULL AND original_expires_at <= NOW()
                  AND original_stored_path IS NOT NULL
            """)
            expired = cur.fetchall()
            for file_id, original_path in expired:
                if original_path and os.path.exists(original_path):
                    try:
                        os.remove(original_path)
                        logger.info(f"Deleted expired original file: {original_path}")
                    except Exception as e:
                        logger.warning(f"Failed to delete expired file {original_path}: {e}")
                cur.execute("UPDATE user_files SET original_stored_path = NULL WHERE id = %s", (file_id,))
            conn.commit()

def cleanup_old_anon_temp_files(days=1):
    now = time.time()
    for item in os.listdir(TEMP_ROOT):
        item_path = os.path.join(TEMP_ROOT, item)
        if os.path.isdir(item_path):
            if (now - os.path.getctime(item_path)) > days * 86400:
                shutil.rmtree(item_path)
                logger.info(f"Removed old anonymous temp dir: {item_path}")

def schedule_project_deletion_cleanup():
    cutoff = utc_now() - timedelta(days=3)
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT id FROM projects WHERE status = 'archived' AND archived_at < %s", (cutoff,))
            to_delete = cur.fetchall()
            for (project_id,) in to_delete:
                logger.info(f"Auto-deleting archived project {project_id} after 3 days")
                cur.execute("DELETE FROM projects WHERE id = %s", (project_id,))
            conn.commit()

def cleanup_expired_recycle_bin():
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT original_stored_path FROM recycle_bin WHERE expires_at <= NOW()")
            paths = cur.fetchall()
            for row in paths:
                if row[0] and os.path.exists(row[0]):
                    try:
                        os.remove(row[0])
                        logger.info(f"Deleted expired recycle file: {row[0]}")
                    except Exception as e:
                        logger.warning(f"Failed to delete expired file {row[0]}: {e}")
            cur.execute("DELETE FROM recycle_bin WHERE expires_at <= NOW()")
            cur.execute("SELECT stored_path FROM project_recycle_bin WHERE expires_at <= NOW()")
            paths = cur.fetchall()
            for row in paths:
                if row[0] and os.path.exists(row[0]):
                    try:
                        os.remove(row[0])
                        logger.info(f"Deleted expired project recycle file: {row[0]}")
                    except Exception as e:
                        logger.warning(f"Failed to delete expired file {row[0]}: {e}")
            cur.execute("DELETE FROM project_recycle_bin WHERE expires_at <= NOW()")
            conn.commit()
            logger.info("Cleaned up expired recycle bin items")

# ======================== Credit Check Routes ========================
credit_tasks = {}  # in‑memory store for running tasks

@app.route('/start_credit_check', methods=['POST'])
def start_credit_check():
    data = request.get_json()
    companies = data.get('companies', [])
    urls = data.get('urls', [])
    if not companies or not urls:
        return jsonify({"error": "Need companies and urls"}), 400

    task_id = str(uuid.uuid4())
    user_id = get_user_id()
    with app.app_context():
        download_url = url_for('download_credit_report', task_id=task_id, _external=True)

    with _credit_tasks_lock:
        credit_tasks[task_id] = {
            'status': 'running',
            'progress': 0,
            'total': len(companies),
            'captcha_needed': False,
            'captcha_image': None,
            'captcha_task': None,
            'captcha_solution': None,
            'reload_captcha': False,
            'download_url': download_url,
            'error': None,
            'waiting': False,
            'resume': False
        }

    threading.Thread(target=_run_credit_check,
                     args=(task_id, companies, urls, user_id, True),
                     daemon=True).start()
    return jsonify({"task_id": task_id})

def _run_credit_check(task_id, companies, urls, user_id, manual_mode=True):
    checker = CreditChecker()
    screenshots = {}

    try:
        for idx, company in enumerate(companies):
            company_shots = []
            for url in urls:
                logger.info(f"Processing {company} at {url}")

                # Navigate and fill company name, also sets zoom
                checker.navigate_and_fill(company, url)

                # Handle CAPTCHA if present
                if checker._is_captcha_present():
                    captcha_img = checker.get_captcha_element_screenshot()
                    if captcha_img:
                        with _credit_tasks_lock:
                            credit_tasks[task_id]['captcha_needed'] = True
                            credit_tasks[task_id]['captcha_image'] = captcha_img.getvalue()
                            credit_tasks[task_id]['captcha_solution'] = None
                            credit_tasks[task_id]['reload_captcha'] = False

                        # Wait for user to solve CAPTCHA via modal
                        captcha_start = time.time()
                        captcha_timeout = 300  # 5 minutes max wait
                        while True:
                            if time.time() - captcha_start > captcha_timeout:
                                logger.warning(f"CAPTCHA wait timed out for task {task_id}")
                                with _credit_tasks_lock:
                                    credit_tasks[task_id]['status'] = 'error'
                                    credit_tasks[task_id]['error'] = '验证码等待超时'
                                return
                            with _credit_tasks_lock:
                                solution = credit_tasks[task_id].get('captcha_solution')
                                reload_flag = credit_tasks[task_id].get('reload_captcha', False)

                            if solution is not None:
                                break
                            if reload_flag:
                                with _credit_tasks_lock:
                                    credit_tasks[task_id]['reload_captcha'] = False
                                # Refresh the CAPTCHA image on the page
                                checker.refresh_captcha()
                                time.sleep(1)
                                new_img = checker.get_captcha_element_screenshot()
                                if new_img:
                                    with _credit_tasks_lock:
                                        credit_tasks[task_id]['captcha_image'] = new_img.getvalue()
                            time.sleep(1)

                        # Submit the CAPTCHA solution
                        with _credit_tasks_lock:
                            solution = credit_tasks[task_id]['captcha_solution']
                        checker.submit_captcha(solution)

                        # Clear CAPTCHA flags
                        with _credit_tasks_lock:
                            credit_tasks[task_id]['captcha_needed'] = False
                            credit_tasks[task_id]['captcha_image'] = None
                        time.sleep(3)   # wait for page to reload

                # Wait for user to confirm results
                if manual_mode:
                    with _credit_tasks_lock:
                        credit_tasks[task_id]['waiting'] = True
                        credit_tasks[task_id]['resume'] = False

                    manual_start = time.time()
                    manual_timeout = 600  # 10 minutes max wait
                    while True:
                        if time.time() - manual_start > manual_timeout:
                            logger.warning(f"Manual confirm wait timed out for task {task_id}")
                            with _credit_tasks_lock:
                                credit_tasks[task_id]['waiting'] = False
                            break
                        with _credit_tasks_lock:
                            resume = credit_tasks[task_id].get('resume', False)
                        if resume:
                            break
                        time.sleep(1)

                    with _credit_tasks_lock:
                        credit_tasks[task_id]['waiting'] = False
                else:
                    time.sleep(3)   # fallback delay for full-auto (unused)

                # Capture screenshot
                shot = checker.capture_viewport()
                company_shots.append(shot)
                logger.info(f"Screenshot captured for {company} at {url}")

            screenshots[company] = company_shots
            with _credit_tasks_lock:
                credit_tasks[task_id]['progress'] = idx + 1

        # ========== Generate Word Document ==========
        doc = Document()
        for section in doc.sections:
            section.top_margin = Cm(0.3)
            section.bottom_margin = Cm(0.3)
            section.left_margin = Cm(0.3)
            section.right_margin = Cm(0.3)

        first = True
        for company, shots in screenshots.items():
            if not first:
                doc.add_page_break()
            first = False
            doc.add_heading(company, level=1)
            for shot in shots:
                shot.seek(0)
                doc.add_picture(shot, width=Inches(7.2))
                doc.add_paragraph()

        doc_buffer = BytesIO()
        doc.save(doc_buffer)
        doc_buffer.seek(0)

        # Save report file
        os.makedirs(CREDIT_REPORTS_DIR, exist_ok=True)
        file_name = f"credit_report_{task_id}.docx"
        file_path = os.path.join(CREDIT_REPORTS_DIR, file_name)
        with open(file_path, 'wb') as f:
            f.write(doc_buffer.getvalue())

        # Insert into database
        with get_db_connection() as conn:
            with conn.cursor() as cur:
                cur.execute(
                    "INSERT INTO credit_check_reports (user_id, task_id, file_path, companies_count) VALUES (%s, %s, %s, %s)",
                    (user_id, task_id, file_path, len(companies))
                )
                conn.commit()

        # Mark task as completed
        with _credit_tasks_lock:
            credit_tasks[task_id]['status'] = 'completed'
        logger.info(f"Credit check task {task_id} finished successfully")

    except Exception as e:
        logger.error(f"Credit check task {task_id} failed: {e}", exc_info=True)
        with _credit_tasks_lock:
            credit_tasks[task_id]['status'] = 'error'
            credit_tasks[task_id]['error'] = str(e)
    finally:
        checker.close()

@app.route('/credit_check_status/<task_id>')
def credit_check_status(task_id):
    with _credit_tasks_lock:
        task = credit_tasks.get(task_id)
        if not task:
            return jsonify({"error": "Task not found"}), 404
        # Make a copy to release lock quickly
        result = {
            'status': task['status'],
            'progress': task['progress'],
            'total': task['total'],
            'captcha_needed': task.get('captcha_needed', False),
            'download_url': task.get('download_url'),
            'error': task.get('error'),
            'waiting': task.get('waiting', False)
        }
    return jsonify(result)

@app.route('/get_captcha_image/<task_id>')
def get_captcha_image(task_id):
    with _credit_tasks_lock:
        task = credit_tasks.get(task_id)
        if not task or not task.get('captcha_needed') or not task.get('captcha_image'):
            return "No captcha image", 404
        img_bytes = task['captcha_image']
    return send_file(BytesIO(img_bytes), mimetype='image/png')

@app.route('/reload_captcha/<task_id>', methods=['POST'])
def reload_captcha(task_id):
    with _credit_tasks_lock:
        if task_id in credit_tasks:
            credit_tasks[task_id]['reload_captcha'] = True
    return jsonify({"status": "reloading"})

@app.route('/solve_captcha/<task_id>', methods=['POST'])
def solve_captcha(task_id):
    data = request.get_json()
    solution = data.get('solution', '')
    with _credit_tasks_lock:
        if task_id in credit_tasks:
            credit_tasks[task_id]['captcha_solution'] = solution
    return jsonify({"status": "ok"})

@app.route('/credit_check_resume/<task_id>', methods=['POST'])
def credit_check_resume(task_id):
    with _credit_tasks_lock:
        if task_id in credit_tasks:
            credit_tasks[task_id]['resume'] = True
    return jsonify({"status": "ok"})

@app.route('/download_credit_report/<task_id>')
def download_credit_report(task_id):
    user_id = get_user_id()
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(
                "SELECT file_path FROM credit_check_reports WHERE task_id = %s AND user_id = %s",
                (task_id, user_id)
            )
            row = cur.fetchone()
            if not row:
                return jsonify({"error": "Report not found or access denied"}), 404
            file_path = row[0]
            if os.path.exists(file_path):
                return send_file(file_path, as_attachment=True, download_name=f"credit_report_{task_id}.docx")
            return jsonify({"error": "Report file missing"}), 404

@app.route('/list_credit_reports')
def list_credit_reports():
    user_id = get_user_id()
    with get_db_connection() as conn:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("""
                SELECT id, task_id, companies_count, created_at
                FROM credit_check_reports
                WHERE user_id = %s
                ORDER BY created_at DESC
            """, (user_id,))
            reports = cur.fetchall()
    return jsonify({"reports": reports})

@app.route('/delete_credit_report/<int:report_id>', methods=['POST'])
def delete_credit_report(report_id):
    user_id = get_user_id()
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT file_path FROM credit_check_reports WHERE id = %s AND user_id = %s", (report_id, user_id))
            row = cur.fetchone()
            if not row:
                return jsonify({"error": "Report not found"}), 404
            file_path = row[0]
            if os.path.exists(file_path):
                os.remove(file_path)
            cur.execute("DELETE FROM credit_check_reports WHERE id = %s", (report_id,))
            conn.commit()
    return jsonify({"success": True})

# ---------- Shutdown helpers (called by app factory via atexit) ----------

def shutdown_agent():
    global _agent
    _agent = None

def shutdown_db_pool():
    db_pool.closeall()
    logger.info("Database pool closed.")

def query_shadowbot(query):
    """ShadowBot is currently disabled."""
    logger.debug(f"query_shadowbot called with: {query}")
    return None