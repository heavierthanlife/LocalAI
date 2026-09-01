"""File processing, text extraction, and similarity (auto-extracted)."""
import os, re, io, time, hashlib, logging, tempfile, subprocess, html, difflib
import threading
import numpy as np
import pandas as pd
import openpyxl
from datetime import datetime
from io import BytesIO
from PIL import Image
import imagehash
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from markitdown import MarkItDown
import docx
from pptx import Presentation
from flask import session

try:
    import pymupdf as fitz
except ImportError:
    import fitz

from app.config import is_valid_extracted_text, ALLOWED_EXTENSIONS, logger
from app.database import get_db_connection
from app.utils.helpers import safe_error_response, utc_now
import app.globals as g
from app.services.semantic import get_semantic_model
from app.services.vl_model import vl_model
from app.services.ocr import ocr_manager, run_ocr
from app.services.file_cache import add_to_cache
from app.services.session_manager import get_cached_image_description, cache_image_description

# ── VL circuit breaker ────────────────────────────────────────────────────
# After N consecutive VL failures (e.g. misconfigured endpoint returning 404),
# skip all further VL calls in this process — prevents dozens of slow failing
# requests from stalling analysis for minutes. A per-file cumulative counter
# also trips the breaker when many failures accumulate even if not consecutive.
_VL_FAIL_THRESHOLD = 3
_VL_CUMULATIVE_THRESHOLD = 8
_vl_fail_lock = threading.Lock()
_vl_consecutive_fails = 0
_vl_broken = False
_vl_file_fails = 0


def _vl_note_success():
    global _vl_consecutive_fails, _vl_broken
    with _vl_fail_lock:
        _vl_consecutive_fails = 0
        _vl_broken = False


def _vl_reset_file():
    """Reset per-file cumulative failure counter (call at start of a file)."""
    global _vl_file_fails
    with _vl_fail_lock:
        _vl_file_fails = 0


def _vl_note_failure() -> bool:
    """Record a VL failure. Returns True if the breaker just tripped."""
    global _vl_consecutive_fails, _vl_broken, _vl_file_fails
    with _vl_fail_lock:
        _vl_consecutive_fails += 1
        _vl_file_fails += 1
        if (_vl_consecutive_fails >= _VL_FAIL_THRESHOLD
                or _vl_file_fails >= _VL_CUMULATIVE_THRESHOLD) and not _vl_broken:
            _vl_broken = True
            logger.warning(
                f"VL failed {_vl_consecutive_fails}x consec / {_vl_file_fails}x total — "
                f"circuit breaker OPEN, skipping remaining VL calls")
            return True
        return False


def _vl_available():
    """vl_model.is_available() AND breaker not open."""
    if _vl_broken:
        return False
    return vl_model.is_available()


# ── Image sampling log (per-process) ──────────────────────────────────────
# Collects {filename, samples} for the clearance report's 图片随机抽检说明.
_IMAGE_SAMPLING_LOG = []
_IMAGE_SAMPLING_LOCK = threading.Lock()


def take_image_sampling_log():
    """Atomically drain and return the accumulated sampling log."""
    global _IMAGE_SAMPLING_LOG
    with _IMAGE_SAMPLING_LOCK:
        out = list(_IMAGE_SAMPLING_LOG)
        _IMAGE_SAMPLING_LOG = []
        return out


def _trunc(s, n):
    s = (s or '').replace('\n', ' ').strip()
    return s[:n] if len(s) > n else s


def _describe_sampled_images(items, filename):
    """Describe a sampled list of (image_bytes, anchor) with batched multi-image
    VL calls (VL_BATCH default 5) across a bounded thread pool (VL_PARALLEL=2).
    Returns (description_lines, sample_info_rows).

    items: list of dicts {blob, chapter, prev, next, prompt}
    """
    import random
    from concurrent.futures import ThreadPoolExecutor, as_completed

    max_img = int(os.getenv('VL_MAX_IMAGES_PER_FILE', '20'))
    parallel = int(os.getenv('VL_PARALLEL', '2'))
    batch = int(os.getenv('VL_BATCH', '5'))
    total = len(items)
    chosen = random.sample(items, min(max_img, total))

    lines = []
    rows = []
    lock = threading.Lock()

    def _run_batch(batch_items, base_seq):
        """One multi-image VL call; returns list of per-item results."""
        if _vl_broken:
            return [None] * len(batch_items)
        blobs = [it['blob'] for it in batch_items]
        descs = vl_model.describe_images_batch(
            blobs, '请逐一描述下列每张图片。每张图严格用单独一行，行首是序号：图1、图2、…，'
                   '然后是本图内容摘要（40字内，含主要文字/数字）。不要添加任何多余段落。'
                   '如果某张图看不清，该行写：图N：无法识别。')
        out = []
        for i, it in enumerate(batch_items):
            desc = (descs[i] if descs else None)
            if desc and not desc.startswith('⚠'):
                _vl_note_success()
                out.append({'seq': base_seq + i, 'chapter': it.get('chapter', ''),
                            'prev': it.get('prev', ''), 'next': it.get('next', ''), 'desc': desc})
            elif descs is None:
                # whole API call failed — counts toward the breaker
                _vl_note_failure()
                out.append(None)
            else:
                # model returned empty/unparseable for this image — API OK, skip silently
                out.append(None)
        return out

    batches = [chosen[i:i + batch] for i in range(0, len(chosen), batch)]
    with ThreadPoolExecutor(max_workers=parallel) as pool:
        futures = [pool.submit(_run_batch, b, i * batch + 1) for i, b in enumerate(batches)]
        for fut in as_completed(futures):
            res = fut.result() or []
            with lock:
                for r in res:
                    if r is None:
                        continue
                    rows.append(r)
                    lines.append(
                        f"[抽检图片{r['seq']}/{len(chosen)} 位置:{_trunc(r.get('chapter', ''), 40)}]: {r['desc']}")

    rows.sort(key=lambda r: r['seq'])
    if total > len(chosen):
        lines.append(f"[图片抽检说明] 本文件共 {total} 张图片，随机抽检 {len(chosen)} 张；其余 {total - len(chosen)} 张未逐一识别以控制耗时。")
    if _vl_broken and not lines:
        lines.append("[图片抽检说明] VL 连续/累计失败已熔断，未识别图片内容。")
    return lines, rows


def describe_images_in_file(file_bytes, filename, page_texts=None):
    """Describe images in a document — RANDOMLY samples up to
    VL_MAX_IMAGES_PER_FILE (default 20) images, in parallel (VL_PARALLEL=2),
    and records each sampled image's location + surrounding text.

    Returns (text, sample_info):
      text:  assembled description text (for the extracted document text)
      sample_info: [{'seq','chapter','prev','next','desc'}, ...] for the
                   report's 图片随机抽检说明 section.
    """
    if not _vl_available():
        return "", []
    _vl_reset_file()
    ext = os.path.splitext(filename)[1].lower()
    items = []
    try:
        if ext == '.pdf':
            doc = fitz.open(stream=BytesIO(file_bytes), filetype="pdf")
            try:
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    for img_idx, img in enumerate(page.get_images(full=True)):
                        base_image = doc.extract_image(img[0])
                        page_txt = (page.get_text() or '').strip()
                        items.append({
                            'blob': base_image["image"],
                            'chapter': f"第{page_num + 1}页",
                            'prev': _trunc(page_txt, 10),
                            'next': '',
                            'prompt': f"Describe this image from page {page_num+1} of the PDF.",
                        })
            finally:
                doc.close()
        elif ext in ['.docx', '.docm']:
            import docx
            doc = docx.Document(BytesIO(file_bytes))
            # map rId -> image rel
            rels = doc.part.rels
            import re as _re
            pat = _re.compile(r'r:embed="(rId\d+)"')
            paras = doc.paragraphs
            seen = set()
            for pi, para in enumerate(paras):
                for m in pat.finditer(para._element.xml):
                    rid = m.group(1)
                    if rid in seen or rid not in rels or 'image' not in rels[rid].target_ref:
                        continue
                    seen.add(rid)
                    cur_txt = para.text.strip()
                    prev_txt = ''
                    for j in range(pi - 1, -1, -1):
                        t = paras[j].text.strip()
                        if t:
                            prev_txt = t
                            break
                    nxt_txt = ''
                    for j in range(pi + 1, len(paras)):
                        t = paras[j].text.strip()
                        if t:
                            nxt_txt = t
                            break
                    items.append({
                        'blob': rels[rid].target_part.blob,
                        'chapter': _trunc(cur_txt or prev_txt or f"第{pi + 1}段", 30),
                        'prev': _trunc(prev_txt, 10),
                        'next': _trunc(nxt_txt, 10),
                        'prompt': "Describe this image from the Word document.",
                    })
                    if _vl_broken:
                        break
        elif ext in ['.pptx', '.pptm']:
            from pptx import Presentation
            prs = Presentation(BytesIO(file_bytes))
            slide_num = 1
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == 13:
                        items.append({
                            'blob': shape.image.blob,
                            'chapter': f"第{slide_num}页幻灯片",
                            'prev': '', 'next': '',
                            'prompt': f"Describe this image from slide {slide_num}.",
                        })
                slide_num += 1
    except Exception as e:
        logger.error(f"Error collecting images from {filename}: {e}")
        return "", []

    lines, rows = _describe_sampled_images(items, filename)
    return "\n".join(lines), rows


# Agent is managed by app.services.agent.get_agent (imported as _svc_get_agent)
# The local wrapper handles max_tokens via session state
def get_agent(max_tokens=None):
    from app.services.agent import get_agent as _agent_get
    return _agent_get(max_tokens)


# Text Preprocessing & Similarity
TECH_STD_PATTERNS = [
    r'GB/T\s*\d+\.?\d*', r'GB\s*\d+\.?\d*', r'ISO\s*\d+', r'IEC\s*\d+',
    r'IEEE\s*\d+', r'DIN\s*\d+', r'BS\s*\d+', r'EN\s*\d+', r'ASME\s*\d+',
    r'API\s*\d+', r'ASTM\s*\d+', r'JJG\s*\d+', r'JB/T\s*\d+', r'HG/T\s*\d+',
    r'SY/T\s*\d+', r'DL/T\s*\d+', r'NB/T\s*\d+', r'SH/T\s*\d+', r'YS/T\s*\d+',
    r'FZ/T\s*\d+', r'QB/T\s*\d+', r'CJ/T\s*\d+', r'JG/T\s*\d+', r'GA/T\s*\d+',
    r'HS/T\s*\d+', r'行业标准', r'国家标准', r'技术规范'
]
SENSITIVE_PATTERNS = [
    r'(公司|集团|有限|股份|组织|委员会|协会|研究院|大学|学院)',
    r'(北京|上海|广州|深圳|杭州|南京|武汉|成都|重庆|天津|西安)',
    r'(项目|工程|系统|平台|软件|硬件|方案)',
    r'(张|王|李|刘|陈|杨|赵|黄|周|吴|徐|孙|马|朱|胡|林|郭|何|高)',
    r'(一等奖|二等奖|三等奖|金奖|银奖|优秀奖)',
    r'\d{17}[\dXx]',
    r'1[3-9]\d{9}',
    r'\d{18}',
    r'证书编号[：:]\s*\w+',
]


# ── jieba-aware TfidfVectorizer helper ──
def _make_vectorizer(stop_words=None, **kwargs):
    """Create a TfidfVectorizer with jieba Chinese tokenizer.

    ``stop_words`` (set) is applied inside the tokenizer (not sklearn's stop_words,
    which is English-oriented). None → use DEFAULT_STOP_WORDS (FIX-013).
    English text uses the default analyzer.
    """
    from app.services.text_utils import tokenize_for_tfidf, has_chinese
    from app.services.stop_words import DEFAULT_STOP_WORDS
    from sklearn.feature_extraction.text import TfidfVectorizer

    if stop_words is None:
        stop_words = DEFAULT_STOP_WORDS

    def _tokenizer(raw_text):
        if has_chinese(raw_text):
            return tokenize_for_tfidf(raw_text, stop_words=stop_words)
        return raw_text  # English text: let default analyzer handle it

    return TfidfVectorizer(tokenizer=_tokenizer, token_pattern=None, **kwargs)


def preprocess_text_for_similarity(text, template_text=None):
    """Preprocess text for similarity: strip noise, tokenize Chinese with jieba, filter.

    Chinese text: jieba segmentation → space-joined for TfidfVectorizer.
    English text: whitespace split with length filter.
    FIX-013: filters high-frequency stop words; if template_text (招标文件) is
    provided, its top-50 high-frequency words are added to the stop set.
    """
    if not text:
        return ""
    from app.services.text_utils import tokenize_for_tfidf, has_chinese, top_keywords
    from app.services.stop_words import DEFAULT_STOP_WORDS

    stop_words = set(DEFAULT_STOP_WORDS)
    if template_text and template_text.strip():
        for w, _ in top_keywords(template_text, top_k=50):
            stop_words.add(w)

    text = re.sub(r'[^\w\u4e00-\u9fff\s]', '', text)
    if has_chinese(text):
        filtered = tokenize_for_tfidf(text, min_len=2, stop_words=stop_words)
    else:
        words = text.split()
        filtered = ' '.join(w for w in words if len(w) >= 3 and w.lower() not in stop_words)
    for pat in TECH_STD_PATTERNS:
        text = re.sub(pat, '', text, flags=re.IGNORECASE)
    text = re.sub(r'^目录|^第[一二三四五六七八九十]+章', '', text, flags=re.MULTILINE)
    for pat in SENSITIVE_PATTERNS:
        text = re.sub(pat, '', text)
    # Preserve jieba-tokenized space-separated words for similarity calculation
    if filtered.strip():
        return filtered
    text = re.sub(r'\s+', ' ', text).strip()
    return text

def remove_template_content(text, template_text, threshold=0.85):
    """Remove template boilerplate using paragraph-level jieba-optimized TF-IDF."""
    if not template_text or not text:
        return text
    paras = [p.strip() for p in text.split('\n') if len(p.strip()) > 10]
    template_paras = [p.strip() for p in template_text.split('\n') if len(p.strip()) > 10]
    if not paras or not template_paras:
        return text
    all_paras = paras + template_paras
    vectorizer = _make_vectorizer(stop_words=None, lowercase=True).fit(all_paras)
    vecs = vectorizer.transform(all_paras)
    para_vecs = vecs[:len(paras)]
    template_vecs = vecs[len(paras):]
    sim_matrix = cosine_similarity(para_vecs, template_vecs)
    keep_mask = np.max(sim_matrix, axis=1) < threshold
    kept_paras = [p for i, p in enumerate(paras) if keep_mask[i]]
    if not kept_paras:
        return "[Template content fully matched] " + text
    return '\n'.join(kept_paras)

def extract_keywords(text, top_k=20):
    """Extract keywords using jieba-optimized TF-IDF (works for Chinese)."""
    if not text.strip():
        return []
    vectorizer = _make_vectorizer(stop_words=None, max_features=top_k)
    try:
        tfidf = vectorizer.fit_transform([text])
        feature_names = vectorizer.get_feature_names_out()
        scores = tfidf.toarray()[0]
        keyword_score = sorted(zip(feature_names, scores), key=lambda x: x[1], reverse=True)
        return [kw for kw, _ in keyword_score[:top_k]]
    except Exception:
        return []

def keyword_overlap_similarity(text1, text2):
    kw1 = set(extract_keywords(text1, 20))
    kw2 = set(extract_keywords(text2, 20))
    if not kw1 and not kw2:
        return 0.0
    # Guard against template/boilerplate-only texts: if either side has too few
    # distinguishing keywords, a high Jaccard is spurious (both sides just share
    # generic 招标 boilerplate). Return a low floor instead of an inflated value.
    if len(kw1) < 4 or len(kw2) < 4:
        return 0.05
    inter = len(kw1 & kw2)
    union = len(kw1 | kw2)
    return inter / union if union > 0 else 0.0

def compute_similarity_with_numbers(text1, text2, template_text=None):
    clean1 = preprocess_text_for_similarity(text1, template_text)
    clean2 = preprocess_text_for_similarity(text2, template_text)
    if template_text:
        clean1 = remove_template_content(clean1, template_text)
        clean2 = remove_template_content(clean2, template_text)
    if not clean1.strip() or not clean2.strip():
        return 0.0, text1, text2, []
    vectorizer = _make_vectorizer(stop_words=None, lowercase=True)
    tfidf = vectorizer.fit_transform([clean1, clean2])
    sim = cosine_similarity(tfidf[0:1], tfidf[1:2])[0][0]
    escaped1 = html.escape(text1)
    escaped2 = html.escape(text2)
    matcher = difflib.SequenceMatcher(None, escaped1, escaped2)
    matching_blocks = matcher.get_matching_blocks()
    segments1 = []
    last_idx = 0
    match_counter = 1
    blocks_detail = []
    for block in matching_blocks:
        i, j, size = block
        if size == 0 or size <= 6:
            continue
        blocks_detail.append({
            "id": match_counter,
            "pos1": i,
            "pos2": j,
            "size": size,
            "text1_snippet": escaped1[i:i + min(size, 100)] + ("..." if size > 100 else ""),
            "text2_snippet": escaped2[j:j + min(size, 100)] + ("..." if size > 100 else "")
        })
        if i > last_idx:
            segments1.append(('text', escaped1[last_idx:i]))
        match_text = escaped1[i:i + size]
        color_class = 'match-highlight-long' if size > 100 else 'match-highlight-short'
        marker = f"<sup><small>[{match_counter}]</small></sup> "
        segments1.append(('match', match_text, marker, color_class))
        last_idx = i + size
        match_counter += 1
    if last_idx < len(escaped1):
        segments1.append(('text', escaped1[last_idx:]))
    segments2 = []
    last_idx = 0
    match_counter = 1
    for block in matching_blocks:
        i, j, size = block
        if size == 0 or size <= 6:
            continue
        if j > last_idx:
            segments2.append(('text', escaped2[last_idx:j]))
        match_text = escaped2[j:j + size]
        color_class = 'match-highlight-long' if size > 100 else 'match-highlight-short'
        marker = f"<sup><small>[{match_counter}]</small></sup> "
        segments2.append(('match', match_text, marker, color_class))
        last_idx = j + size
        match_counter += 1
    if last_idx < len(escaped2):
        segments2.append(('text', escaped2[last_idx:]))

    def build_html(segments):
        parts = []
        for seg in segments:
            if seg[0] == 'text':
                parts.append(seg[1])
            else:
                _, text, marker, color_class = seg
                parts.append(marker + f'<span class="{color_class}">{text}</span>')
        return ''.join(parts)

    html1 = build_html(segments1)
    html2 = build_html(segments2)
    return sim, html1, html2, blocks_detail

def compute_batch_semantic_similarity(texts, lang_code=None):
    """Compute pairwise semantic similarity matrix with language-aware model selection.
    
    Uses the new multi-model system: auto-detects language or uses provided lang_code.
    Falls back to legacy distiluse model if needed.
    """
    from app.services.semantic import compute_batch_semantic_similarity as _sem
    return _sem(texts, lang_code=lang_code)

def file_attr_similarity(meta1, meta2):
    score = 0.0
    if meta1.get('author') and meta2.get('author') and meta1['author'] == meta2['author']:
        score += 50
    try:
        if meta1.get('creationDate') and meta2.get('creationDate'):
            date1 = re.sub(r'D:', '', meta1['creationDate'])[:14]
            date2 = re.sub(r'D:', '', meta2['creationDate'])[:14]
            if date1 == date2:
                score += 30
    except Exception:
        pass
    name1 = meta1.get('filename', '')
    name2 = meta2.get('filename', '')
    if name1 and name2:
        common = len(set(name1.lower()) & set(name2.lower()))
        total = len(set(name1.lower()) | set(name2.lower()))
        if total > 0:
            score += (common / total) * 20
    return min(score, 100.0)

def _extract_doc_ole(file_bytes):
    """Pure-Python .doc extraction using olefile to read WordDocument stream."""
    try:
        import olefile
        ole = olefile.OleFileIO(io.BytesIO(file_bytes))
        # Try to read the WordDocument stream
        if ole.exists('WordDocument'):
            data = ole.openstream('WordDocument').read()
            # Extract text between Unicode markers
            text = data.decode('utf-16-le', errors='ignore')
            # Clean up non-printable chars
            import re
            text = re.sub(r'[^\u4e00-\u9fff\w\s.,;:!?()\-\u0020-\u007e]', '', text)
            text = re.sub(r'\s+', ' ', text).strip()
            if len(text) > 50:
                ole.close()
                return text
            ole.close()
    except Exception as e:
        logger.debug(f"OLE extraction failed: {e}")
    return None

def _convert_doc_via_soffice(file_bytes):
    """Convert .doc to .docx using LibreOffice headless (skill-recommended)."""
    try:
        with tempfile.NamedTemporaryFile(suffix='.doc', delete=False) as f:
            f.write(file_bytes)
            doc_path = f.name
        out_dir = tempfile.mkdtemp()
        result = subprocess.run(
            ['soffice', '--headless', '--convert-to', 'docx', '--outdir', out_dir, doc_path],
            capture_output=True, timeout=60
        )
        if result.returncode == 0:
            for fn in os.listdir(out_dir):
                if fn.endswith('.docx'):
                    with open(os.path.join(out_dir, fn), 'rb') as f:
                        return f.read()
    except (FileNotFoundError, subprocess.TimeoutExpired, Exception) as e:
        logger.debug(f"soffice conversion failed: {e}")
    finally:
        try: os.unlink(doc_path)
        except OSError: pass
        try: shutil.rmtree(out_dir, ignore_errors=True)
        except OSError: pass
    return None

def extract_text_from_doc_crossplatform(file_bytes):
    # Try 1: LibreOffice headless convert .doc → .docx (skill-recommended)
    docx_bytes = _convert_doc_via_soffice(file_bytes)
    if docx_bytes:
        try:
            doc = docx.Document(io.BytesIO(docx_bytes))
            text = '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
            if text and len(text) > 50:
                return text
        except Exception as e:
            logger.debug(f"docx parse after soffice failed: {e}")

    # Try 2: olefile pure-Python
    text = _extract_doc_ole(file_bytes)
    if text:
        return text

    with tempfile.NamedTemporaryFile(suffix='.doc', delete=False) as f:
        f.write(file_bytes)
        temp_doc = f.name
    try:
        # Try antiword
        try:
            result = subprocess.run(['antiword', temp_doc], capture_output=True, timeout=30)
            if result.returncode == 0 and result.stdout:
                for enc in ('gbk', 'gb18030', 'utf-8'):
                    try:
                        decoded = result.stdout.decode(enc)
                        if decoded.strip():
                            return decoded
                    except UnicodeDecodeError:
                        continue
                return result.stdout.decode('utf-8', errors='replace')
        except FileNotFoundError:
            logger.warning("antiword not installed")
        except Exception as e:
            logger.warning(f"antiword error: {e}")

        # Try catdoc
        try:
            result = subprocess.run(['catdoc', '-a', temp_doc], capture_output=True, timeout=30)
            if result.returncode == 0 and result.stdout:
                for enc in ('gbk', 'gb18030', 'utf-8'):
                    try:
                        decoded = result.stdout.decode(enc)
                        if decoded.strip():
                            return decoded
                    except UnicodeDecodeError:
                        continue
                return result.stdout.decode('utf-8', errors='replace')
        except FileNotFoundError:
            logger.warning("catdoc not installed")
        except Exception as e:
            logger.warning(f"catdoc error: {e}")

        return None
    finally:
        if os.path.exists(temp_doc):
            os.unlink(temp_doc)

def extract_text_from_doc(file_bytes):
    """Extract text from .doc using Win32 COM (Windows only, requires Word installed)."""
    import sys
    if sys.platform != 'win32':
        return None  # Win32 COM only available on Windows
    import pythoncom
    import win32com.client as win32
    import tempfile
    import os
    import time

    pythoncom.CoInitialize()
    temp_doc = None
    temp_txt = None
    word = None
    doc = None
    try:
        fd, temp_doc = tempfile.mkstemp(suffix='.doc')
        os.close(fd)
        with open(temp_doc, 'wb') as f:
            f.write(file_bytes)

        fd, temp_txt = tempfile.mkstemp(suffix='.txt')
        os.close(fd)

        word = win32.Dispatch("Word.Application")
        word.Visible = False
        word.DisplayAlerts = 0
        doc = word.Documents.Open(temp_doc, AddToRecentFiles=False, Visible=False)
        doc.SaveAs2(temp_txt, FileFormat=2)
        doc.Close(SaveChanges=False)
        word.Quit(SaveChanges=False)

        time.sleep(0.1)

        for enc in ('gbk', 'gb18030', 'utf-8', 'gb2312', 'latin-1'):
            try:
                with open(temp_txt, 'r', encoding=enc) as f:
                    text = f.read()
                if text.strip():
                    logger.info(f"Successfully decoded .doc with {enc}")
                    return text
            except UnicodeDecodeError:
                continue
        with open(temp_txt, 'r', encoding='utf-8', errors='replace') as f:
            text = f.read()
        return text
    except Exception as e:
        logger.error(f"win32com .doc extraction failed: {e}")
        return None
    finally:
        if doc:
            try:
                doc.Close(False)
            except Exception:
                pass
        if word:
            try:
                word.Quit(False)
            except Exception:
                pass
        if temp_doc and os.path.exists(temp_doc):
            try:
                os.unlink(temp_doc)
            except Exception:
                pass
        if temp_txt and os.path.exists(temp_txt):
            try:
                os.unlink(temp_txt)
            except Exception:
                pass
        pythoncom.CoUninitialize()

def detect_excel_format(file_bytes):
    if len(file_bytes) < 8:
        return None
    if file_bytes[:8] == b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1':
        return 'xls'
    if file_bytes[:2] == b'PK':
        return 'xlsx'
    return None

def extract_text_from_xls(file_bytes):
    try:
        import xlrd
        workbook = xlrd.open_workbook(file_contents=file_bytes)
        text_parts = []
        for sheet in workbook.sheets():
            sheet_text = []
            for row in range(sheet.nrows):
                row_text = " ".join(str(cell.value) for cell in sheet.row(row) if cell.value)
                if row_text.strip():
                    sheet_text.append(row_text)
            if sheet_text:
                text_parts.append(f"--- Sheet: {sheet.name} ---\n" + "\n".join(sheet_text))
        return "\n\n".join(text_parts) if text_parts else "[No text in Excel]"
    except Exception as e:
        return f"[Excel parsing error (old format): {e}]"

def detect_word_format(file_bytes):
    if len(file_bytes) < 8:
        return None
    if file_bytes[:8] == b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1':
        return 'doc'
    if file_bytes[:2] == b'PK':
        return 'docx'
    return None

def extract_images_from_file(file_storage):
    images = []
    ext = os.path.splitext(file_storage.filename)[1].lower()
    file_bytes = file_storage.read()
    file_storage.seek(0)
    if ext == '.pdf':
        doc = None
        try:
            doc = fitz.open(stream=BytesIO(file_bytes), filetype="pdf")
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                img_list = page.get_images(full=True)
                for img in img_list:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image = Image.open(BytesIO(image_bytes))
                    images.append(image)
        except Exception:
            pass
        finally:
            if doc:
                doc.close()
    return images

def image_similarity(images1, images2):
    if not images1 or not images2:
        return 0.0
    max_sim = 0.0
    for img1 in images1:
        hash1 = imagehash.phash(img1)
        for img2 in images2:
            hash2 = imagehash.phash(img2)
            sim = 1 - (hash1 - hash2) / 64.0
            max_sim = max(max_sim, sim)
    return max_sim * 100

def extract_metadata(file_storage):
    meta = {'filename': file_storage.filename}
    ext = os.path.splitext(file_storage.filename)[1].lower()
    file_bytes = file_storage.read()
    file_storage.seek(0)
    if ext == '.pdf':
        try:
            doc = fitz.open(stream=BytesIO(file_bytes), filetype="pdf")
            info = doc.metadata
            meta['author'] = info.get('author', '')
            meta['creator'] = info.get('creator', '')
            meta['producer'] = info.get('producer', '')
            meta['creationDate'] = info.get('creationDate', '')
        except Exception:
            pass
    elif ext in ['.docx', '.docm']:
        try:
            doc = docx.Document(BytesIO(file_bytes))
            core_props = doc.core_properties
            meta['author'] = core_props.author or ''
            meta['created'] = core_props.created
            meta['modified'] = core_props.modified
        except Exception:
            pass
    return meta

def truncate_filename(filename, max_len=40):
    if len(filename) <= max_len:
        return filename
    name, ext = os.path.splitext(filename)
    if len(ext) > 10:
        ext = ext[:10]
    available = max_len - len(ext) - 3
    if available < 1:
        ext = ext[:max_len]
        return ext
    truncated_name = name[:available] + '...'
    return truncated_name + ext

def detect_file_real_type(file_bytes: bytes) -> str:
    if len(file_bytes) < 8:
        return None
    if file_bytes[:8] == b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1':
        return 'doc'
    if file_bytes[:4] == b'PK\x03\x04':
        return 'docx'
    return None


# Advanced extraction using fallbacks
# KREUZBERG_AVAILABLE, UNSTRUCTURED_AVAILABLE, KREUZBERG_SIZE_LIMIT imported from app.globals
KREUZBERG_AVAILABLE = g.KREUZBERG_AVAILABLE
UNSTRUCTURED_AVAILABLE = g.UNSTRUCTURED_AVAILABLE
KREUZBERG_SIZE_LIMIT = g.KREUZBERG_SIZE_LIMIT

def extract_text_advanced(file_bytes: bytes, filename: str, file_size: int) -> str:
    ext = os.path.splitext(filename)[1].lower()
    office_extensions = {'.docx', '.xlsx', '.pptx', '.xlsm', '.xlsb', '.pptm', '.doc'}
    if ext not in office_extensions:
        return None

    # Try MarkItDown first
    try:
        from markitdown import MarkItDown
        md = MarkItDown()
        result = md.convert_stream(BytesIO(file_bytes), file_extension=ext.lstrip('.'))
        text = result.text_content
        if text and text.strip():
            logger.info(f"MarkItDown extracted {len(text)} chars from {filename}")
            return text
    except Exception as e:
        logger.warning(f"MarkItDown failed for {filename}: {e}")

    # For .doc, try legacy
    if ext == '.doc':
        text = extract_text_from_doc_crossplatform(file_bytes)
        if text and text.strip():
            return text
        text = extract_text_from_doc(file_bytes)
        if text and text.strip():
            return text

    # For .docx, fallback to python-docx
    if ext == '.docx':
        try:
            doc = docx.Document(BytesIO(file_bytes))
            text = "\n".join([para.text for para in doc.paragraphs])
            for table in doc.tables:
                for row in table.rows:
                    row_text = "\t".join([cell.text for cell in row.cells])
                    text += "\n" + row_text
            if text.strip():
                return text
        except Exception as e:
            logger.warning(f"python-docx fallback failed for {filename}: {e}")

    return None


def _extract_docx_paragraphs(file_bytes: bytes) -> list:
    """Extract structured paragraphs from DOCX using python-docx.

    Returns list of {text, level, style} dicts:
    - text: paragraph text
    - level: heading level (0 for body, 1-6 for headings)
    - style: 'heading' or 'body'
    """
    paragraphs = []
    try:
        doc = docx.Document(BytesIO(file_bytes))
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style_name = (para.style.name or '').lower() if para.style else ''
            if 'heading' in style_name:
                try:
                    level = int(''.join(c for c in style_name if c.isdigit()) or '1')
                except ValueError:
                    level = 1
                paragraphs.append({'text': text, 'level': level, 'style': 'heading'})
            else:
                paragraphs.append({'text': text, 'level': 0, 'style': 'body'})
    except Exception as e:
        logger.warning(f"DOCX structured extraction failed: {e}")
    return paragraphs


def _extract_pdf_paragraphs(file_bytes: bytes) -> list:
    """Extract structured paragraphs from PDF.

    Returns list of {text, level, style} dicts. PDF doesn't have native heading
    styles, so all paragraphs are returned with level 0 (body).
    """
    paragraphs = []
    try:
        doc = fitz.open(stream=file_bytes, filetype='pdf')
        for page in doc:
            text = page.get_text()
            for line in text.split('\n'):
                line = line.strip()
                if line and len(line) > 2:
                    paragraphs.append({'text': line, 'level': 0, 'style': 'body'})
    except Exception as e:
        logger.warning(f"PDF structured extraction failed: {e}")
    return paragraphs


def extract_structured_text(file_storage) -> list:
    """Extract structured paragraphs from a file.

    Returns list of {text, level, style} dicts for use by document_parser.
    Supports .docx (with heading levels), .pdf (plain text), .doc (plain text).
    """
    filename = file_storage.filename
    if not filename:
        return []

    file_bytes = file_storage.read()
    file_storage.seek(0)
    ext = os.path.splitext(filename)[1].lower()

    # WPS aliases
    wps_map = {'.wps': '.doc', '.et': '.xls', '.dps': '.ppt'}
    if ext in wps_map:
        ext = wps_map[ext]

    if ext == '.docx':
        return _extract_docx_paragraphs(file_bytes)
    elif ext == '.pdf':
        return _extract_pdf_paragraphs(file_bytes)
    elif ext == '.doc':
        text = extract_text_from_doc_crossplatform(file_bytes) or extract_text_from_doc(file_bytes) or ''
        return [{'text': line, 'level': 0, 'style': 'body'} for line in text.split('\n') if line.strip()]
    else:
        return []


def extract_text_from_file(file_storage):
    filename = file_storage.filename
    if not filename:
        return None, {}

    file_bytes = file_storage.read()
    file_storage.seek(0)
    file_hash = hashlib.sha256(file_bytes).hexdigest()

    ext = os.path.splitext(filename)[1].lower()
    wps_map = {'.wps': '.doc', '.et': '.xls', '.dps': '.ppt'}
    original_ext = ext
    if ext in wps_map:
        ext = wps_map[ext]

    text = None
    page_texts = {}

    if ext in ['.txt', '.md', '.text', '.csv', '.json']:
        try:
            text = file_bytes.decode('utf-8')
        except UnicodeDecodeError:
            text = file_bytes.decode('utf-8', errors='replace')
        page_texts = {1: text}

    elif ext in ['.html', '.htm']:
        try:
            from bs4 import BeautifulSoup
            html_text = file_bytes.decode('utf-8', errors='replace')
            soup = BeautifulSoup(html_text, 'html.parser')
            text = soup.get_text(separator='\n', strip=True)
            if not text or not text.strip():
                text = html_text
            page_texts = {1: text}
        except Exception as e:
            logger.warning(f"HTML parsing failed, falling back to raw text: {e}")
            text = file_bytes.decode('utf-8', errors='replace')
            page_texts = {1: text}

    elif ext == '.pdf':
        try:
            doc = fitz.open(stream=BytesIO(file_bytes), filetype="pdf")
            full_text = []
            page_texts = {}
            has_text = False
            try:
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    page_text = page.get_text()
                    if page_text and page_text.strip():
                        has_text = True
                        full_text.append(page_text)
                        page_texts[page_num + 1] = page_text
                extracted = "\n".join(full_text).strip()
                if has_text and len(extracted) > 50:
                    text = extracted
                else:
                    logger.info("PDF appears to be scanned (no text). Starting OCR...")
                    if not ocr_manager.is_available():
                        if vl_model.is_available():
                            logger.info("Using VL model for scanned PDF")
                            extracted = ""
                            for page_num in range(len(doc)):
                                page = doc.load_page(page_num)
                                pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
                                img_bytes = pix.tobytes("png")
                                description = vl_model.describe_pdf_page(img_bytes, page_num + 1)
                                extracted += f"\n\n--- 第{page_num + 1}页 (VL分析) ---\n{description}"
                            page_texts = {i + 1: "" for i in range(len(doc))}
                            text = extracted
                        else:
                            text = "[无法提取PDF文本，且OCR/VL不可用]"
                    else:
                        ocr_results = []
                        ocr_page_texts = {}
                        zoom = 2.0
                        mat = fitz.Matrix(zoom, zoom)
                        for page_num in range(len(doc)):
                            page = doc.load_page(page_num)
                            pix = page.get_pixmap(matrix=mat, alpha=False)
                            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                            max_dim = 2000
                            if max(img.size) > max_dim:
                                ratio = max_dim / max(img.size)
                                new_size = (int(img.size[0] * ratio), int(img.size[1] * ratio))
                                img = img.resize(new_size, Image.Resampling.LANCZOS)
                            img_np = np.array(img)
                            page_text = run_ocr(img_np)
                            if page_text:
                                ocr_results.append(page_text)
                                ocr_page_texts[page_num + 1] = page_text
                            else:
                                ocr_results.append("")
                                ocr_page_texts[page_num + 1] = ""
                        if any(t.strip() for t in ocr_results):
                            extracted = "\n\n".join(ocr_results)
                            page_texts = ocr_page_texts
                            text = extracted
                        else:
                            if vl_model.is_available():
                                extracted = ""
                                for page_num in range(len(doc)):
                                    page = doc.load_page(page_num)
                                    pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
                                    img_bytes = pix.tobytes("png")
                                    description = vl_model.describe_pdf_page(img_bytes, page_num + 1)
                                    extracted += f"\n\n--- 第{page_num + 1}页 (VL分析) ---\n{description}"
                                page_texts = {i + 1: "" for i in range(len(doc))}
                                text = extracted
                            else:
                                text = "[No text detected in PDF even after OCR and VL not available]"
            finally:
                doc.close()
        except Exception as e:
            logger.error(f"PDF error: {e}", exc_info=True)
            text = safe_error_response("无法解析PDF文件，请确保文件未损坏。", log_error=e)

    elif ext in ['.docx', '.docm', '.dotx', '.dotm']:
        real_format = detect_word_format(file_bytes)
        if real_format == 'doc':
            text = extract_text_from_doc(file_bytes)
            if text:
                page_texts = {1: text}
            else:
                text = "[无法从 .doc 文件中提取文本。请转换为 .docx 格式后重试，或安装 antiword/catdoc。]"
        else:
            try:
                def extract_docx_text(byte_data):
                    doc = docx.Document(BytesIO(byte_data))
                    full_text = "\n".join([para.text for para in doc.paragraphs])
                    for table in doc.tables:
                        for row in table.rows:
                            row_text = "\t".join([cell.text for cell in row.cells])
                            full_text += "\n" + row_text
                    del doc
                    return full_text if full_text.strip() else "[No text in Word document]"

                text = extract_docx_text(file_bytes)
                page_texts = {1: text}
            except Exception as e:
                logger.error(f"DOCX parsing error: {e}")
                text = safe_error_response("无法解析Word文档，请转换为DOCX格式或检查文件。", log_error=e)

    elif ext in ['.xlsx', '.xlsm', '.xltx', '.xltm', '.xlsb']:
        wb = None
        try:
            wb = openpyxl.load_workbook(BytesIO(file_bytes), read_only=True, data_only=True)
            text_parts = []
            for sheet in wb.worksheets:
                sheet_text = []
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        sheet_text.append(row_text)
                if sheet_text:
                    text_parts.append(f"--- Sheet: {sheet.title} ---\n" + "\n".join(sheet_text))
            full_text = "\n\n".join(text_parts) if text_parts else "[No text in Excel]"
            text = full_text
            page_texts = {1: full_text}
        except Exception as e:
            logger.warning(f"openpyxl failed for {filename}: {e}. Trying fallback methods.")
            real_format = detect_excel_format(file_bytes)
            if real_format == 'xls':
                text = extract_text_from_xls(file_bytes)
                if text and not text.startswith("["):
                    page_texts = {1: text}
                else:
                    file_storage.seek(0)
                    md = MarkItDown()
                    result = md.convert(BytesIO(file_bytes), file_extension=original_ext.lstrip('.'))
                    text = result.text_content
                    if text and text.strip():
                        page_texts = {1: text}
                    else:
                        text = safe_error_response("无法解析Excel文档，请转换为xlsx格式或检查文件。", log_error=e)
            else:
                file_storage.seek(0)
                md = MarkItDown()
                result = md.convert(BytesIO(file_bytes), file_extension=original_ext.lstrip('.'))
                text = result.text_content
                if text and text.strip():
                    page_texts = {1: text}
                else:
                    text = safe_error_response("无法解析Excel文档，请转换为xlsx格式或检查文件。", log_error=e)
        finally:
            if wb:
                wb.close()

    elif ext in ['.pptx', '.pptm', '.potx', '.ppsx']:
        prs = None
        try:
            prs = Presentation(BytesIO(file_bytes))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            full_text = "\n".join(text_runs)
            text = full_text if full_text.strip() else "[No text in PowerPoint]"
            page_texts = {1: full_text}
        except Exception as e:
            logger.error(f"PPTX parsing error: {e}")
            text = safe_error_response("无法解析PowerPoint文件，请检查文件。", log_error=e)
        finally:
            if prs is not None:
                del prs

    elif ext == '.xls':
        xls = None
        try:
            xls = pd.ExcelFile(BytesIO(file_bytes), engine='xlrd')
            text_parts = []
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                sheet_text = df.to_string(index=False, header=True)
                if sheet_text.strip():
                    text_parts.append(f"--- Sheet: {sheet_name} ---\n{sheet_text}")
            full_text = "\n\n".join(text_parts) if text_parts else "[No text in Excel]"
            text = full_text
            page_texts = {1: full_text}
        except Exception as e:
            logger.error(f"XLS parsing error: {e}")
            text = f"[Excel parsing failed: {str(e)}]"
        finally:
            if xls is not None:
                try:
                    xls.close()
                except Exception:
                    pass
    elif ext == '.doc':
        # Try cross-platform methods first (olefile, soffice, antiword, catdoc)
        text = extract_text_from_doc_crossplatform(file_bytes)
        if not (text and text.strip()):
            # Fall back to win32com (Windows + Word required)
            text = extract_text_from_doc(file_bytes)
        if text and text.strip():
            page_texts = {1: text}
        else:
            text = "[无法从 .doc 文件中提取文本。请转换为 .docx 格式后重试。]"
    elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
        if ocr_manager.is_available():
            try:
                image = Image.open(BytesIO(file_bytes))
                if image.mode != 'RGB':
                    image = image.convert('RGB')
                max_dim = 2000
                if max(image.size) > max_dim:
                    ratio = max_dim / max(image.size)
                    new_size = (int(image.size[0] * ratio), int(image.size[1] * ratio))
                    image = image.resize(new_size, Image.Resampling.LANCZOS)
                img_np = np.array(image)
                extracted_text = run_ocr(img_np)
                if extracted_text:
                    text = extracted_text
                    page_texts = {1: extracted_text}
            except Exception as e:
                logger.warning(f"OCR failed for image: {e}")
        if not text and vl_model.is_available():
            logger.info("Using VL model for image description")
            description = vl_model.describe_image(file_bytes)
            text = description
            page_texts = {1: description}
        if not text:
            logger.warning("OCR and VL both failed to extract text from image")
            text = safe_error_response("无法从图片提取文本，请确保图片清晰或使用其他格式。")
    else:
        try:
            file_storage.seek(0)
            md = MarkItDown()
            result = md.convert(BytesIO(file_bytes), file_extension=original_ext.lstrip('.'))
            text = result.text_content
            if text and text.strip():
                page_texts = {1: text}
            else:
                text = "[No text extracted by MarkItDown]"
        except Exception as e:
            logger.error(f"MarkItDown parsing failed for {original_ext}: {e}")
            text = safe_error_response(f"不支持的文件格式: {original_ext}", log_error=e)

    if not text or text.startswith("["):
        return text, page_texts

    analyze_images = session.get('analyze_images', True)
    if analyze_images:
        cached_desc = get_cached_image_description(file_hash)
        if cached_desc:
            image_desc = cached_desc
        else:
            image_desc, sample_info = describe_images_in_file(file_bytes, filename, page_texts if page_texts else None)
            if image_desc:
                cache_image_description(file_hash, image_desc)
            if sample_info:
                # stash for clearance report's 图片随机抽检说明 section
                with _IMAGE_SAMPLING_LOCK:
                    _IMAGE_SAMPLING_LOG.append({
                        'filename': filename,
                        'samples': sample_info,
                    })
        if image_desc:
            text += "\n\n--- Image Descriptions ---\n" + image_desc

    def clean_report_headers(text):
        if not text:
            return text
        lines = text.split('\n')
        cleaned_lines = []
        skip = False
        for line in lines:
            stripped = line.strip()
            # Check if this line is a header to skip
            is_header = any(stripped.startswith(prefix) for prefix in [
                '--- Sheet:', '技术标规律性分析检查结果', '标段名称：', '投标单位个数：',
                '检查结果：', '检查规则：', '相似度计算说明：', '一、标书围串风险分析结果',
                '二、分析结果详情', '签字：', '日期：', '技术标规律性分析详情'
            ])
            if is_header:
                skip = True
                continue
            if skip:
                # Only stop skipping when we hit a blank line (section separator)
                if stripped == '':
                    skip = False
                continue
            cleaned_lines.append(line)
        return '\n'.join(cleaned_lines)

    if text:
        text = clean_report_headers(text)

    return text, page_texts


# ── Path-based streaming extraction (large-file friendly) ────────────────

class _BytesFileShim:
    """Minimal FileStorage-like wrapper so legacy extractors can work on a
    disk path without loading the whole file twice."""

    def __init__(self, abs_path, filename=None):
        self._abs_path = abs_path
        self.filename = filename or os.path.basename(abs_path)
        self._buf = None

    def read(self):
        if self._buf is None:
            with open(self._abs_path, 'rb') as f:
                self._buf = f.read()
        return self._buf

    def seek(self, pos):
        # Legacy code calls seek(0) after read(); drop cache to re-read
        if pos == 0:
            self._buf = None
        return 0


def _extract_pdf_streaming(abs_path):
    """Page-by-page PDF text extraction straight from disk.

    fitz.open(filename=...) reads pages lazily, so peak memory is bounded by
    one rendered page instead of the whole file.
    Returns (text, page_texts) with the same semantics as extract_text_from_file.
    """
    text = None
    page_texts = {}
    try:
        doc = fitz.open(filename=abs_path)
    except Exception as e:
        logger.error(f"PDF open failed ({abs_path}): {e}", exc_info=True)
        return safe_error_response("无法解析PDF文件，请确保文件未损坏。", log_error=e), {}

    try:
        full_text = []
        has_text = False
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            page_text = page.get_text()
            if page_text and page_text.strip():
                has_text = True
                full_text.append(page_text)
                page_texts[page_num + 1] = page_text
            # Page objects are not retained; fitz reclaims them on next load_page
        extracted = "\n".join(full_text).strip()

        if has_text and len(extracted) > 50:
            text = extracted
        else:
            logger.info("PDF appears to be scanned (no text). Starting page-by-page OCR...")
            if ocr_manager.is_available():
                ocr_results = []
                ocr_page_texts = {}
                mat = fitz.Matrix(2.0, 2.0)
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    pix = page.get_pixmap(matrix=mat, alpha=False)
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    max_dim = 2000
                    if max(img.size) > max_dim:
                        ratio = max_dim / max(img.size)
                        img = img.resize((int(img.size[0] * ratio), int(img.size[1] * ratio)),
                                         Image.Resampling.LANCZOS)
                    page_ocr = run_ocr(np.array(img))
                    ocr_results.append(page_ocr or "")
                    ocr_page_texts[page_num + 1] = page_ocr or ""
                    del pix, img  # free the page bitmap before the next page
                if any(t.strip() for t in ocr_results):
                    text = "\n\n".join(ocr_results)
                    page_texts = ocr_page_texts
                elif _vl_available():
                    text = ""
                    for page_num in range(len(doc)):
                        if _vl_broken:
                            break
                        page = doc.load_page(page_num)
                        pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
                        desc = vl_model.describe_pdf_page(pix.tobytes("png"), page_num + 1)
                        if desc:
                            _vl_note_success()
                        else:
                            if _vl_note_failure():
                                break
                        text += f"\n\n--- 第{page_num + 1}页 (VL分析) ---\n{desc}"
                    page_texts = {i + 1: "" for i in range(len(doc))}
                else:
                    text = "[No text detected in PDF even after OCR and VL not available]"
            elif _vl_available():
                text = ""
                for page_num in range(len(doc)):
                    if _vl_broken:
                        break
                    page = doc.load_page(page_num)
                    pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
                    desc = vl_model.describe_pdf_page(pix.tobytes("png"), page_num + 1)
                    if desc:
                        _vl_note_success()
                    else:
                        if _vl_note_failure():
                            break
                    text += f"\n\n--- 第{page_num + 1}页 (VL分析) ---\n{desc}"
                page_texts = {i + 1: "" for i in range(len(doc))}
            else:
                text = "[无法提取PDF文本，且OCR/VL不可用]"
    except Exception as e:
        logger.error(f"PDF streaming error ({abs_path}): {e}", exc_info=True)
        return safe_error_response("无法解析PDF文件，请确保文件未损坏。", log_error=e), {}
    finally:
        try:
            doc.close()
        except Exception:
            pass

    return text, page_texts


def extract_text_from_path(abs_path, filename=None):
    """Extract text from a file on disk without loading it fully into memory.

    Large PDFs are processed page-by-page (constant memory); plain-text formats
    are streamed in chunks; other formats (docx/xlsx/pptx/images — typically
    small) delegate to the existing bytes-based extractor.
    Returns (text, page_texts).
    """
    filename = filename or os.path.basename(abs_path)
    ext = os.path.splitext(filename)[1].lower()
    wps_map = {'.wps': '.doc', '.et': '.xls', '.dps': '.ppt'}
    ext = wps_map.get(ext, ext)

    # PDF — the big-file case — streams page by page
    if ext == '.pdf':
        text, page_texts = _extract_pdf_streaming(abs_path)
        if text and not text.startswith("["):
            text = _clean_report_headers_static(text)
        return text, page_texts

    # Plain text formats — stream in chunks
    if ext in ['.txt', '.md', '.text', '.csv', '.json']:
        chunks = []
        try:
            with open(abs_path, 'rb') as f:
                while True:
                    chunk = f.read(8 * 1024 * 1024)
                    if not chunk:
                        break
                    chunks.append(chunk)
        except OSError as e:
            logger.error(f"Text read failed ({abs_path}): {e}")
            return safe_error_response("无法读取文件。", log_error=e), {1: ''}
        data = b''.join(chunks)
        try:
            text = data.decode('utf-8')
        except UnicodeDecodeError:
            text = data.decode('utf-8', errors='replace')
        return text, {1: text}

    if ext in ['.html', '.htm']:
        shim = _BytesFileShim(abs_path, filename)
        return extract_text_from_file(shim)

    # Everything else delegates via the shim (docx/xlsx/pptx/images/etc.)
    shim = _BytesFileShim(abs_path, filename)
    return extract_text_from_file(shim)


def _clean_report_headers_static(text):
    """Header cleanup shared by the streaming extractor (mirror of the inner
    clean_report_headers in extract_text_from_file)."""
    if not text:
        return text
    lines = text.split('\n')
    cleaned_lines = []
    skip = False
    for line in lines:
        stripped = line.strip()
        is_header = any(stripped.startswith(prefix) for prefix in [
            '--- Sheet:', '技术标规律性分析检查结果', '标段名称：', '投标单位个数：',
            '检查结果：', '检查规则：', '相似度计算说明：', '一、标书围串风险分析结果',
            '二、分析结果详情', '签字：', '日期：', '技术标规律性分析详情'
        ])
        if is_header:
            skip = True
            continue
        if skip:
            if stripped == '':
                skip = False
            continue
        cleaned_lines.append(line)
    return '\n'.join(cleaned_lines)


def extract_metadata_from_path(abs_path, filename=None):
    """Metadata extraction from a file on disk (PDF opened from path)."""
    filename = filename or os.path.basename(abs_path)
    meta = {'filename': filename}
    ext = os.path.splitext(filename)[1].lower()
    if ext == '.pdf':
        try:
            doc = fitz.open(filename=abs_path)
            info = doc.metadata or {}
            meta['author'] = info.get('author', '')
            meta['creator'] = info.get('creator', '')
            meta['producer'] = info.get('producer', '')
            meta['creationDate'] = info.get('creationDate', '')
            doc.close()
        except Exception:
            pass
        return meta
    # Other formats reuse the bytes path via shim
    shim = _BytesFileShim(abs_path, filename)

    class _ShimNameOnly(_BytesFileShim):
        pass
    return extract_metadata(shim)


def get_or_extract_file_analysis(file_storage, file_type, user_id, thread_id=None, project_id=None):
    filename = file_storage.filename
    file_bytes = file_storage.read()
    file_storage.seek(0)
    file_hash = hashlib.sha256(file_bytes).hexdigest()
    file_size = len(file_bytes)

    # Guard against empty/hash-of-empty bytes (indicates file was already consumed)
    EMPTY_HASH = 'e3b0c44298fc1c149afbf4c8996fb92427ae41e4649b934ca495991b7852b855'
    if file_hash == EMPTY_HASH or file_size == 0:
        logger.warning(f"get_or_extract_file_analysis: file {filename} appears empty (already consumed?)")
        return ""

    # Composite key: hash + size + user_id (per-user cache, prevents cross-user leak)
    cache_key = f"{file_hash}_{file_size}_{user_id}"

    with get_db_connection() as conn:
        with conn.cursor() as cur:
            # Check existing analysis (scoped to user and thread if available)
            cur.execute("""
                SELECT id, extracted_text FROM file_analysis
                WHERE file_hash = %s AND file_type = %s AND file_size = %s
                  AND user_id = %s AND deleted_at IS NULL
            """, (file_hash, file_type, file_size, user_id))
            row = cur.fetchone()
            if row:
                analysis_id, extracted_text = row
                # If cached text is invalid marker, force re‑extraction
                if extracted_text == "[INVALID_EXTRACTION]":
                    logger.warning(f"Cached extraction for {filename} marked as invalid. Re‑extracting.")
                    extracted_text, _ = extract_text_from_file(file_storage)
                    if not extracted_text or extracted_text.startswith("["):
                        extracted_text = ""
                    # Validate again
                    if not is_valid_extracted_text(extracted_text):
                        extracted_text = "[INVALID_EXTRACTION]"
                    # Update cache
                    cur.execute("""
                        UPDATE file_analysis
                        SET extracted_text = %s, usage_count = usage_count + 1, last_used_at = NOW()
                        WHERE id = %s
                    """, (extracted_text, analysis_id))
                    conn.commit()
                else:
                    # Normal valid cache – update usage count
                    cur.execute("""
                        UPDATE file_analysis
                        SET usage_count = usage_count + 1, last_used_at = NOW()
                        WHERE id = %s
                    """, (analysis_id,))
                    conn.commit()
                return extracted_text if extracted_text != "[INVALID_EXTRACTION]" else ""
            else:
                # No cache – extract and validate
                extracted_text, _ = extract_text_from_file(file_storage)
                if not extracted_text or extracted_text.startswith("["):
                    extracted_text = ""
                if not is_valid_extracted_text(extracted_text):
                    extracted_text = ""
                # Don't cache failed extractions — mark as deleted to avoid stale cache
                if not extracted_text:
                    cur.execute("DELETE FROM file_analysis WHERE file_hash = %s AND user_id = %s",
                                (file_hash, user_id))
                    conn.commit()
                    return ""
                # Store in cache
                cur.execute("""
                    INSERT INTO file_analysis (file_hash, file_type, file_size, original_filename, user_id, thread_id, project_id, extracted_text, usage_count, last_used_at)
                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s, 1, NOW())
                    RETURNING id
                """, (file_hash, file_type, file_size, filename, user_id, thread_id, project_id, extracted_text))
                analysis_id = cur.fetchone()[0]
                conn.commit()

                # Also update user_files.content if thread_id provided
                if thread_id:
                    cur.execute("""
                        SELECT id FROM user_files
                        WHERE user_id = %s AND file_hash = %s AND thread_id = %s
                        ORDER BY created_at DESC LIMIT 1
                    """, (user_id, file_hash, thread_id))
                    uf_row = cur.fetchone()
                    if uf_row:
                        cur.execute("""
                            UPDATE user_files SET content = %s WHERE id = %s
                        """, (extracted_text if extracted_text != "[INVALID_EXTRACTION]" else "", uf_row[0]))
                        conn.commit()

                return extracted_text if extracted_text != "[INVALID_EXTRACTION]" else ""

# Batch comparison helper functions
