"""Generate .docx / .xlsx from markdown content for project AI assistant downloads."""

import io
import os
import logging
from datetime import datetime
from typing import Optional, Tuple

logger = logging.getLogger(__name__)

try:
    from docx import Document
    from docx.shared import Inches, Pt, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl
    from openpyxl.styles import Font, Alignment, Border, Side, PatternFill
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False


def markdown_to_docx(md_text: str, title: str = "AI生成文档") -> bytes:
    """Convert markdown text to a .docx file. Returns bytes."""
    if not HAS_DOCX:
        raise RuntimeError("python-docx not installed")

    doc = Document()
    # Set default font
    style = doc.styles['Normal']
    font = style.font
    font.name = 'SimSun'
    font.size = Pt(11)

    # Title
    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title_para.add_run(title)
    run.bold = True
    run.font.size = Pt(16)
    doc.add_paragraph()  # spacer

    # Parse markdown lines
    lines = md_text.split('\n')
    for line in lines:
        line = line.rstrip()
        if not line:
            doc.add_paragraph()
            continue

        # Headers
        if line.startswith('### '):
            p = doc.add_paragraph()
            run = p.add_run(line[4:])
            run.bold = True
            run.font.size = Pt(13)
        elif line.startswith('## '):
            p = doc.add_paragraph()
            run = p.add_run(line[3:])
            run.bold = True
            run.font.size = Pt(14)
        elif line.startswith('# '):
            p = doc.add_paragraph()
            run = p.add_run(line[2:])
            run.bold = True
            run.font.size = Pt(15)
        # Bullet points
        elif line.strip().startswith('- ') or line.strip().startswith('* '):
            text = line.strip()[2:]
            p = doc.add_paragraph(text, style='List Bullet')
        # Numbered list
        elif line.strip() and line.strip()[0].isdigit() and '. ' in line.strip()[:4]:
            text = line.strip().split('. ', 1)[1] if '. ' in line.strip() else line.strip()
            p = doc.add_paragraph(text, style='List Number')
        # Bold markers
        elif '**' in line:
            p = doc.add_paragraph()
            parts = line.split('**')
            for i, part in enumerate(parts):
                run = p.add_run(part)
                if i % 2 == 1:
                    run.bold = True
        else:
            doc.add_paragraph(line)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def markdown_to_xlsx(md_text: str, sheet_name: str = "AI生成") -> bytes:
    """Convert markdown table(s) to .xlsx. Falls back to plain text if no table found."""
    if not HAS_XLSX:
        raise RuntimeError("openpyxl not installed")

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = sheet_name[:31]  # sheet name max 31 chars

    thin_border = Border(
        left=Side(style='thin'), right=Side(style='thin'),
        top=Side(style='thin'), bottom=Side(style='thin')
    )
    header_fill = PatternFill(start_color='D9E1F2', end_color='D9E1F2', fill_type='solid')
    header_font = Font(bold=True, size=11)

    # Try to find markdown tables
    lines = md_text.split('\n')
    tables = []
    current_table = []
    in_table = False

    for line in lines:
        stripped = line.strip()
        if stripped.startswith('|') and stripped.endswith('|'):
            # Skip separator lines like |---|---|
            if all(c in '|-: ' for c in stripped.replace('|', '').strip()):
                continue
            cells = [c.strip() for c in stripped.split('|')[1:-1]]
            current_table.append(cells)
            in_table = True
        else:
            if in_table and current_table:
                tables.append(current_table)
                current_table = []
            in_table = False
    if in_table and current_table:
        tables.append(current_table)

    if tables:
        # Write first table found
        row_idx = 1
        for table in tables[0]:
            for col_idx, cell in enumerate(table, 1):
                c = ws.cell(row=row_idx, column=col_idx, value=cell)
                c.border = thin_border
                if row_idx == 1:
                    c.fill = header_fill
                    c.font = header_font
            row_idx += 1
        # Auto-width
        for col in ws.columns:
            max_len = 0
            for cell in col:
                if cell.value:
                    max_len = max(max_len, len(str(cell.value)))
            ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 40)
    else:
        # No table found: put the whole text in one cell
        ws.cell(row=1, column=1, value="AI生成内容").font = header_font
        ws.cell(row=2, column=1, value=md_text[:30000])
        ws.column_dimensions['A'].width = 80

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


def markdown_to_pptx(md_text: str, title: str = "AI生成演示") -> bytes:
    """Convert markdown to a .pptx presentation. Each ## heading becomes a slide."""
    if not HAS_DOCX:
        raise RuntimeError("python-pptx not installed")
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = datetime.now().strftime('%Y-%m-%d')

    # Parse sections (## headings → slides)
    lines = md_text.split('\n')
    current_slide = None
    current_body = None
    current_text = []

    for line in lines:
        line = line.rstrip()
        if line.startswith('## ') or line.startswith('# '):
            # Save previous slide content
            if current_body and current_text:
                tf = current_body.text_frame
                tf.text = '\n'.join(current_text)[:3000]
                for para in tf.paragraphs:
                    para.font.size = Pt(14)
            # Start new slide
            current_slide = prs.slides.add_slide(prs.slide_layouts[1])
            current_slide.shapes.title.text = line.lstrip('#').strip()[:80]
            current_body = current_slide.placeholders[1]
            current_text = []
        elif line.strip():
            current_text.append(line)

    # Save last slide
    if current_body and current_text:
        tf = current_body.text_frame
        tf.text = '\n'.join(current_text)[:3000]

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


def generate_file(md_text: str, fmt: str, title: str = "AI生成") -> Tuple[bytes, str, str]:
    """Generate a file from markdown text.
    
    Returns: (file_bytes, filename, mime_type)
    """
    from app.services.prompt_safety import validate_markdown_structure

    ts = datetime.now().strftime('%Y%m%d_%H%M%S')
    safe_title = title[:30].replace('/', '_').replace('\\', '_')

    # ── Pre-check markdown structure ──
    valid, issues = validate_markdown_structure(md_text)
    if not valid:
        logger.warning(f"Markdown structure issues in generated file '{safe_title}': {issues}")
        # Auto-fix common issues: close unclosed fences
        fence_count = md_text.count('```')
        if fence_count % 2 != 0:
            md_text = md_text + '\n```'
            logger.info("Auto-closed unclosed code fence in markdown")
        # Re-check after fix
        valid, issues = validate_markdown_structure(md_text)
        if not valid:
            logger.warning(f"Markdown issues persist after auto-fix: {issues}")

    if fmt == 'docx':
        data = markdown_to_docx(md_text, title)
        return data, f"{safe_title}_{ts}.docx", 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
    elif fmt == 'xlsx':
        data = markdown_to_xlsx(md_text, title)
        return data, f"{safe_title}_{ts}.xlsx", 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    elif fmt == 'pptx':
        data = markdown_to_pptx(md_text, title)
        return data, f"{safe_title}_{ts}.pptx", 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    else:
        raise ValueError(f"Unsupported format: {fmt}. Use 'docx', 'xlsx', or 'pptx'.")
